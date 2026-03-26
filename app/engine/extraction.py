"""
Document extraction engine.

Extracts COMPANY data from pitch decks and documents.
Does NOT extract VoLo deal terms (check size, ownership, etc.) —
those are user inputs in the deal configuration step.
"""

import io
import json
import os
import re
import logging
from typing import Optional

logger = logging.getLogger(__name__)


# ── OCR fallback for image-heavy PDFs ─────────────────────────────────────────

def extract_text_pdf_with_ocr(file_bytes: bytes) -> str:
    text = ""
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        page_count = len(reader.pages)
    except Exception as e:
        logger.warning(f"pypdf failed: {e}")
        page_count = 0

    text_density = len(text.strip()) / max(page_count, 1) if page_count else 0

    if text_density < 100:
        logger.info("Low text density from pypdf, trying pdfplumber")
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                parts = []
                for page in pdf.pages:
                    t = page.extract_text() or ""
                    parts.append(t)
                    for table in (page.extract_tables() or []):
                        for row in table:
                            parts.append(" | ".join(str(c or "") for c in row))
                plumber_text = "\n".join(parts)
                if len(plumber_text.strip()) > len(text.strip()):
                    text = plumber_text
                    text_density = len(text.strip()) / max(page_count, 1)
        except ImportError:
            logger.warning("pdfplumber not available")
        except Exception as e:
            logger.warning(f"pdfplumber failed: {e}")

    if text_density < 50 and page_count > 0:
        logger.info("Very low text density, attempting OCR")
        try:
            ocr_text = _ocr_pdf(file_bytes)
            if len(ocr_text.strip()) > len(text.strip()):
                text = ocr_text
        except Exception as e:
            logger.warning(f"OCR failed: {e}")

    return text[:80000]


def _ocr_pdf(file_bytes: bytes) -> str:
    try:
        from pdf2image import convert_from_bytes
        import pytesseract
    except ImportError:
        logger.warning("pdf2image or pytesseract not installed")
        return ""
    try:
        images = convert_from_bytes(file_bytes, dpi=200, first_page=1, last_page=30)
    except Exception as e:
        logger.warning(f"pdf2image conversion failed: {e}")
        return ""
    parts = []
    for i, img in enumerate(images):
        try:
            page_text = pytesseract.image_to_string(img, lang="eng")
            parts.append(f"--- Page {i+1} ---\n{page_text}")
        except Exception as e:
            logger.warning(f"OCR failed on page {i+1}: {e}")
    return "\n".join(parts)


def extract_tables_from_pdf(file_bytes: bytes) -> list:
    tables = []
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for i, page in enumerate(pdf.pages):
                for tbl in (page.extract_tables() or []):
                    if not tbl or len(tbl) < 2:
                        continue
                    headers = [str(c or "").strip() for c in tbl[0]]
                    rows = [[str(c or "").strip() for c in row] for row in tbl[1:]]
                    tables.append({"page": i + 1, "headers": headers, "rows": rows})
    except ImportError:
        logger.warning("pdfplumber not installed")
    except Exception as e:
        logger.warning(f"Table extraction failed: {e}")
    return tables


def extract_text_pptx_enhanced(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts = [f"--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_parts.append(t)
            if shape.has_table:
                tbl = shape.table
                for row in tbl.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_parts.append(" | ".join(cells))
        if len(slide_parts) > 1:
            parts.extend(slide_parts)
    return "\n".join(parts)[:80000]


# ── Two-pass Claude extraction ────────────────────────────────────────────────

_ENERGY_FEW_SHOT = """Example extraction for an energy/climate tech company:
{
  "name": "SolarFlex Inc",
  "technology_description": "Flexible thin-film solar panels for commercial rooftops, achieving 22% efficiency at 40% lower cost than rigid panels",
  "stage": "Series A",
  "commercial_launch_yr": 2022,
  "tam_claim": 85000000000,
  "tam_units": "USD",
  "sam_claim": 12000000000,
  "market_geography": "North America, Europe",
  "unit_definition": "MW of flexible solar panel capacity",
  "unit_volumes_projected": [5, 15, 40, 100, 200, 350, 500, 700, 900, 1100],
  "revenue_projections": [2.5, 8, 22, 60, 130, 240, 370, 530, 700, 900],
  "current_revenue": 1.2,
  "revenue_units": "millions USD",
  "growth_rate_claim": "3x year-over-year",
  "unit_economics": {"cost_per_unit": 0.45, "price_per_unit": 0.72, "unit": "$/W", "gross_margin": "38%"},
  "funding_raised": 18000000,
  "team_size": 45,
  "competitive_landscape": "First Solar, SunPower, traditional rigid panel manufacturers",
  "displaced_technology": "conventional electricity generation from natural gas and coal",
  "trl_indicators": "pilot manufacturing line operational, 3 commercial installations completed",
  "key_risks": ["manufacturing scale-up", "supply chain for thin-film materials", "building code approvals"],
  "confidence": {"overall": 0.75, "fields": {"name": 1.0, "tam_claim": 0.6, "revenue_projections": 0.5, "unit_economics": 0.7}},
  "notes": "TAM from management estimate of global commercial rooftop solar market. Revenue projections from deck slide 14."
}"""

_EXTRACT_SCHEMA = """{
  "name": "string or null - company name",
  "technology_description": "string or null - what the company does and how it works",
  "stage": "Pre-Seed|Seed|Series A|Series B|Series C|Growth or null",
  "commercial_launch_yr": "integer year or null - when product first sold commercially",
  "tam_claim": "number or null - total addressable market in raw currency units (not millions)",
  "tam_units": "string - currency like USD, EUR, or unit like MWh, tonnes",
  "sam_claim": "number or null - serviceable addressable market",
  "market_geography": "string or null - target markets/regions",
  "unit_definition": "string or null - what one unit of their product is (MW, vehicle, tonne, etc.)",
  "unit_volumes_projected": "[up to 10 numbers] or null - annual unit sales/deployments projected",
  "revenue_projections": "[up to 10 numbers] or null - annual revenue figures if stated",
  "current_revenue": "number or null - most recent annual revenue",
  "revenue_units": "string - e.g. 'millions USD', 'thousands EUR'",
  "growth_rate_claim": "string or null - stated growth rate",
  "unit_economics": "object or null - cost per unit, price per unit, margin, etc.",
  "funding_raised": "number or null - total funding raised to date in raw currency",
  "team_size": "integer or null - number of employees",
  "competitive_landscape": "string or null - key competitors or alternatives",
  "displaced_technology": "string or null - what existing technology/resource this replaces",
  "trl_indicators": "string or null - evidence of technology readiness (pilots, demos, certifications)",
  "key_risks": "[strings] or null - stated or obvious risks",
  "confidence": {"overall": "float 0-1", "fields": {"field_name": "float 0-1"}},
  "notes": "string - any additional context or caveats"
}"""


def two_pass_extract(text: str, source_name: str, tables: list = None, model: str = None) -> dict:
    """
    Two-pass extraction:
    Pass 1: Raw field extraction from document text.
    Pass 2: Validation, cross-reference, and confidence calibration.
    """
    # Resolve model — default to Haiku, support Anthropic + Refiant providers
    _model = model or "claude-haiku-4-5-20251001"
    _is_refiant = _model.startswith("qwen")

    if _is_refiant:
        api_key = os.environ.get("REFIANT_API_KEY", "")
        if not api_key:
            raise ValueError("REFIANT_API_KEY not set")
        try:
            import anthropic
        except ImportError:
            raise ValueError("anthropic library not installed")
        client = anthropic.Anthropic(api_key=api_key, base_url="https://api.refiant.ai/v1")
    else:
        try:
            import anthropic
        except ImportError:
            raise ValueError("anthropic library not installed")
        api_key = os.environ.get("ANTHROPIC_API_KEY", "")
        if not api_key:
            raise ValueError("ANTHROPIC_API_KEY not set")
        client = anthropic.Anthropic(api_key=api_key)

    table_context = ""
    if tables:
        table_context = "\n\nEXTRACTED TABLES:\n"
        for t in tables[:10]:
            table_context += f"Page {t['page']}: {' | '.join(t['headers'])}\n"
            for row in t['rows'][:5]:
                table_context += f"  {' | '.join(row)}\n"

    pass1_prompt = (
        f"You are a venture capital analyst extracting company information "
        f"from a startup pitch deck or investor document.\n\n"
        f"Your job is to extract ONLY what the company states about itself — "
        f"market claims, financials, technology, team, and traction. "
        f"Do NOT guess at investor-side deal terms like ownership percentage "
        f"or investment amounts.\n\n"
        f"Source: {source_name}\n\nTEXT:\n{text[:60000]}{table_context}\n\n"
        f"Here is an example of a good extraction:\n{_ENERGY_FEW_SHOT}\n\n"
        f"Now extract fields from the document above. Return ONLY a valid JSON object:\n{_EXTRACT_SCHEMA}\n\n"
        f"RULES:\n"
        f"- Return ONLY JSON, no markdown fences\n"
        f"- Use null for fields you cannot determine from the document\n"
        f"- tam_claim should be in raw currency units (e.g. 85000000000 for $85B)\n"
        f"- revenue_projections in the same units as stated in the document\n"
        f"- Confidence: 0.0=guessing, 0.5=inferred from context, 1.0=explicitly stated\n"
        f"- Do NOT fabricate numbers — if it's not in the document, use null\n"
    )

    msg1 = client.messages.create(
        model=_model,
        max_tokens=3000,
        messages=[{"role": "user", "content": pass1_prompt}],
    )
    raw1 = msg1.content[0].text.strip()
    raw1 = re.sub(r"^```[a-z]*\n?", "", raw1).rstrip("`").strip()
    try:
        pass1_result = json.loads(raw1)
    except json.JSONDecodeError:
        m = re.search(r"\{.*\}", raw1, re.DOTALL)
        if m:
            pass1_result = json.loads(m.group())
        else:
            raise ValueError("Could not parse JSON from Pass 1")

    pass2_prompt = (
        f"You previously extracted these fields from a startup document:\n"
        f"{json.dumps(pass1_result, indent=2)}\n\n"
        f"Now validate and correct the extraction:\n"
        f"1. If tam_claim seems unreasonably large (>$10T) or small (<$100K), flag it\n"
        f"2. If revenue_projections don't show a plausible trajectory, flag it\n"
        f"3. Verify the revenue_units match the magnitude of revenue figures\n"
        f"4. Check that unit_volumes_projected and revenue_projections are consistent\n"
        f"5. Verify commercial_launch_yr is between 2010-2040\n"
        f"6. Recalibrate confidence scores honestly\n"
        f"7. Add _validation_warnings as an array of strings for anything questionable\n\n"
        f"Return the corrected JSON object (same schema, with _validation_warnings added).\n"
        f"Return ONLY JSON, no markdown fences or explanation."
    )

    msg2 = client.messages.create(
        model=_model,
        max_tokens=3000,
        messages=[{"role": "user", "content": pass2_prompt}],
    )
    raw2 = msg2.content[0].text.strip()
    raw2 = re.sub(r"^```[a-z]*\n?", "", raw2).rstrip("`").strip()
    try:
        result = json.loads(raw2)
    except json.JSONDecodeError:
        m = re.search(r"\{.*\}", raw2, re.DOTALL)
        result = json.loads(m.group()) if m else pass1_result

    result = validate_extraction(result)
    return result


def validate_extraction(result: dict) -> dict:
    """Domain-constraint validation with warnings."""
    warnings = result.get("_validation_warnings", [])

    tam = result.get("tam_claim")
    if tam is not None:
        if tam < 0:
            warnings.append("tam_claim negative — set to null")
            result["tam_claim"] = None
        elif tam > 1e13:
            warnings.append(f"tam_claim={tam} suspiciously large (>$10T) — verify units")

    launch = result.get("commercial_launch_yr")
    if launch is not None and (launch < 2000 or launch > 2045):
        warnings.append(f"commercial_launch_yr={launch} out of range")
        result["commercial_launch_yr"] = None

    for field in ("revenue_projections", "unit_volumes_projected"):
        vols = result.get(field)
        if vols is not None:
            if not isinstance(vols, list):
                result[field] = None
                warnings.append(f"{field} not a list")
            else:
                cleaned = []
                for v in vols:
                    try:
                        cleaned.append(float(v or 0))
                    except (TypeError, ValueError):
                        cleaned.append(0)
                result[field] = cleaned

    team = result.get("team_size")
    if team is not None and (team < 0 or team > 100000):
        warnings.append(f"team_size={team} seems wrong")
        result["team_size"] = None

    if warnings:
        result["_validation_warnings"] = warnings

    return result


# ── Deal terms extraction from prior IC memos ────────────────────────────────

_DEAL_TERMS_SCHEMA = """{
  "company_name": "string or null - company name",
  "check_size_m": "number or null - investment check size in millions USD",
  "pre_money_m": "number or null - pre-money valuation in millions USD",
  "post_money_m": "number or null - post-money valuation in millions USD",
  "round_size_m": "number or null - total round size in millions USD",
  "ownership_pct": "number or null - ownership percentage acquired (0-100)",
  "entry_stage": "Pre-Seed|Seed|Series A|Series B|Series C|Growth or null",
  "investment_date": "string or null - date of investment (YYYY-MM-DD or YYYY)",
  "fund_year": "integer or null - fund year at time of investment (1-indexed)",
  "board_seat": "boolean or null - whether a board seat was obtained",
  "pro_rata_rights": "boolean or null - whether pro-rata rights were secured",
  "liquidation_preference": "string or null - liquidation preference terms (e.g. 1x non-participating)",
  "key_commitments": "[strings] or null - specific commitments or milestones company agreed to",
  "revenue_at_investment": "number or null - revenue at time of investment (in millions USD)",
  "arr_at_investment": "number or null - ARR at time of investment (in millions USD)",
  "confidence": {"overall": "float 0-1", "fields": {"field_name": "float 0-1"}},
  "notes": "string - any caveats about extracted values"
}"""

_DEAL_TERMS_FEW_SHOT = """Example extraction from a prior IC memo:
{
  "company_name": "CleanGrid Technologies",
  "check_size_m": 3.0,
  "pre_money_m": 25.0,
  "post_money_m": 30.0,
  "round_size_m": 5.0,
  "ownership_pct": 10.0,
  "entry_stage": "Seed",
  "investment_date": "2023-06",
  "fund_year": 2,
  "board_seat": true,
  "pro_rata_rights": true,
  "liquidation_preference": "1x non-participating preferred",
  "key_commitments": [
    "Achieve 5 pilot deployments by Q4 2024",
    "Reach $1M ARR within 18 months",
    "Complete UL certification for residential product",
    "Hire VP Engineering and VP Sales"
  ],
  "revenue_at_investment": 0.1,
  "arr_at_investment": null,
  "confidence": {"overall": 0.85, "fields": {"check_size_m": 0.95, "pre_money_m": 0.90, "round_size_m": 0.80, "entry_stage": 1.0}},
  "notes": "Check size explicitly stated on page 2. Pre-money derived from stated post-money minus round size. Revenue figure from executive summary."
}"""


def extract_deal_terms(text: str, source_name: str = "IC Memo", model: str = None) -> dict:
    """
    Extract prior investment deal terms from an IC memo or investment document.

    Returns structured deal terms including check size, valuation, ownership,
    stage, and key commitments.
    """
    _model = model or "claude-haiku-4-5-20251001"
    _is_refiant = _model.startswith("qwen")

    if _is_refiant:
        api_key = os.environ.get("REFIANT_API_KEY", "")
        if not api_key:
            raise ValueError("REFIANT_API_KEY not set")
        base_url = "https://api.refiant.ai/v1"
    else:
        api_key = os.environ.get("ANTHROPIC_API_KEY", "")
        if not api_key:
            raise ValueError("ANTHROPIC_API_KEY not set")
        base_url = None

    # Truncate text to fit context
    max_chars = 60000
    if len(text) > max_chars:
        text = text[:max_chars] + "\n\n[... truncated ...]"

    system_prompt = (
        "You are a financial document analyst specializing in venture capital investment memos. "
        "Extract INVESTOR-SIDE deal terms from this investment committee memo. "
        "Focus on: check size, valuation, ownership, stage, round size, and any "
        "specific commitments or milestones the portfolio company agreed to meet. "
        "Return ONLY valid JSON matching the schema below. "
        "If a value cannot be found, use null. "
        "Convert all monetary values to millions USD."
    )

    user_msg = (
        f"Extract deal terms from this document.\n\n"
        f"Source: {source_name}\n\n"
        f"Output JSON schema:\n{_DEAL_TERMS_SCHEMA}\n\n"
        f"Example:\n{_DEAL_TERMS_FEW_SHOT}\n\n"
        f"Document text:\n{text}"
    )

    try:
        if _is_refiant:
            import openai
            client = openai.OpenAI(api_key=api_key, base_url=base_url)
            resp = client.chat.completions.create(
                model=_model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_msg},
                ],
                max_tokens=2000,
                temperature=0.0,
            )
            raw = resp.choices[0].message.content.strip()
        else:
            import anthropic
            client = anthropic.Anthropic(api_key=api_key)
            resp = client.messages.create(
                model=_model,
                max_tokens=2000,
                temperature=0.0,
                system=system_prompt,
                messages=[{"role": "user", "content": user_msg}],
            )
            raw = resp.content[0].text.strip()

        # Parse JSON from response (handle markdown code fences)
        if raw.startswith("```"):
            raw = re.sub(r"^```\w*\n?", "", raw)
            raw = re.sub(r"\n?```$", "", raw)
        result = json.loads(raw)

        # Validate and clean numeric fields
        for field in ["check_size_m", "pre_money_m", "post_money_m", "round_size_m",
                       "ownership_pct", "revenue_at_investment", "arr_at_investment"]:
            val = result.get(field)
            if val is not None:
                try:
                    result[field] = float(val)
                except (TypeError, ValueError):
                    result[field] = None

        if result.get("fund_year") is not None:
            try:
                result["fund_year"] = int(result["fund_year"])
            except (TypeError, ValueError):
                result["fund_year"] = None

        # Derive missing values where possible
        if result.get("post_money_m") and result.get("round_size_m") and not result.get("pre_money_m"):
            result["pre_money_m"] = result["post_money_m"] - result["round_size_m"]
        if result.get("pre_money_m") and result.get("round_size_m") and not result.get("post_money_m"):
            result["post_money_m"] = result["pre_money_m"] + result["round_size_m"]
        if result.get("check_size_m") and result.get("post_money_m") and not result.get("ownership_pct"):
            result["ownership_pct"] = round(result["check_size_m"] / result["post_money_m"] * 100, 2)

        result["_extraction_source"] = source_name
        return result

    except json.JSONDecodeError as e:
        logger.warning("Deal terms extraction JSON parse failed: %s", e)
        return {"error": f"JSON parse error: {e}", "_extraction_source": source_name}
    except Exception as e:
        logger.warning("Deal terms extraction failed: %s", e)
        return {"error": str(e), "_extraction_source": source_name}
