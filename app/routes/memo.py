"""
Investment Memo routes — generate comprehensive investment memos from report data + data room documents.

POST   /api/memo/templates              — create a memo template
GET    /api/memo/templates              — list templates
GET    /api/memo/templates/{id}         — get template
PUT    /api/memo/templates/{id}         — update template
DELETE /api/memo/templates/{id}         — delete template

POST   /api/memo/documents/upload       — upload data room document(s)
GET    /api/memo/documents/{session_id} — list documents for a session
DELETE /api/memo/documents/{doc_id}     — remove a document

POST   /api/memo/generate               — generate memo (LLM call)
GET    /api/memo/history                — list generated memos
GET    /api/memo/history/{id}           — get a generated memo
DELETE /api/memo/history/{id}           — delete a generated memo
GET    /api/memo/history/{id}/docx      — export as Word document
GET    /api/memo/reports                — list available deal reports for selection
"""

import asyncio
import io
import json
import logging
import os
import re
import tempfile
import threading
import time
import uuid
from datetime import datetime
from pathlib import Path
from typing import Optional, List

from fastapi import APIRouter, Depends, File, Form, HTTPException, Query, UploadFile
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
from pydantic import BaseModel

from ..auth import CurrentUser, get_current_user, get_optional_user, decode_token
from ..database import get_db, get_model_preferences, MODEL_DEFAULTS

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/memo", tags=["memo"])

# ── In-memory memo job store ─────────────────────────────────────────────────
# Memo generation moved to a background worker so proxy/CDN HTTP timeouts
# (~2–3 min on Replit/Cloudflare) can't cut off a 3–5 min LLM pipeline.
# The HTTP endpoint returns a job_id immediately; the frontend polls
# /api/memo/generate/status/{job_id} for real per-call progress and the
# final result payload.
_MEMO_JOBS: dict[str, dict] = {}
_MEMO_LOCK = threading.Lock()

# Retention: jobs older than this are dropped from memory on access.
_MEMO_JOB_TTL_SECONDS = 3600  # 1 hour


def _prune_old_memo_jobs() -> None:
    """Drop memo jobs older than the TTL to keep memory bounded."""
    now = time.time()
    with _MEMO_LOCK:
        stale = [jid for jid, j in _MEMO_JOBS.items()
                 if now - j.get("_created_mono", now) > _MEMO_JOB_TTL_SECONDS]
        for jid in stale:
            _MEMO_JOBS.pop(jid, None)

# ── Upload directory ─────────────────────────────────────────────────────────
# Memo document uploads (PDFs, XLSXs, term sheets, screenshots, etc.) live
# on the local filesystem; the DB stores only the path + metadata.
#
# Default lives inside the source tree at <repo>/data/memo_uploads/, fine for
# local dev. On Replit Reserved-VM deployments the source tree is replaced
# on every redeploy, which would wipe every previously-uploaded file. Set
# VOLO_UPLOADS_DIR to a path *outside* the source tree (e.g.
# /home/runner/.volo/uploads) on the deployment to make uploads persist.
# The directory is created automatically and any legacy files are migrated
# on first run with the new path.
import shutil as _shutil

_LEGACY_UPLOAD_DIR = Path(__file__).resolve().parent.parent.parent / "data" / "memo_uploads"
_env_uploads = os.environ.get("VOLO_UPLOADS_DIR", "").strip()
if _env_uploads:
    UPLOAD_DIR = Path(_env_uploads)
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    # One-time migration: copy any existing legacy files over so historical
    # data-room uploads remain accessible. Skip files that already exist
    # at the destination so re-runs are safe.
    if _LEGACY_UPLOAD_DIR.exists() and _LEGACY_UPLOAD_DIR != UPLOAD_DIR:
        try:
            n_copied = 0
            for src in _LEGACY_UPLOAD_DIR.rglob("*"):
                if not src.is_file():
                    continue
                rel = src.relative_to(_LEGACY_UPLOAD_DIR)
                dst = UPLOAD_DIR / rel
                if dst.exists():
                    continue
                dst.parent.mkdir(parents=True, exist_ok=True)
                _shutil.copy2(src, dst)
                n_copied += 1
            # Rewrite DB rows that still point at the legacy absolute path so
            # /view, /delete, and memo generation all find the migrated files.
            try:
                from ..database import get_db as _get_db
                _legacy_prefix = str(_LEGACY_UPLOAD_DIR)
                _new_prefix = str(UPLOAD_DIR)
                _db = _get_db()
                try:
                    cur = _db.execute(
                        "UPDATE memo_documents "
                        "SET file_path = ? || substr(file_path, ? + 1) "
                        "WHERE file_path LIKE ? || '%'",
                        (_new_prefix, len(_legacy_prefix), _legacy_prefix),
                    )
                    _db.commit()
                    n_rows = cur.rowcount or 0
                finally:
                    _db.close()
            except Exception as _db_err:
                n_rows = -1
                print(f"[VoLo Engine] WARN: memo_documents path rewrite failed: {_db_err}", flush=True)
            print(f"[VoLo Engine] Memo uploads migrated: {n_copied} files copied, "
                  f"{n_rows} DB rows rewritten -> {UPLOAD_DIR}", flush=True)
        except OSError as e:
            print(f"[VoLo Engine] WARN: memo upload migration failed: {e}", flush=True)
else:
    UPLOAD_DIR = _LEGACY_UPLOAD_DIR
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

# Maximum file size: 50 MB
MAX_FILE_SIZE = 50 * 1024 * 1024

# Allowed file extensions for data room documents
ALLOWED_EXTENSIONS = {
    '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.csv', '.txt',
    '.pptx', '.ppt', '.md', '.html', '.json', '.png', '.jpg', '.jpeg', '.gif', '.webp', '.svg'
}

DOC_CATEGORIES = [
    'financial_model', 'pitch_deck', 'term_sheet', 'cap_table',
    'legal', 'ip_patent', 'customer_reference', 'market_research',
    'technical_diligence', 'team_bios', 'board_materials', 'screenshots', 'other'
]


# ═══════════════════════════════════════════════════════════════════════════════
#  DEFAULT VoLo TEMPLATE — auto-seeded for new users
# ═══════════════════════════════════════════════════════════════════════════════

VOLO_DEFAULT_TEMPLATE_NAME = "VoLo Earth Ventures — Standard IC Memo"

VOLO_DEFAULT_TEMPLATE = """# Investment Overview

(LOGO) One-Liner

| Field | Value |
|-------|-------|
| VEV Focus Area | |
| Investment Stage & Vehicle | |
| Key Terms | |
| Capital Amount | |
| Expected VEV Equity | |
| Board Participation | |

## Portfolio Themes
- Low cost, next gen energy storage
- Uses abundant materials
- Digitization of high skill activities
- System-level enabling technology
- Can be sold with zero-dollar year 1 premium
- Increase productivity and profitability of at-scale manufacturing base
- Rapidly scalable after achieving derisking milestone(s)

# High Level Opportunities
- [Key investment thesis points]

# High Level Risks
- [Key risk factors]

# Company Overview
- Set a clear tone: What is the company's 'Aha' or 'breakthrough'?
- Contextualize Opportunity: Why does the market need this? How big is the market, what is the gap, and why now?
- Why this categorical approach to the problem? What bottlenecks do they solve?
- Differentiation
- Validate: Traction + Defensibility (high level; tech details go in Traction)
- Entice: Preview 'how big this can get'? Why can they uniquely succeed?

# Market
- Start broad, with total demand and market size
- Describe market dynamics of current market composition by major segments
- Serviceable market — where does this company have a right to win?

# Business Model
- How does the company make money? Unit Economics
- Scalability — map of capex needs and financing plans
- Broad GTM
- Highlight relevant strengths from RVM
- Forward Projections

# Team
- Quick overview of relevant background (operating + founding experience, combined skills)
- Execution to date in relation to VoLo's experience with the team
- Additional signals that they will be able to execute against venture trajectory/return
- Include relevant factors of RVM

# Traction
| Category | Details |
|----------|---------|
| Technical | - |
| Commercial | - |
| Customer Feedback | - |

# Competitive Position
- Broad framework to categorize competition: zoom out → in
- Quick overview on why/where other categories are limited; include incumbents
- Lift & contextualize competitive landscape slide (if provided)

# Carbon Impact
- Theory of Change
- Carbon Economics: Include image from RVM or Carbon calc

# Technology: Overview, IP and Moat
- Technical brief: differentiation, what is derisked to date
- Derisking ahead: remaining categorical step changes, how VoLo got comfortable
- Highlight relevant strengths from RVM

# Financing Overview: Round and Exit Planning

## Deal Structure
Overview & VoLo Role — include total raise
VoLo Earth Value-add and Additionality
Committed Syndicate

## Capital Position and Funding History
| Metric | Value |
|--------|-------|
| Total Raised to Date | |
| Current Cash in Bank | |
| Current Burn Rate | |
| Runway Provided by Financing | |

## Round Pricing
- Price, Implied multiple (TTM x NTM)
- Pricing & Comps: public comps chart, comps discussion (describe categories & rationale)
- Private Comps

## Exit Planning
- Basic MOIC with revenue multiple range x discounted revenue projection (3-5 years)
- Acquisition: likely acquirers, strategic positioning
- IPO potential (if applicable)
- Downside Protection
- RVM Exit screenshot with assumptions

## Key Milestones & Use of Funds
- High level milestones for follow-on round accountability
- Brief discussion of feasibility

## Growth Investor Insights
"""


def seed_default_template(owner_id: int):
    """Ensure the default VoLo template exists for a user. Idempotent."""
    conn = get_db()
    try:
        existing = conn.execute(
            "SELECT id FROM memo_templates WHERE owner_id=? AND name=?",
            (owner_id, VOLO_DEFAULT_TEMPLATE_NAME),
        ).fetchone()
        if existing:
            return existing["id"]
        cur = conn.execute(
            """INSERT INTO memo_templates (owner_id, name, description, content, is_default)
               VALUES (?, ?, ?, ?, 1)""",
            (owner_id, VOLO_DEFAULT_TEMPLATE_NAME,
             "Standard VoLo Earth Ventures Investment Committee memo template",
             VOLO_DEFAULT_TEMPLATE),
        )
        conn.commit()
        return cur.lastrowid
    finally:
        conn.close()


# ═══════════════════════════════════════════════════════════════════════════════
#  TEMPLATE CRUD
# ═══════════════════════════════════════════════════════════════════════════════

class TemplateCreate(BaseModel):
    name: str
    description: str = ""
    content: str = ""


class TemplateUpdate(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    content: Optional[str] = None


@router.post("/templates")
async def create_template(req: TemplateCreate, user: CurrentUser = Depends(get_current_user)):
    conn = get_db()
    try:
        cur = conn.execute(
            """INSERT INTO memo_templates (owner_id, name, description, content)
               VALUES (?, ?, ?, ?)""",
            (user.id, req.name, req.description, req.content),
        )
        conn.commit()
        return {"id": cur.lastrowid, "name": req.name}
    finally:
        conn.close()


@router.get("/templates")
async def list_templates(user: CurrentUser = Depends(get_current_user)):
    # Auto-seed default template if user has none
    seed_default_template(user.id)
    conn = get_db()
    try:
        rows = conn.execute(
            "SELECT id, name, description, is_default, created_at, updated_at FROM memo_templates WHERE owner_id=? ORDER BY is_default DESC, updated_at DESC",
            (user.id,),
        ).fetchall()
        return [dict(r) for r in rows]
    finally:
        conn.close()


@router.get("/templates/{template_id}")
async def get_template(template_id: int, user: CurrentUser = Depends(get_current_user)):
    conn = get_db()
    try:
        row = conn.execute(
            "SELECT * FROM memo_templates WHERE id=? AND owner_id=?",
            (template_id, user.id),
        ).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Template not found")
        return dict(row)
    finally:
        conn.close()


@router.put("/templates/{template_id}")
async def update_template(template_id: int, req: TemplateUpdate, user: CurrentUser = Depends(get_current_user)):
    conn = get_db()
    try:
        existing = conn.execute(
            "SELECT id FROM memo_templates WHERE id=? AND owner_id=?",
            (template_id, user.id),
        ).fetchone()
        if not existing:
            raise HTTPException(status_code=404, detail="Template not found")

        updates = []
        params = []
        if req.name is not None:
            updates.append("name=?")
            params.append(req.name)
        if req.description is not None:
            updates.append("description=?")
            params.append(req.description)
        if req.content is not None:
            updates.append("content=?")
            params.append(req.content)
        if updates:
            updates.append("updated_at=datetime('now')")
            params.extend([template_id, user.id])
            conn.execute(
                f"UPDATE memo_templates SET {', '.join(updates)} WHERE id=? AND owner_id=?",
                params,
            )
            conn.commit()
        return {"ok": True}
    finally:
        conn.close()


@router.delete("/templates/{template_id}")
async def delete_template(template_id: int, user: CurrentUser = Depends(get_current_user)):
    conn = get_db()
    try:
        conn.execute("DELETE FROM memo_templates WHERE id=? AND owner_id=?", (template_id, user.id))
        conn.commit()
        return {"ok": True}
    finally:
        conn.close()


# ═══════════════════════════════════════════════════════════════════════════════
#  DOCUMENT UPLOAD & MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════════════

def _extract_text_from_file(file_path: Path, file_name: str) -> str:
    """Best-effort text extraction from uploaded files."""
    ext = Path(file_name).suffix.lower()
    text = ""

    try:
        if ext in ('.txt', '.md', '.csv', '.json', '.html'):
            text = file_path.read_text(errors='replace')[:100_000]

        elif ext == '.pdf':
            pdf_bytes = file_path.read_bytes()
            # Try pypdf first (always available)
            try:
                import pypdf
                reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
                pages = []
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        pages.append(f"[Page {i+1}]\n{page_text}")
                text = "\n\n".join(pages)
            except Exception as e1:
                logger.warning("pypdf extraction failed: %s", e1)
                text = ""

            # If pypdf got thin results, try pdfplumber
            page_count = 0
            try:
                import pypdf as _p
                page_count = len(_p.PdfReader(io.BytesIO(pdf_bytes)).pages)
            except Exception:
                pass
            text_density = len(text.strip()) / max(page_count, 1) if page_count else 0

            if text_density < 100 and page_count > 0:
                try:
                    import pdfplumber
                    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
                        plumber_parts = []
                        for i, page in enumerate(pdf.pages):
                            t = page.extract_text() or ""
                            if t.strip():
                                plumber_parts.append(f"[Page {i+1}]\n{t}")
                            for table in (page.extract_tables() or []):
                                for row in table:
                                    plumber_parts.append(" | ".join(str(c or "") for c in row))
                        plumber_text = "\n\n".join(plumber_parts)
                        if len(plumber_text.strip()) > len(text.strip()):
                            text = plumber_text
                except Exception as e2:
                    logger.warning("pdfplumber fallback failed: %s", e2)

            text = text[:200_000]
            if not text.strip():
                text = "[PDF extraction returned no text — document may be image-only]"

        elif ext == '.docx':
            try:
                from docx import Document
                doc = Document(str(file_path))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                text = "\n".join(paragraphs)[:200_000]
            except ImportError:
                text = "[DOCX text extraction requires python-docx]"

        elif ext in ('.xlsx', '.xls'):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(file_path), data_only=True, read_only=True)
                parts = []
                for ws in wb.worksheets[:5]:  # Limit to first 5 sheets
                    parts.append(f"=== Sheet: {ws.title} ===")
                    row_count = 0
                    for row in ws.iter_rows(values_only=True):
                        vals = [str(c) if c is not None else '' for c in row]
                        if any(v.strip() for v in vals):
                            parts.append('\t'.join(vals))
                            row_count += 1
                            if row_count > 200:
                                parts.append(f"... (truncated, {ws.max_row} total rows)")
                                break
                wb.close()
                text = "\n".join(parts)[:200_000]
            except ImportError:
                text = "[Excel text extraction requires openpyxl]"

        elif ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                parts = []
                for i, slide in enumerate(prs.slides):
                    parts.append(f"=== Slide {i+1} ===")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    parts.append(para.text.strip())
                text = "\n".join(parts)[:200_000]
            except ImportError:
                text = "[PPTX text extraction requires python-pptx]"

    except Exception as e:
        text = f"[Extraction error: {str(e)[:200]}]"

    return text


@router.post("/documents/upload")
async def upload_documents(
    files: List[UploadFile] = File(...),
    session_id: str = Form(""),
    category: str = Form("other"),
    user: CurrentUser = Depends(get_current_user),
):
    if not session_id:
        session_id = str(uuid.uuid4())[:12]

    results = []
    conn = get_db()
    try:
        for file in files:
            ext = Path(file.filename or "unknown").suffix.lower()
            if ext not in ALLOWED_EXTENSIONS:
                results.append({"file": file.filename, "error": f"File type {ext} not allowed"})
                continue

            # Read file content
            content = await file.read()
            if len(content) > MAX_FILE_SIZE:
                results.append({"file": file.filename, "error": "File exceeds 50MB limit"})
                continue

            # Save to disk
            safe_name = re.sub(r'[^\w\-.]', '_', file.filename or 'unnamed')
            dest_dir = UPLOAD_DIR / session_id
            dest_dir.mkdir(parents=True, exist_ok=True)
            dest_path = dest_dir / safe_name

            # Handle name collision
            counter = 1
            while dest_path.exists():
                stem = Path(safe_name).stem
                dest_path = dest_dir / f"{stem}_{counter}{ext}"
                counter += 1

            dest_path.write_bytes(content)

            # Extract text
            extracted = _extract_text_from_file(dest_path, safe_name)

            # Insert DB record
            cur = conn.execute(
                """INSERT INTO memo_documents
                   (owner_id, memo_session_id, file_name, file_type, file_size, extracted_text, doc_category, file_path)
                   VALUES (?, ?, ?, ?, ?, ?, ?, ?)""",
                (user.id, session_id, file.filename, ext, len(content), extracted, category, str(dest_path)),
            )
            conn.commit()
            doc_id = cur.lastrowid

            # For PDFs: extract embedded images and store them as separate image records
            # so they can be injected inline into the memo
            img_records_added = 0
            if ext == '.pdf':
                pdf_images = _extract_images_from_pdf(dest_path, dest_dir)
                for img in pdf_images:
                    try:
                        img_size = Path(img["file_path"]).stat().st_size
                        if img_size < 5000:   # skip tiny icons / artifacts
                            continue
                        conn.execute(
                            """INSERT INTO memo_documents
                               (owner_id, memo_session_id, file_name, file_type, file_size, extracted_text, doc_category, file_path)
                               VALUES (?, ?, ?, ?, ?, ?, ?, ?)""",
                            (user.id, session_id, img["file_name"], ".png",
                             img_size, "", category, img["file_path"]),
                        )
                        img_records_added += 1
                    except Exception as img_err:
                        logger.debug("Failed to store extracted image %s: %s", img["file_name"], img_err)
                if img_records_added:
                    conn.commit()

            results.append({
                "id": doc_id,
                "file": file.filename,
                "size": len(content),
                "extracted_chars": len(extracted),
                "category": category,
                "images_extracted": img_records_added,
            })

        return {"session_id": session_id, "documents": results}
    finally:
        conn.close()


@router.get("/documents/{session_id}")
async def list_documents(session_id: str, user: CurrentUser = Depends(get_current_user)):
    conn = get_db()
    try:
        rows = conn.execute(
            "SELECT id, file_name, file_type, file_size, doc_category, uploaded_at FROM memo_documents WHERE owner_id=? AND memo_session_id=? ORDER BY uploaded_at",
            (user.id, session_id),
        ).fetchall()
        return [dict(r) for r in rows]
    finally:
        conn.close()


@router.get("/documents/{doc_id}/view")
async def view_document(
    doc_id: int,
    token: Optional[str] = Query(None),
    user: Optional[CurrentUser] = Depends(get_optional_user),
):
    """Serve an uploaded document file for inline viewing/verification.
    Supports token via query param for <img> tags and new-tab access where
    the browser cannot send an Authorization header."""
    # Primary: use bearer-auth user.  Fallback: decode token from query param.
    # This fallback is essential for <img src="...?token=..."> — browsers never
    # send Authorization headers for image requests.
    effective_user = user
    if not effective_user and token:
        payload = decode_token(token)
        if payload:
            effective_user = CurrentUser(uid=int(payload["sub"]), username=payload["user"], role=payload["role"])
    if not effective_user:
        raise HTTPException(status_code=401, detail="Authentication required")

    conn = get_db()
    try:
        row = conn.execute(
            "SELECT file_name, file_path, file_type FROM memo_documents WHERE id=? AND owner_id=?",
            (doc_id, effective_user.id),
        ).fetchone()
        if not row or not row["file_path"]:
            raise HTTPException(404, "Document not found")
        fpath = Path(row["file_path"])
        if not fpath.exists():
            raise HTTPException(404, "File not found on disk")

        content_types = {
            '.pdf': 'application/pdf',
            '.png': 'image/png',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.txt': 'text/plain',
            '.md': 'text/plain',
            '.csv': 'text/csv',
            '.html': 'text/html',
        }
        ct = content_types.get(row["file_type"], 'application/octet-stream')
        return FileResponse(str(fpath), media_type=ct, filename=row["file_name"])
    finally:
        conn.close()


@router.delete("/documents/{doc_id}")
async def delete_document(doc_id: int, user: CurrentUser = Depends(get_current_user)):
    conn = get_db()
    try:
        row = conn.execute("SELECT file_path FROM memo_documents WHERE id=? AND owner_id=?", (doc_id, user.id)).fetchone()
        if row and row["file_path"]:
            try:
                Path(row["file_path"]).unlink(missing_ok=True)
            except Exception:
                pass
        conn.execute("DELETE FROM memo_documents WHERE id=? AND owner_id=?", (doc_id, user.id))
        conn.commit()
        return {"ok": True}
    finally:
        conn.close()


@router.post("/session/clear")
async def clear_memo_session(
    session_id: str = Query(...),
    user: CurrentUser = Depends(get_current_user),
):
    """Wipe a memo session: delete every uploaded data-room document for the
    given session_id (DB row + file on disk). Used by the "Start New Memo"
    button to give the analyst a clean slate between memos without making
    them remove files one by one."""
    conn = get_db()
    deleted = 0
    try:
        rows = conn.execute(
            "SELECT id, file_path FROM memo_documents WHERE owner_id=? AND memo_session_id=?",
            (user.id, session_id),
        ).fetchall()
        for r in rows:
            if r["file_path"]:
                try:
                    Path(r["file_path"]).unlink(missing_ok=True)
                except Exception:
                    pass
        conn.execute(
            "DELETE FROM memo_documents WHERE owner_id=? AND memo_session_id=?",
            (user.id, session_id),
        )
        deleted = len(rows)
        conn.commit()
    finally:
        conn.close()
    return {"ok": True, "deleted": deleted}


# ─────────────────────────────────────────────────────────────────────────────
#  RE-PROCESS IMAGES — extract embedded images from existing uploaded PDFs
# ─────────────────────────────────────────────────────────────────────────────

@router.post("/documents/reprocess-images")
async def reprocess_images(session_id: str = Query(...), user: CurrentUser = Depends(get_current_user)):
    """Re-extract images from all PDFs in a session. Useful for documents uploaded
    before automatic image extraction was enabled."""
    conn = get_db()
    try:
        pdf_rows = conn.execute(
            "SELECT id, file_name, file_path, doc_category FROM memo_documents "
            "WHERE owner_id=? AND memo_session_id=? AND file_type='.pdf'",
            (user.id, session_id),
        ).fetchall()

        added = 0
        for row in pdf_rows:
            if not row["file_path"]:
                continue
            fpath = Path(row["file_path"])
            if not fpath.exists():
                continue
            dest_dir = fpath.parent
            imgs = _extract_images_from_pdf(fpath, dest_dir, max_images=4)
            for img in imgs:
                try:
                    img_size = Path(img["file_path"]).stat().st_size
                    if img_size < 5000:
                        continue
                    # Don't double-insert if already there
                    exists = conn.execute(
                        "SELECT id FROM memo_documents WHERE file_path=?", (img["file_path"],)
                    ).fetchone()
                    if exists:
                        continue
                    conn.execute(
                        """INSERT INTO memo_documents
                           (owner_id, memo_session_id, file_name, file_type, file_size, extracted_text, doc_category, file_path)
                           VALUES (?,?,?,?,?,?,?,?)""",
                        (user.id, session_id, img["file_name"], ".png",
                         img_size, "", row["doc_category"], img["file_path"]),
                    )
                    added += 1
                except Exception as e:
                    logger.debug("reprocess image insert error: %s", e)
        conn.commit()
        return {"ok": True, "images_added": added, "pdfs_processed": len(pdf_rows)}
    finally:
        conn.close()


# ═══════════════════════════════════════════════════════════════════════════════
#  DEAL TERMS EXTRACTION — Extract prior deal terms from IC memo documents
# ═══════════════════════════════════════════════════════════════════════════════

@router.post("/documents/extract-deal-terms")
async def extract_deal_terms_endpoint(
    session_id: str = Query(...),
    user: CurrentUser = Depends(get_current_user),
):
    """
    Extract prior investment deal terms from uploaded prior IC memo documents.

    Looks for documents with category 'prior_ic_memo' in the given session,
    extracts text, and uses LLM to pull structured deal terms (check size,
    pre-money, ownership, stage, commitments, etc.).
    """
    from ..engine.extraction import extract_deal_terms

    conn = get_db()
    try:
        # Find prior IC memo documents
        rows = conn.execute(
            """SELECT id, file_name, extracted_text, doc_category
               FROM memo_documents
               WHERE owner_id=? AND memo_session_id=? AND doc_category='prior_ic_memo'
               ORDER BY uploaded_at""",
            (user.id, session_id),
        ).fetchall()

        if not rows:
            # Also check for term_sheet or board_materials categories as fallback
            rows = conn.execute(
                """SELECT id, file_name, extracted_text, doc_category
                   FROM memo_documents
                   WHERE owner_id=? AND memo_session_id=? AND doc_category IN ('term_sheet', 'board_materials')
                   ORDER BY uploaded_at""",
                (user.id, session_id),
            ).fetchall()

        if not rows:
            raise HTTPException(status_code=404, detail="No prior IC memo or term sheet documents found in this session")

        # Combine text from all prior IC memo documents
        combined_text = ""
        source_files = []
        for row in rows:
            text = row["extracted_text"] or ""
            if text.strip():
                combined_text += f"\n\n=== Document: {row['file_name']} (Category: {row['doc_category']}) ===\n\n{text}"
                source_files.append(row["file_name"])

        if not combined_text.strip():
            raise HTTPException(status_code=422, detail="No extractable text found in prior IC memo documents")

        # Get model preferences
        model_prefs = get_model_preferences(user.id)
        model = model_prefs.get("extraction", "claude-haiku-4-5-20251001")

        # Run extraction in thread to avoid blocking
        result = await asyncio.to_thread(
            extract_deal_terms,
            combined_text,
            source_name=", ".join(source_files),
            model=model,
        )

        result["_source_documents"] = source_files
        result["_session_id"] = session_id

        return result

    finally:
        conn.close()


# ═══════════════════════════════════════════════════════════════════════════════
#  REPORT DATA — List available reports for attaching to memos
# ═══════════════════════════════════════════════════════════════════════════════

@router.get("/reports")
async def list_reports_for_memo(
    user: CurrentUser = Depends(get_current_user),
    q: str = "",
):
    """Return deal reports across the whole team for memo selection.

    Reports are shared inside the firm — any signed-in user can pull any
    completed report into a memo. Each row carries the author's username so
    the picker can show "by joseph" / "by jack". Optional `q` filters on
    company name (case-insensitive substring)."""
    conn = get_db()
    try:
        sql = (
            "SELECT r.id, r.company_name, r.archetype, r.entry_stage, "
            "       r.status, r.created_at, r.owner_id, r.custom_title, "
            "       COALESCE(u.username, '(unknown)') AS author "
            "FROM deal_reports r "
            "LEFT JOIN users u ON u.id = r.owner_id "
        )
        params: tuple = ()
        if q.strip():
            sql += "WHERE LOWER(r.company_name) LIKE ? "
            params = (f"%{q.strip().lower()}%",)
        sql += "ORDER BY r.created_at DESC LIMIT 200"
        rows = conn.execute(sql, params).fetchall()
        return [dict(r) for r in rows]
    finally:
        conn.close()


# ═══════════════════════════════════════════════════════════════════════════════
#  MEMO GENERATION — LLM call
# ═══════════════════════════════════════════════════════════════════════════════

class MemoGenerateRequest(BaseModel):
    report_id: Optional[int] = None
    template_id: Optional[int] = None
    template_text: Optional[str] = None          # Inline template override
    session_id: Optional[str] = None             # data room document session (manual uploads)
    library_id: Optional[int] = None             # deal document library (Google Drive sync)
    additional_instructions: str = ""
    company_name: str = ""
    links: List[str] = []                        # URLs / reference links
    model_override: Optional[str] = None         # Force a specific model
    investment_type: str = "first"               # "first" or "followon"
    locked_sections: dict = {}                   # {section_key: markdown_text} — skips LLM, uses text as-is
    # ── Memo engine version ─────────────────────────────────────────────────
    # "v1" — legacy: per-doc Pass 1 extraction → per-section Pass 2 stitched
    #        from buckets → Pass 3 synthesis. Each file processed in isolation;
    #        folder structure ignored; long PDFs truncated; Excel limited to
    #        first 5 sheets.
    # "v2" — manifest-aware: single Pass 1 reads the whole data room (folder
    #        tree + previews + Excel sheet samples) and outputs a JSON manifest
    #        classifying every file. Stage 2 does targeted deep reads of
    #        primary docs (no truncation, agent-picked Excel sheets, OCR
    #        fallback). Stage 3 writes each section with the FULL data-room
    #        corpus in context, using prompt caching to keep costs bounded.
    engine_version: str = "v1"


def _format_round_terms(ov: dict) -> list:
    """Render the CURRENT round's authoritative deal terms from deal_overview.

    deal_overview is the single source of truth for round structure (current
    pre/post-money, this round's check + ownership, and a SEPARATE prior_investments
    list). Surfacing it explicitly — and labeling priors as sunk context — keeps the
    memo (especially the Pass-3 synthesis sections) anchored on the round the memo is
    FOR, instead of inferring valuation from the raw input dump which, for a follow-on,
    foregrounds the prior round and yields the wrong post-money.
    """
    if not ov:
        return []
    is_fo = ov.get("investment_type") == "followon" or ov.get("is_blended_followon")
    out = ["\n## THIS ROUND — DEAL TERMS (current round; authoritative)"]
    out.append(f"- Investment type: {'Follow-on' if is_fo else 'First check'}")
    out.append(f"- Stage: {ov.get('entry_stage', 'N/A')}")
    out.append(f"- Pre-money: ${ov.get('pre_money_millions', 'N/A')}M")
    out.append(f"- Round size: ${ov.get('round_size_millions', 'N/A')}M")
    out.append(f"- Post-money: ${ov.get('post_money_millions', 'N/A')}M")
    out.append(f"- VoLo new check (this round): ${ov.get('check_size_millions', 'N/A')}M")
    out.append(f"- Ownership from this check: {ov.get('new_check_ownership_pct', 'N/A')}%")
    if ov.get("entry_year"):
        out.append(f"- Investment (entry) year: {ov.get('entry_year')}")
    if is_fo:
        out.append(
            f"- Combined VoLo ownership after this round (priors + new): "
            f"{ov.get('combined_entry_ownership_pct', ov.get('entry_ownership_pct', 'N/A'))}%"
        )
        out.append(
            f"- Total VoLo invested across all rounds (prior + new): "
            f"${ov.get('total_invested_millions', ov.get('total_exposure_m', 'N/A'))}M"
        )
        priors = ov.get("prior_investments", []) or []
        if priors:
            out.append(
                "- PRIOR rounds (sunk context only — NOT this memo's round; never use "
                "these as the headline deal terms):"
            )
            for p in priors:
                _t = p.get("type", "priced")
                if _t in ("convertible", "safe") or p.get("cap_m"):
                    terms = f"cap ${p.get('cap_m', 'N/A')}M, discount {p.get('discount_pct', 'N/A')}%"
                else:
                    terms = f"pre-money ${p.get('effective_pre_m', 'N/A')}M"
                out.append(
                    f"    Round {p.get('round_num', '?')} "
                    f"({p.get('stage', '?')}, yr {p.get('year', '?')}, {_t}): "
                    f"VoLo ${p.get('check_m', 'N/A')}M @ {terms}, "
                    f"ownership {p.get('ownership_pct', 'N/A')}%"
                )
    out.append(
        "- NOTE: This IC memo is for the CURRENT round above — use its pre-/post-money, "
        "check size, and ownership as the headline deal terms."
    )
    if is_fo:
        out.append(
            "- FOLLOW-ON — ADDRESS BOTH: VoLo is already an investor here. The memo must cover BOTH "
            "the current round (the decision at hand) AND explicitly recognize and address VoLo's "
            "prior investment shown above — that VoLo previously invested, the combined ownership and "
            "total invested across all rounds, and how this new check builds on that existing position. "
            "Do NOT present a prior round's valuation as the current headline terms, and do NOT omit the "
            "prior investment from the memo."
        )
    return out


def _financing_summary_block(ov: dict) -> str:
    """Deterministic, clearly-separated financing summary built straight from the
    report's deal_overview, so the round numbers can never drift from the model.
    Prepended to the Financing Overview section to give an explicit prior-vs-
    current separation, on top of the financing detail woven through the memo."""
    if not ov:
        return ""

    def _m(v):
        return f"${v:,.1f}M" if isinstance(v, (int, float)) else "N/A"

    def _pct(v):
        return f"{v:.1f}%" if isinstance(v, (int, float)) else "N/A"

    def _cell(v):
        # Markdown pipe-table cells split naively on '|'; neutralise any '|' or
        # newline in free-text values (e.g. a stage label) so the table can't break.
        return str(v).replace("|", "/").replace("\n", " ").replace("\r", " ").strip()

    is_fo = ov.get("investment_type") == "followon" or ov.get("is_blended_followon")
    priors = ov.get("prior_investments", []) or []

    lines = ["### Financing Summary", ""]

    # ── Current round ──────────────────────────────────────────────────────
    lines += [
        "**Current Round — this investment**", "",
        "| Term | Value |",
        "| --- | --- |",
        f"| Stage | {_cell(ov.get('entry_stage', 'N/A'))} |",
        f"| Pre-money | {_m(ov.get('pre_money_millions'))} |",
        f"| Round size | {_m(ov.get('round_size_millions'))} |",
        f"| Post-money | {_m(ov.get('post_money_millions'))} |",
        f"| VoLo check (this round) | {_m(ov.get('check_size_millions'))} |",
        f"| Ownership from this check | {_pct(ov.get('new_check_ownership_pct'))} |",
    ]
    if ov.get("entry_year"):
        lines.append(f"| Investment year | {ov.get('entry_year')} |")
    lines.append("")

    # ── Prior financing ────────────────────────────────────────────────────
    lines += ["**Prior Financing — VoLo's existing position**", ""]
    if is_fo and priors:
        lines += [
            "| Round | Stage | Year | VoLo check | Round terms | VoLo ownership |",
            "| --- | --- | --- | --- | --- | --- |",
        ]
        for p in priors:
            _t = (p.get("type") or "priced").lower()
            if _t in ("convertible", "safe") or p.get("cap_m"):
                terms = f"cap {_m(p.get('cap_m'))}, disc {_pct(p.get('discount_pct'))}"
            else:
                terms = f"pre-money {_m(p.get('effective_pre_m'))}"
            lines.append(
                f"| {_cell(p.get('round_num', '?'))} | {_cell(p.get('stage', '?'))} | {_cell(p.get('year', '?'))} | "
                f"{_m(p.get('check_m'))} | {terms} | {_pct(p.get('ownership_pct'))} |"
            )
        lines += [
            "",
            "*Combined VoLo ownership after this round:* "
            f"**{_pct(ov.get('combined_entry_ownership_pct') or ov.get('entry_ownership_pct'))}**  ·  "
            "*Total VoLo invested across all rounds (prior + new):* "
            f"**{_m(ov.get('total_invested_millions') or ov.get('total_exposure_m'))}**",
        ]
    else:
        lines.append("*None — this is VoLo's first check into the company.*")
    lines.append("")
    return "\n".join(lines)


def _return_profile_block(report: dict) -> str:
    """Deterministic 'This Check vs Total Position' return summary, mirroring the
    deal report's dual-view header. Built from hero_metrics + total_position_forecast
    so it can't drift. THIS CHECK (return on the new capital) is the decision and the
    focus; TOTAL POSITION is the blended return on prior + new and is flagged as
    context only — it bakes in markup already accrued on the prior investment, so it
    is NOT the new check's prospective return. Keeps the prior investment visible
    without re-confusing it with the current decision. Scales match the report UI:
    MOIC = raw multiple; IRR / P(>3x) / survival = fractions (×100); ownership %s
    are already percent-valued."""
    if not report:
        return ""
    ov   = report.get("deal_overview", {}) or {}
    hero = report.get("hero_metrics", {}) or {}
    tpf  = report.get("total_position_forecast") or {}
    if not hero:
        return ""

    def _x(v):   # MOIC multiple
        return f"{v:.2f}x" if isinstance(v, (int, float)) else "N/A"

    def _p(v):   # fraction -> percent (IRR, P(>3x), survival)
        return f"{v * 100:.1f}%" if isinstance(v, (int, float)) else "N/A"

    def _pp(v):  # already a percent number (ownership)
        return f"{v:.1f}%" if isinstance(v, (int, float)) else "N/A"

    def _m(v, d=1):
        return f"${v:.{d}f}M" if isinstance(v, (int, float)) else "N/A"

    is_fo = bool(tpf) and (
        ov.get("investment_type") == "followon"
        or ov.get("is_blended_followon")
        or hero.get("basis") == "current_round"
    )
    check_m = ov.get("check_size_millions", hero.get("new_check_m"))
    own_new = ov.get("new_check_ownership_pct", hero.get("new_check_ownership_pct"))

    # Full metric set (fix 14): show expected (mean) alongside median/conditional
    # MOIC and the loss/upside probabilities so the mean is never read in
    # isolation, and append a caveat when the mean is tail-driven.
    sim       = report.get("simulation", {}) or {}
    moic_cond = sim.get("moic_conditional", {}) or {}
    prob      = sim.get("probability", sim.get("probability_buckets", {})) or {}

    def _returns_table(src, extra=False):
        rows = [
            "| Metric | Value |",
            "| --- | --- |",
            f"| Expected (mean) MOIC | {_x(src.get('expected_moic'))} |",
        ]
        if extra and isinstance(moic_cond.get("p50"), (int, float)):
            rows.append(f"| Median MOIC (P50, all paths) | {_x(moic_cond.get('p50'))} |")
        if extra and isinstance(moic_cond.get("mean"), (int, float)):
            rows.append(f"| Conditional MOIC (mean of survivors) | {_x(moic_cond.get('mean'))} |")
        rows += [
            f"| Expected IRR | {_p(src.get('expected_irr'))} |",
            f"| P(>3x) | {_p(src.get('p_gt_3x'))} |",
        ]
        if extra and isinstance(prob.get("gt_5x"), (int, float)):
            rows.append(f"| P(>5x) | {_p(prob.get('gt_5x'))} |")
        rows.append(f"| Survival rate | {_p(src.get('survival_rate'))} |")
        if extra and isinstance(prob.get("total_loss"), (int, float)):
            rows.append(f"| P(total loss) | {_p(prob.get('total_loss'))} |")
        rows.append("")
        return rows

    def _tail_caveat():
        exp_moic = hero.get("expected_moic")
        cond_p50 = moic_cond.get("p50")
        if (isinstance(exp_moic, (int, float)) and isinstance(cond_p50, (int, float))
                and cond_p50 > 0 and exp_moic >= 2.0 * cond_p50):
            p10x = prob.get("gt_10x")
            tail = f" (P(>10x) is {_p(p10x)})" if isinstance(p10x, (int, float)) else ""
            return ("*The expected (mean) MOIC is materially above the median outcome, so it is driven by "
                    f"a small number of extreme right-tail paths{tail}; the typical (median) result is lower.*")
        return ""

    lines = ["### Return Profile", ""]
    if is_fo:
        lines += [f"**This Check — the investment decision** (return on the new {_m(check_m)} at {_pp(own_new)})", ""]
        lines += _returns_table(hero, extra=True)
        _tc = _tail_caveat()
        if _tc:
            lines += [_tc, ""]
        lines += ["*The forward, decision-relevant return on the new capital deployed in this round.*", ""]
        lines += [
            f"**Total Position — all VoLo capital deployed** "
            f"(prior + new: {_m(tpf.get('total_invested_m'), 2)} at {_pp(tpf.get('combined_ownership_pct'))} combined)",
            "",
        ]
        lines += _returns_table(tpf)
        lines += [
            "*Context only — Total Position includes markup already accrued on VoLo's prior investment, "
            "so it OVERSTATES the new check's prospective return. Use This Check for the follow-on decision; "
            "Total Position shows how this check sits within VoLo's existing position.*",
            "",
        ]
    else:
        lines += [f"**Return on this investment** ({_m(check_m)} at {_pp(own_new)})", ""]
        lines += _returns_table(hero, extra=True)
        _tc = _tail_caveat()
        if _tc:
            lines += [_tc, ""]
    return "\n".join(lines)


def _verdict_token(text: str):
    """Classify a section's verdict as invest / pass / conditional / hold /
    monitor from an EXPLICIT recommendation phrase — a recommendation cue followed
    (within a short window) by the verdict word in any common form (invest /
    investing / investment, pass / passing / decline, monitor / monitoring, hold).
    Returns None when there is no explicit cue, so ordinary prose ('pass-through
    economics', 'conditional on closing', 'would not pass up') never classifies.
    Only an explicit invest-vs-pass split is treated as a contradiction downstream."""
    if not text:
        return None
    t = text.lower()
    # canonical verdict <- verb-family regex. Order matters: 'conditional' is a
    # qualifier and is checked first so "conditional invest" reads as conditional.
    fams = [
        ("conditional", r'condition(?:al|ally)?'),
        ("pass",        r'pass(?:ing|ed)?|declin(?:e|ing|ed)'),
        ("invest",      r'invest(?:ing|ed|ment|s)?|proceed(?:ing)?|participat(?:e|ing)'),
        ("monitor",     r'monitor(?:ing|ed)?|watch(?:list)?|track'),
        ("hold",        r'hold(?:ing)?'),
    ]
    cue = (r'(?:recommendation|we\s+(?:would\s+)?recommend|recommend(?:ed|ing|ation)?'
           r'|verdict|our\s+(?:call|recommendation)|decision)')
    for canon, verb in fams:
        # cue, then a small gap (connectors like "is to", "a", "against"), then verb
        if re.search(cue + r'\b[:\s*—–-]{0,4}[^.\n]{0,30}?\b(?:' + verb + r')\b', t):
            return canon
    return None


def _consistency_qa_block(report: dict, inputs: dict, section_texts: dict, memo_md: str) -> str:
    """SOFT pre-render backstop (fixes 1/4/13/14/33). Prevention is handled
    upstream (canonical facts + prompt discipline); this only LABELS anything
    that slipped through — it never blocks rendering. Reuses the existing
    qa_review engine (deterministic report checks + LLM prose review) and adds a
    couple of cross-section text checks. Returns a markdown block, or '' if clean."""
    errors = []       # genuine inconsistencies to reconcile
    warnings = []     # review items
    own_signal = False  # a targeted (non-generic) finding worth surfacing a block for

    # Fix 13 — the two synthesis surfaces must not give OPPOSITE verdicts. Only an
    # explicit invest-vs-pass split is a true contradiction; 'conditional' (and
    # hold/monitor) are qualifiers compatible with either, so they never flag.
    _rec = _verdict_token(section_texts.get("recommendation", ""))
    _ov = _verdict_token(section_texts.get("investment_overview", ""))
    if _rec and _ov and {_rec, _ov} == {"invest", "pass"}:
        errors.append(f"Recommendation mismatch — the Investment Overview implies '{_ov.upper()}' but the Recommendation section reads '{_rec.upper()}'. Reconcile to one verdict.")
        own_signal = True

    # Fix 4 — false diligence-gap phrasing (backstop; prompts prevent it).
    _placeholders = ("not provided", "no data available", "data room did not contain",
                     "not disclosed", "no information was provided")
    _hits = sorted({p for p in _placeholders if p in (memo_md or "").lower()})
    if _hits:
        warnings.append("Placeholder 'missing data' phrasing present (" + ", ".join(_hits) + "). Confirm the datum is genuinely absent; if it was extracted elsewhere, state the specific remaining gap instead.")
        own_signal = True

    # Reuse the existing QA engine (fix 14 MOIC/IRR plausibility, ownership /
    # post-money / follow-on math, + the LLM prose contradiction pass).
    try:
        from ..engine.qa_review import run_qa_review
        qa = run_qa_review(report, inputs or {}, memo_markdown=memo_md, run_llm=True)
        for f in (qa.get("findings") or []):
            sev = f.get("severity", "info")
            _txt = f"{f.get('title', '')}: {f.get('detail', '')}".strip().strip(':').strip()
            if not _txt:
                continue
            if sev == "error":
                errors.append(_txt)
            elif sev == "warning":
                warnings.append(_txt)
    except Exception as e:
        logger.warning(f"run_qa_review failed in memo QA pass: {e}")

    # Only surface a block when something actionable is present: a hard error, or
    # one of our targeted signals. A memo whose only findings are generic QA
    # warnings (e.g. an intentionally-omitted section) stays clean — prevention,
    # not noise.
    if not errors and not own_signal:
        return ""

    lines = ["## Internal QA & Consistency Review",
             "*Automated pre-render checks — internal. Resolve these before circulating the memo.*", ""]
    if errors:
        lines.append("**Unresolved inconsistencies:**")
        lines += [f"- {t}" for t in errors]
        lines.append("")
    if warnings:
        lines.append("**Warnings to review:**")
        lines += [f"- {t}" for t in warnings]
        lines.append("")
    return "\n".join(lines)


def _extract_context_section(context: str, heading_prefix: str) -> str:
    """Return the block of a report-context string whose heading line starts with
    `heading_prefix` (e.g. '## THIS ROUND'), up to the next '## ' heading.

    Used to lift a specific labeled block (like the current-round deal terms) out
    of the full report context so it can be injected where the full context is not
    passed through verbatim (e.g. the Pass-3 synthesis writer's quant filter).
    """
    if not context or not heading_prefix:
        return ""
    lines = context.split("\n")
    start = None
    for i, ln in enumerate(lines):
        if ln.strip().startswith(heading_prefix):
            start = i
            break
    if start is None:
        return ""
    out = [lines[start]]
    for ln in lines[start + 1:]:
        if ln.strip().startswith("## "):
            break
        out.append(ln)
    return "\n".join(out).strip()


def _build_report_context(report_row) -> str:
    """Build a structured context string from a deal report."""
    if not report_row:
        return ""

    report = json.loads(report_row["report_json"] or "{}")
    inputs = json.loads(report_row["inputs_json"] or "{}")
    extraction = json.loads(report_row["extraction_json"] or "{}")

    parts = [f"# DEAL REPORT: {report_row['company_name']}"]
    parts.append(f"Archetype: {report_row['archetype']}")
    parts.append(f"Entry Stage: {report_row['entry_stage']}")

    # THIS ROUND — authoritative current-round deal terms from deal_overview,
    # placed high so the memo (esp. synthesis sections) anchors on the current
    # round rather than inferring valuation from the raw input dump below (which
    # for a follow-on foregrounds the prior round → wrong post-money).
    parts.extend(_format_round_terms(report.get("deal_overview", {})))

    # ── CANONICAL FACTS (fixes 6/3/9): stamp the authoritative, single-source
    #    values so every section uses the SAME number and cannot drift between
    #    sections, and label provenance so a model output is never presented as a
    #    company/founder figure.
    _ov = report.get("deal_overview", {}) or {}
    _trl = _ov.get("trl", inputs.get("trl"))
    _trl_label = _ov.get("trl_label")
    _fm_block = report.get("financial_model", {}) or {}
    _fm_has = bool(_fm_block.get("has_data"))
    # A founder revenue curve is genuine (the company's own projection) whenever
    # founder_comparison has data — that is True both for an uploaded Excel model
    # AND for founder projections extracted from the deck (the "manual" path in
    # deal_report). Only when NO founder curve exists at all is the series a pure
    # sim baseline. Keying solely on financial_model.has_data would wrongly disavow
    # real deck-extracted founder numbers as "sim-derived" (fix 9).
    _fc_has = bool((report.get("founder_comparison") or {}).get("has_data"))
    parts.append("\n## CANONICAL FACTS — use these EXACT values in every section; do NOT recompute or vary them")
    if _trl not in (None, "", "N/A"):
        _trl_line = f"- TRL (technology readiness level): {_trl}"
        if _trl_label:
            _trl_line += f" ({_trl_label})"
        parts.append(_trl_line + " — use this TRL verbatim everywhere; do not state a different TRL in any section.")
    parts.append(f"- Company stage: {report_row['entry_stage']}")
    if _fm_has:
        parts.append('- Founder revenue model: the company uploaded a financial model; the "Founder" revenue series below IS the company\'s own uploaded projection — attribute it to the company/management.')
    elif _fc_has:
        parts.append('- Founder revenue model: no Excel model was uploaded, but the "Founder" revenue series below comes from the founders\' OWN projections in the company materials (the deck). It IS a company/management projection — attribute it to the founders/management (note it is deck-derived, not an audited model), and do NOT call it a simulation output.')
    else:
        parts.append('- Founder revenue model: NONE provided. Any "founder" revenue curve below is a MODEL baseline, NOT a company/management projection — do not attribute it to the founders or call it a founder/management forecast; label it "sim-derived baseline (no founder projection provided)".')
    parts.append("- Provenance of figures below: SIMULATION RESULTS = model-generated Monte-Carlo outputs (NOT company projections or company-reported numbers); FINANCIAL MODEL = figures extracted from the company's uploaded model; deal terms / INPUT PARAMETERS = as entered. Never present a simulation output as a founder projection or a company-reported figure.")

    # Input parameters. Skip the structured carbon/volume blobs AND the round
    # structure fields — the latter are stated authoritatively under THIS ROUND
    # above; re-dumping raw pre_money/check/prior_* here invites the synthesis
    # writer to conflate a prior round with the current one.
    _ctx_skip_input_keys = {
        'trl', 'trl_label',  # TRL is stated once, authoritatively, in CANONICAL FACTS
        'volume', 'op_carbon', 'emb_carbon', 'portfolio',
        'pre_money_millions', 'check_size_millions', 'round_size_m', 'post_money_millions',
        'investment_type', 'prior_investments',
        'prior_first_pre_money_m', 'prior_first_check_m', 'prior_first_round_size_m',
        'prior_first_entry_year', 'prior_first_entry_stage', 'prior_first_moic_distribution',
        'follow_on_is_pro_rata', 'follow_on_target_ownership_pct',
    }
    if inputs:
        parts.append("\n## INPUT PARAMETERS")
        for k, v in inputs.items():
            if v and k not in _ctx_skip_input_keys:
                parts.append(f"- {k}: {v}")

    # Simulation results — hero_metrics live at top level, probability inside simulation
    sim = report.get("simulation", {})
    hero = report.get("hero_metrics", {}) or sim.get("hero_metrics", {})
    if sim or hero:
        parts.append("\n## SIMULATION RESULTS (model-generated Monte-Carlo outputs — NOT company projections or company-reported figures)")
        if hero:
            parts.append(f"- Expected MOIC: {hero.get('expected_moic', 'N/A')}x")
            parts.append(f"- P(>3x): {hero.get('p_gt_3x', 'N/A')}")
            parts.append(f"- Expected IRR: {hero.get('expected_irr', 'N/A')}")
            parts.append(f"- Survival Rate: {hero.get('survival_rate', 'N/A')}")

        prob = sim.get("probability", sim.get("probability_buckets", {}))
        if prob:
            parts.append(f"- P(Total Loss): {prob.get('total_loss', 'N/A')}")
            parts.append(f"- P(>1x): {prob.get('gt_1x', 'N/A')}")
            parts.append(f"- P(>5x): {prob.get('gt_5x', 'N/A')}")
            parts.append(f"- P(>10x): {prob.get('gt_10x', 'N/A')}")

        moic_cond = sim.get("moic_conditional", {})
        if moic_cond:
            parts.append(f"- MOIC Conditional Mean: {moic_cond.get('mean', 'N/A')}x")
            parts.append(f"- MOIC Conditional P50: {moic_cond.get('p50', 'N/A')}x")

    # Adoption — key is adoption_analysis, curve is a dict keyed by year
    adoption = report.get("adoption_analysis", report.get("adoption", {}))
    if adoption:
        parts.append("\n## MARKET ADOPTION (model S-curve; revenue values are COMPANY revenue in $M by CALENDAR YEAR — this is company revenue, not total-market adoption, and not the same as the addressable market)")
        scurve = adoption.get("scurve", {})
        if scurve:
            parts.append(f"- Bass p (innovation): {scurve.get('bass_p_mean', 'N/A')}")
            parts.append(f"- Bass q (imitation): {scurve.get('bass_q_mean', 'N/A')}")
        # adoption_curve can be columnar {years:[...], p10:[...], p50:[...], founder:[...]}
        # OR row-based {year: {median_rev_m,...}} OR a list of row dicts
        ac = adoption.get("adoption_curve", adoption.get("divergence_table", {}))
        if isinstance(ac, dict) and ac:
            # Columnar format: has a 'years' key whose value is a list
            if isinstance(ac.get("years"), list):
                years_list = ac.get("years", [])
                p50_list   = ac.get("p50", ac.get("sim_median_m", []))
                founder_list = ac.get("founder", ac.get("founder_rev_m", []))
                parts.append("- Adoption Curve (Sim P50 by year):")
                for i, yr in enumerate(years_list[:8]):
                    p50_val      = p50_list[i]      if i < len(p50_list)      else "?"
                    founder_val  = founder_list[i]  if i < len(founder_list)  else "N/A"
                    founder_str  = f"${founder_val:.1f}M" if isinstance(founder_val, (int, float)) else "N/A"
                    p50_str      = f"${p50_val:.1f}M"     if isinstance(p50_val,     (int, float)) else str(p50_val)
                    parts.append(f"  Year {yr}: Founder {founder_str}, Sim P50 {p50_str}")
            else:
                # Row-keyed format: {year_str: {median_rev_m, ...}}
                parts.append("- Adoption Curve (Founder vs Sim by year):")
                for yr, vals in list(ac.items())[:8]:
                    if not isinstance(vals, dict):
                        continue
                    parts.append(f"  Year {yr}: Founder ${vals.get('founder_rev_m', vals.get('founder','?'))}M, "
                                 f"Sim Median ${vals.get('sim_median_m', vals.get('median','?'))}M")
        elif isinstance(ac, list) and ac:
            parts.append("- Divergence Table (Founder vs Sim):")
            for row in ac[:8]:
                parts.append(f"  Year {row.get('year','?')}: Founder ${row.get('founder_rev_m','?')}M, "
                             f"Sim Med ${row.get('sim_median_m','?')}M")

    # Carbon impact — key is carbon_impact (not carbon)
    carbon = report.get("carbon_impact", report.get("carbon", {}))
    co = carbon.get("outputs", {})
    ci_inter = carbon.get("intermediates", {})
    ci_inp = carbon.get("carbon_inputs", {})
    if co:
        parts.append("\n## CARBON IMPACT")
        if ci_inp:
            parts.append(f"- Unit definition: {ci_inp.get('unit_definition', 'N/A')}")
            parts.append(f"- Unit service life: {ci_inp.get('unit_service_life_yrs', 'N/A')} years")
            parts.append(f"- Displaced resource: {ci_inp.get('displaced_resource', 'N/A')}")
            parts.append(f"- Baseline lifetime prod: {ci_inp.get('baseline_lifetime_prod', 'N/A')} {ci_inp.get('specific_production_units', '')}")
            parts.append(f"- Range improvement: {ci_inp.get('range_improvement', 'N/A')}")
            vols = ci_inp.get('year_volumes', [])
            if vols:
                parts.append(f"- Volume forecast: {', '.join(str(round(v,1)) for v in vols[:10])}")
        if ci_inter:
            parts.append(f"- JD (displaced vol/unit): {ci_inter.get('jd', 'N/A')}")
            ann_op = ci_inter.get('annual_operating', [])
            if ann_op:
                parts.append(f"- Annual operating tCO2: {', '.join(f'{v:.1f}' for v in ann_op[:10])}")
            ann_lc = ci_inter.get('annual_lifecycle', [])
            if ann_lc:
                parts.append(f"- Annual lifecycle tCO2: {', '.join(f'{v:.1f}' for v in ann_lc[:10])}")
        parts.append(f"- Total Lifecycle tCO2: {co.get('company_tonnes', 'N/A')}")
        parts.append(f"- VoLo Pro-Rata: {co.get('volo_prorata', 'N/A')}")
        parts.append(f"- Risk Divisor: {carbon.get('risk_divisor_used', 'N/A')} ({carbon.get('risk_divisor_source', '')})")
        parts.append(f"- Risk-Adjusted: {co.get('volo_risk_adj', 'N/A')}")
        parts.append(f"- t/$ (unadjusted): {co.get('tonnes_per_dollar', 'N/A')}")
        parts.append(f"- t/$ (Risk-Adj): {co.get('risk_adj_tpd', 'N/A')}")

    # Portfolio impact
    pi = report.get("portfolio_impact", {})
    if pi:
        parts.append("\n## PORTFOLIO IMPACT")
        parts.append(f"- TVPI Base Mean: {pi.get('tvpi_base_mean', 'N/A')}x")
        parts.append(f"- TVPI New Mean: {pi.get('tvpi_new_mean', 'N/A')}x")
        parts.append(f"- TVPI Lift: {pi.get('tvpi_mean_lift', 'N/A')}x")
        parts.append(f"- IRR Base Mean: {pi.get('irr_base_mean', 'N/A')}")
        parts.append(f"- IRR New Mean: {pi.get('irr_new_mean', 'N/A')}")
        # Percentile context + Carta anchor so the memo never conflates the MEAN with the top decile.
        parts.append(
            f"- TVPI With Deal by percentile: median (P50) {pi.get('tvpi_new_p50', 'N/A')}x, "
            f"upper-quartile (P75) {pi.get('tvpi_new_p75', 'N/A')}x, "
            f"TOP DECILE (P90) {pi.get('tvpi_new_p90', 'N/A')}x"
        )
        parts.append("- Carta VC fund TVPI benchmark (mature fund): median ~1.65x, P75 ~2.2x, TOP DECILE (P90) ~4.0x")
        parts.append(
            "- READ CAREFULLY: 'TVPI Mean' above is the EXPECTED (average) fund outcome and sits BELOW the "
            "top decile by design. The like-for-like comparison to Carta's ~4x top decile is the fund's P90 — "
            "NOT the mean. Never describe the fund as below benchmark by comparing its mean to Carta's top decile."
        )

    # Sensitivity
    sens = report.get("sensitivity", {})
    tornado = sens.get("tornado", [])
    if tornado:
        parts.append("\n## SENSITIVITY ANALYSIS (Top Drivers)")
        for t in tornado[:6]:
            parts.append(f"- {t.get('param', '?')}: Low {t.get('moic_low', '?')}x / High {t.get('moic_high', '?')}x")

    # Check size optimization — key is position_sizing, optimal inside grid_search
    ps = report.get("position_sizing", report.get("check_optimization", {}))
    best = ps.get("grid_search", {}).get("optimal", ps.get("best_check", {}))
    if best:
        parts.append("\n## CHECK SIZE OPTIMIZATION")
        parts.append(f"- Optimal Check: ${best.get('check_m', 'N/A')}M")
        parts.append(f"- Implied Ownership: {best.get('ownership_pct', 'N/A')}%")
        parts.append(f"- Fund P50 Impact: +{best.get('fund_p50_pct_chg', 'N/A')}%")
        parts.append(f"- Recommended Check: ${ps.get('recommended_check_m', 'N/A')}M")

    # Follow-on sizing — present only for follow-on investments. Surfaces the
    # prior rounds so the memo sizes the NEW check against capital already in.
    fops = report.get("followon_position_sizing", {})
    if fops.get("has_data"):
        if not best:
            parts.append("\n## CHECK SIZE OPTIMIZATION")
        _pr = fops.get("prorata", {}) or {}
        _comb = fops.get("combined", {}) or {}
        _bl = _comb.get("blended_stats") or {}
        _combown = _comb.get("combined_ownership_pct", _comb.get("total_ownership_pct_approx", "N/A"))
        parts.append(f"- FOLLOW-ON INVESTMENT (#{fops.get('followon_number', '?')} follow-on; sizing basis: {fops.get('sizing_method', 'N/A')})")
        parts.append(f"  Prior capital already deployed: ${_comb.get('total_prior_m', 'N/A')}M across {_comb.get('n_prior_rounds', '?')} round(s)")
        parts.append(f"  Recommended follow-on check: ${fops.get('recommended_followon_check_m', 'N/A')}M; combined invested (prior + new): ${_comb.get('total_invested_m', 'N/A')}M")
        parts.append(f"  Diluted ownership entering round: {_pr.get('current_diluted_ownership_pct', 'N/A')}% -> combined post-round: {_combown}%")
        parts.append(f"  Pro-rata check to hold ownership flat: ${_pr.get('pro_rata_check_m', 'N/A')}M (recommendation is {_pr.get('recommended_vs_pro_rata', 'N/A')})")
        if _bl.get("blended_moic_p50") is not None:
            parts.append(f"  Blended MOIC on all capital deployed: P10 {_bl.get('blended_moic_p10')}x / P50 {_bl.get('blended_moic_p50')}x / P90 {_bl.get('blended_moic_p90')}x")
        # Two return bases. The HEADLINE is the new check (forward, decision-
        # relevant); the total-position figure is secondary context.
        tpf = report.get("total_position_forecast") or {}
        hero = report.get("hero_metrics", {}) or {}
        if tpf:
            parts.append(f"  RETURN BASIS: the headline E[MOIC] ({hero.get('expected_moic', 'N/A')}x) is the NEW CHECK evaluated standalone from this round's valuation — the forward, decision-relevant return on the capital being deployed now.")
            parts.append(f"  TOTAL-POSITION (context only): return on ALL capital deployed (prior + new, ${tpf.get('total_invested_m', 'N/A')}M at {tpf.get('combined_ownership_pct', 'N/A')}% combined) is E[MOIC] {tpf.get('expected_moic', 'N/A')}x, IRR {tpf.get('expected_irr', 'N/A')}. This includes markup already accrued on the priors, so do NOT treat it as the new check's prospective return.")

    # Financial model — pull from report_json["financial_model"]["financials"]
    fm = report.get("financial_model", {})
    fm_fins = fm.get("financials", {})
    if fm.get("has_data") and fm_fins:
        parts.append("\n## FINANCIAL MODEL (from uploaded financial model)")
        parts.append(f"- Source file: {fm.get('file_name', 'uploaded model')}")
        parts.append(f"- Scale: {fm.get('scale_info', 'N/A')}")
        if fm.get("model_summary"):
            parts.append(f"- Summary: {str(fm['model_summary'])[:300]}")
        # Revenue trajectory
        rev = fm_fins.get("revenue", {})
        if rev:
            rev_items = sorted(rev.items())
            rev_str = ", ".join(f"{yr}: ${v/1e6:.1f}M" for yr, v in rev_items if v is not None)
            parts.append(f"- Revenue by year: {rev_str}")
        # EBITDA trajectory
        ebitda = fm_fins.get("ebitda", {})
        if ebitda:
            eb_items = sorted(ebitda.items())
            eb_str = ", ".join(f"{yr}: ${v/1e6:.1f}M" for yr, v in eb_items if v is not None)
            parts.append(f"- EBITDA by year: {eb_str}")
        # Net income
        ni = fm_fins.get("net_income", {})
        if ni:
            ni_items = sorted(ni.items())
            ni_str = ", ".join(f"{yr}: ${v/1e6:.1f}M" for yr, v in ni_items if v is not None)
            parts.append(f"- Net Income by year: {ni_str}")
        # Cash
        cash = fm_fins.get("cash", {})
        if cash:
            cash_items = sorted(cash.items())
            cash_str = ", ".join(f"{yr}: ${v/1e6:.1f}M" for yr, v in cash_items if v is not None)
            parts.append(f"- Cash by year: {cash_str}")

    # Extraction data (legacy records format)
    if extraction:
        records = extraction.get("records", [])
        if records:
            parts.append("\n## FINANCIAL MODEL EXTRACTION")
            for rec in records[:20]:
                parts.append(f"- {rec.get('metric','')} ({rec.get('period','')}): {rec.get('value','')} [source: {rec.get('sheet','')}, row {rec.get('row','')}]")

    # Valuation comps — key is valuation_context
    comps = report.get("valuation_context", report.get("valuation_comps", {}))
    if comps:
        parts.append("\n## VALUATION CONTEXT")
        parts.append("- NOTE: these EV/EBITDA multiples are sector comps derived from a GENERIC public-markets dataset. When you discuss valuation, weight toward the comparables closest to THIS company's business model, customer base, margin profile, and likely exit pathway; do NOT average across unrelated sectors, and flag when the available comps are a poor match rather than blending them.")
        parts.append(f"- IPO EV/EBITDA Mean: {comps.get('ipo_ev_ebitda_mean', 'N/A')}x")
        parts.append(f"- Acquisition EV/EBITDA Mean: {comps.get('acq_ev_ebitda_mean', 'N/A')}x")
        rng = comps.get('suggested_exit_multiple_range', [])
        if rng:
            parts.append(f"- Suggested Exit Multiple Range: {rng[0]}x – {rng[1]}x")
        matches = comps.get("matches", [])
        if matches:
            parts.append("- Comparable Industries:")
            for m in matches[:5]:
                parts.append(f"  {m.get('industry_name','?')}: IPO {m.get('ipo_ev_ebitda','N/A')}x, Acq {m.get('acq_ev_ebitda','N/A')}x EV/EBITDA")

    # Risk assessment
    risk = report.get("risk_assessment", sim.get("risk_assessment", {}))
    if risk:
        parts.append("\n## RISK ASSESSMENT")
        if isinstance(risk, dict):
            for k, v in risk.items():
                if isinstance(v, str):
                    parts.append(f"- {k}: {v[:300]}")
                elif isinstance(v, list):
                    for item in v[:5]:
                        parts.append(f"- {item}")

    return "\n".join(parts)


_IMAGE_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.svg'}


def _extract_images_from_pdf(pdf_path: Path, dest_dir: Path, max_images: int = 6) -> list:
    """Extract embedded images from a PDF and save them as PNG files.

    Returns a list of dicts with 'file_name' and 'file_path' for each saved image.
    Uses pypdf's image extraction + Pillow for format conversion (handles jp2, etc.).
    No poppler required. Limits to max_images to avoid flooding the data room.
    """
    saved = []
    try:
        import pypdf
        from PIL import Image as PILImage
        reader = pypdf.PdfReader(str(pdf_path))
        img_index = 0
        for page_num, page in enumerate(reader.pages):
            if img_index >= max_images:
                break
            try:
                for image_obj in page.images:
                    if img_index >= max_images:
                        break
                    try:
                        raw_data = image_obj.data
                        if len(raw_data) < 5000:   # skip tiny icons/artifacts
                            continue
                        img_name = f"{pdf_path.stem}_p{page_num+1}_img{img_index+1}.png"
                        img_path = dest_dir / img_name
                        # Convert to PNG via Pillow (handles jp2, bmp, tiff, etc.)
                        pil_img = PILImage.open(io.BytesIO(raw_data))
                        if pil_img.mode not in ("RGB", "RGBA"):
                            pil_img = pil_img.convert("RGB")
                        pil_img.save(str(img_path), format="PNG")
                        saved.append({"file_name": img_name, "file_path": str(img_path)})
                        img_index += 1
                    except Exception:
                        pass
            except Exception:
                pass
    except Exception as e:
        logger.debug("PDF image extraction failed for %s: %s", pdf_path.name, e)
    return saved

def _load_raw_documents(session_id: str, owner_id: int) -> list:
    """Load all uploaded documents for a session.
    - Text-bearing files (PDF, DOCX, PPTX, etc.) are included when extracted_text is non-empty.
    - Image files (.png/.jpg/etc.) are always included even if they have no extracted text,
      so they can be auto-placed inline in the rendered memo.
    """
    conn = get_db()
    try:
        rows = conn.execute(
            "SELECT id, file_name, file_type, doc_category, extracted_text FROM memo_documents "
            "WHERE owner_id=? AND memo_session_id=? ORDER BY doc_category, uploaded_at",
            (owner_id, session_id),
        ).fetchall()
        result = []
        for r in rows:
            ext = (r["file_type"] or "").lower()
            has_text = bool((r["extracted_text"] or "").strip())
            is_image = ext in _IMAGE_EXTS
            if has_text or is_image:
                result.append(dict(r))
        return result
    finally:
        conn.close()


# ═══════════════════════════════════════════════════════════════════════════════
#  MEMO SECTION DEFINITIONS — the canonical structure
# ═══════════════════════════════════════════════════════════════════════════════

MEMO_SECTIONS = [
    {
        "key": "investment_overview",
        "title": "Investment Overview",
        "is_synthesis": True,
        "guidance": (
            "One-liner summary of the deal. Then a quick-reference table with: VEV Focus Area, "
            "Investment Stage & Vehicle, Key Terms, Capital Amount, Expected VEV Equity, Board Participation. "
            "Follow with 'Portfolio Themes' — a short bulleted list of why this fits VoLo's thesis "
            "(e.g. low cost next-gen energy storage, rapidly scalable, system-level enabling tech). "
            "For a follow-on, the Key Terms / Capital Amount / Expected VEV Equity must reflect the "
            "CURRENT round, and the overview must also state that VoLo is an existing investor — prior "
            "capital invested and combined ownership across rounds. "
            "This section is written LAST after all other sections are complete."
        ),
    },
    {
        "key": "high_level_opportunities",
        "title": "High Level Opportunities",
        "is_synthesis": True,
        "guidance": (
            "3-6 bullet points capturing the strongest reasons to invest. "
            "Each bullet should be a concise, compelling statement backed by data from the memo sections. "
            "Think: team pedigree, market tailwinds, technical moat, unit economics, carbon impact. "
            "Written LAST from all completed sections."
        ),
    },
    {
        "key": "high_level_risks",
        "title": "High Level Risks",
        "is_synthesis": True,
        "guidance": (
            "3-6 bullet points capturing the most important risks and concerns. "
            "Each bullet should be specific and actionable (not vague). Reference quantified risks "
            "from the sensitivity analysis where possible. Include technology, market, execution, "
            "financing, and regulatory risks. Written LAST from all completed sections."
        ),
    },
    {
        "key": "company_overview",
        "title": "Company Overview",
        "is_synthesis": False,
        "guidance": (
            "Set a clear tone: What is the company's 'Aha' or 'breakthrough'? "
            "Contextualize the opportunity — why does the market need this, how big is the market, "
            "what is the gap, and why now? Why this categorical approach to the problem? What bottlenecks "
            "do they solve? Describe differentiation. Validate with traction and defensibility (high level — "
            "tech details go in Traction). Entice: preview 'how big this can get'. Insert team background, "
            "market size, or notable channel partners for scale if relevant. Reinforce vision."
        ),
        "report_fields": ["inputs"],
    },
    {
        "key": "market",
        "title": "Market",
        "is_synthesis": False,
        "guidance": (
            "Start broad with total demand and market size (TAM). Describe market dynamics and current "
            "composition by major segments. Then narrow to serviceable market (SAM/SOM) — where does this "
            "company have a right to win? Reference the Bass diffusion S-curve model parameters and adoption "
            "timeline from the RVM. Compare founder TAM claims to independent sources."
        ),
        "report_fields": ["inputs", "adoption"],
    },
    {
        "key": "business_model",
        "title": "Business Model",
        "is_synthesis": False,
        "guidance": (
            "How does the company make money? Cover unit economics in detail. Discuss scalability — "
            "map of capex needs and financing plans. Describe the broad go-to-market strategy. "
            "Highlight relevant strengths from the RVM (include RVM output image reference if available). "
            "Cover forward projections — compare founder projections vs simulation median, flag divergences, "
            "include divergence table data."
        ),
        "report_fields": ["extraction", "inputs", "simulation", "adoption"],
    },
    {
        "key": "financials",
        "title": "Financials",
        "is_synthesis": False,
        "guidance": (
            "Present the financial projections and RVM simulation results side-by-side. "
            "Describe the revenue model trajectory: S-curve adoption calibration (Bass p, q parameters), "
            "inflection year, peak adoption year. Compare founder revenue projections against the RVM "
            "simulation median and P25/P75 band — identify years where founder assumptions diverge "
            "significantly from the model and explain the drivers. Summarize key financial metrics: "
            "gross margins, burn rate, runway, unit economics from the financial model if provided. "
            "Discuss capital efficiency and path to profitability. Reference the revenue cone chart and "
            "founder comparison table (injected below)."
        ),
        "report_fields": ["extraction", "simulation", "adoption"],
    },
    {
        "key": "team",
        "title": "Team",
        "is_synthesis": False,
        "guidance": (
            "Quick overview of relevant background — operating + founding experience, combined skills of "
            "founding team. Discuss execution to date in relation to VoLo's experience with the team. "
            "Provide additional signals that they will be able to execute against a venture trajectory/return. "
            "Include relevant factors from the RVM. Key hires needed."
        ),
        "report_fields": [],
    },
    {
        "key": "traction",
        "title": "Traction",
        "is_synthesis": False,
        "guidance": (
            "Organize into three subsections: Technical traction (TRL level, product milestones, IP), "
            "Commercial traction (revenue, pipeline, LOIs, partnerships), and Customer Feedback "
            "(references, NPS, testimonials). Be specific with data points."
        ),
        "report_fields": ["extraction"],
    },
    {
        "key": "competitive_position",
        "title": "Competitive Position",
        "is_synthesis": False,
        "guidance": (
            "Provide a broad framework to categorize competition — zoom out then zoom in. "
            "Quick overview on why/where other categories are limited. Include incumbents where relevant. "
            "If the company provided a competitive landscape slide, lift and contextualize it. "
            "Discuss barriers to entry and moat sustainability."
        ),
        "report_fields": ["inputs"],
    },
    {
        "key": "carbon_impact",
        "title": "Carbon Impact",
        "is_synthesis": False,
        "guidance": (
            "Write a detailed, quantitative Carbon Impact section structured in five parts:\n\n"
            "1. THEORY OF CHANGE: In 1-2 paragraphs explain the physical mechanism of emissions "
            "displacement. What resource does this technology replace (e.g. natural gas, diesel, "
            "grid electricity)? What is the baseline carbon intensity of that resource? "
            "How does each unit of the company's product displace emissions — describe the chain "
            "from unit deployed → lifetime production → displaced resource intensity → tCO₂/unit. "
            "Note that VoLo's mitigation definitions align with GIIN IRIS+ metrics and are designed "
            "to capture the full life-cycle activities contributing impact materiality and financial "
            "materiality per the EU CSRD double materiality framework.\n\n"
            "2. DISPLACEMENT CHAIN CALCULATION: Walk through the math step by step with actual numbers "
            "from the RVM: (a) Displaced volume per unit (JD) = range_improvement × baseline_lifetime_prod. "
            "(b) Carbon intensity at Year 1 (JE) — state the value and note if grid intensity declines. "
            "(c) Annual operating impact = JD × units_deployed × CI_year for each of the 10 forecast years. "
            "(d) Embodied carbon (manufacturing/upstream) — whether modeled and its magnitude vs operating. "
            "(e) Annual lifecycle = operating + embodied. Cite actual tCO₂ values from the report. "
            "This methodology follows the Project Frame Pre-Investment Considerations (PRIME ERP).\n\n"
            "3. 10-YEAR CARBON SUMMARY WITH RVM METRICS: Present all four output metrics with exact values "
            "using VoLo's published terminology from the annual Impact Report: "
            "(a) Total Potential (10yr) — company-level total potential carbon impact over the 10-year horizon. "
            "(b) Fund Pro-Rata Share — attributed per the GHG Protocol equity-method: total × VoLo ownership%%. "
            "(c) P-50 Risk-Adjusted — pro-rata tCO₂ weighted by likelihood-of-success curves fitted to the "
            "company's technical and commercial risk (TRL modifiers applied to Carta benchmark graduation rates "
            "across 5,000 Monte Carlo paths). This is the expected avoided emissions, not a worst-case haircut. "
            "(d) TCPI (Tonnes CO₂ to Paid-In Capital) — P-50 risk-adjusted tCO₂ ÷ VoLo check size. "
            "This is VoLo's primary cross-deal carbon efficiency KPI. Note the survival rate applied and "
            "what TRL milestone would improve it.\n\n"
            "4. PORTFOLIO CONTEXT: Compare TCPI to typical VoLo range (~0.01–0.10 t CO₂/$) as published "
            "in the VoLo Impact Report (fund-level TCPI >2 tonnes per dollar). "
            "Discuss sensitivity: carbon impact at 50%% volume forecast vs base case. "
            "Identify the commercial milestone that would meaningfully improve the survival rate.\n\n"
            "5. CLIMATE THESIS ALIGNMENT: Connect to VoLo's theory of change — that superior risk-adjusted "
            "returns trigger profit-seeking capital, accelerating climate solutions at scale. "
            "Is the carbon impact front-loaded (embodied at manufacturing) or back-loaded (cumulative deployment)? "
            "How does this compare to other technologies in VoLo's portfolio? "
            "Reference the annual impact chart and TCPI attribution waterfall below."
        ),
        "report_fields": ["carbon"],
    },
    {
        "key": "technology_ip_moat",
        "title": "Technology: Overview, IP and Moat",
        "is_synthesis": False,
        "guidance": (
            "Technical brief including: differentiation (how the tech works, why it's better), "
            "what is derisked to date, derisking ahead — remaining categorical step changes and how "
            "VoLo got comfortable with these. Discuss IP/patent landscape and freedom to operate. "
            "Highlight relevant strengths from the RVM."
        ),
        "report_fields": ["inputs"],
    },
    {
        "key": "financing_overview",
        "title": "Financing Overview: Round and Exit Planning",
        "is_synthesis": False,
        "guidance": (
            "PRIOR FINANCING vs CURRENT ROUND: A deterministic 'Financing Summary' with separate "
            "Current Round and Prior Financing tables is injected at the top of this section — do "
            "not restate those raw numbers. In the narrative, clearly DISTINGUISH (a) the CURRENT "
            "round being decided here from (b) VoLo's PRIOR financing. For a follow-on, explicitly "
            "acknowledge VoLo is an existing investor and explain how this round builds on the prior "
            "position (combined ownership, total invested); never conflate a prior round's valuation "
            "with the current round's terms.\n\n"
            "RETURN PROFILE: A deterministic 'Return Profile' block (This Check vs Total Position) is "
            "also injected — do not restate its numbers. Analyse the investment primarily through THIS "
            "CHECK, the return on the new capital being deployed (the actual decision). Treat TOTAL "
            "POSITION as context only: it is the blended return on prior + new capital and includes "
            "markup already accrued on the prior investment, so it overstates the new check's "
            "prospective return. Always keep the two distinct.\n\n"
            "DEAL STRUCTURE: Overview of the round and VoLo's role. Include total raise. "
            "Describe VoLo Earth value-add and additionality. List committed syndicate members.\n\n"
            "CAPITAL POSITION: Total raised to date, current cash in bank, current burn rate, "
            "runway provided by this financing.\n\n"
            "ROUND PRICING: Price, implied multiple (TTM x NTM), comps discussion with public comps "
            "chart (describe categories of comps and rationale with EV/Revenue and weighting). "
            "Include private comps.\n\n"
            "EXIT PLANNING: Basic MOIC with revenue multiple range x discounted revenue projection in "
            "3-5 years. Discuss likely acquirers and strategic positioning. IPO potential if applicable. "
            "Downside protection. Include RVM Exit screenshot reference and assumptions on VoLo check, "
            "subsequent rounds, EBITDA.\n\n"
            "KEY MILESTONES & USE OF FUNDS: High-level milestones to hold company accountable for in "
            "follow-on round. Brief discussion of feasibility.\n\n"
            "GROWTH INVESTOR INSIGHTS: What would make this attractive to growth-stage investors?"
        ),
        "report_fields": ["inputs", "simulation", "valuation_comps", "check_optimization", "portfolio_impact"],
    },
    {
        "key": "cohort_analysis",
        "title": "Cohort Analysis: Survival, Dilution & Risk Calibration",
        "is_synthesis": False,
        "guidance": (
            "This section documents the probabilistic risk model that underlies VoLo's survival-adjusted "
            "carbon and financial metrics. It is derived entirely from Carta benchmark data and TRL "
            "readiness modifiers — NOT part of the PRIME ERP carbon methodology. Write in a clear, "
            "explanatory style aimed at an LP or co-investor who wants to understand how VoLo quantifies "
            "venture risk. Structure as follows:\n\n"
            "1. RISK MODEL OVERVIEW: In 1-2 paragraphs explain that VoLo runs a Monte Carlo simulation "
            "(5,000 paths) that combines stage-by-stage venture survival probabilities from Carta "
            "benchmark data with TRL-based modifiers. The result is a survival rate (probability this "
            "company reaches a meaningful exit) and an outcome distribution across total loss, partial "
            "return, 1-3× MOIC, 3-10× MOIC, and 10×+ MOIC outcomes.\n\n"
            "2. CARTA BASELINE GRADUATION RATES: Present the stage-by-stage graduation rates from Carta "
            "benchmark data for the relevant sector. For each stage (Pre-Seed through Series D), report: "
            "graduation rate (probability of advancing to next stage), median round size, median pre-money "
            "valuation, and median months to graduation. Explain that these represent market-wide venture "
            "outcomes and serve as the baseline before TRL adjustment.\n\n"
            "3. TRL MODIFIER TABLE: Present the full TRL 1-9 modifier lookup table with four parameters "
            "for each TRL level: survival_penalty (additional annual probability of failure), "
            "capital_multiplier (extra dilution from additional funding rounds needed), "
            "extra_bridge_prob (probability of requiring a bridge round at each stage), and "
            "exit_multiple_discount (haircut to exit multiple at each TRL level). Highlight this "
            "company's current TRL level. Explain the decay formula: decay(s) = max(0, 1 - s × 0.3) "
            "where s is the survival_penalty — this is applied multiplicatively at each stage transition.\n\n"
            "4. EFFECTIVE SURVIVAL FUNNEL: Show how the TRL modifier transforms raw Carta graduation "
            "rates into this company's effective stage-by-stage survival probability. For each stage, "
            "show: Carta baseline rate, TRL modifier applied (decay factor), effective rate after "
            "adjustment, and cumulative survival probability from entry. State the final overall "
            "survival_rate and meaningful_exit_rate as headline numbers. Note that this survival rate "
            "is used ONLY in the VoLo Risk-Adjusted tCO₂ metric — it is NOT part of the PRIME ERP "
            "methodology.\n\n"
            "5. OUTCOME DISTRIBUTION: Describe the simulated outcome distribution across the five "
            "buckets. Explain what drives the shape — at low TRL, total_loss dominates; at higher TRL "
            "the distribution shifts toward moderate returns. Reference the outcome distribution chart "
            "injected below.\n\n"
            "6. KEY SENSITIVITIES: How would the survival rate change if TRL advances by 2 levels? "
            "What is the most capital-intensive stage transition and why? How does the sector profile "
            "affect outcomes compared to the broader Carta universe?"
        ),
        "report_fields": ["simulation"],
    },
    {
        "key": "fund_return_model",
        "title": "Fund Return Model",
        "is_synthesis": False,
        "guidance": (
            "This section covers VoLo's fund-level return modeling for this deal. Structure it in two parts:\n\n"
            "PART 1 — PORTFOLIO IMPACT (RVM Section 3): "
            "Describe the marginal impact of adding this deal to the VoLo fund. "
            "Present the simulation methodology: the VCSimulator runs 2,000 portfolio paths, "
            "comparing the base portfolio (with or without committed deals) against a portfolio "
            "that includes this deal. Report the delta in TVPI Mean, TVPI P50, TVPI P75, and IRR "
            "at mean and P50. ALSO report the TOP DECILE (P90) fund TVPI and benchmark it against the "
            "Carta top decile (~4x): the fund's P90 is the like-for-like top-decile comparison, while the "
            "TVPI Mean is the EXPECTED case and sits below the top decile by design — never describe the "
            "fund as below benchmark by comparing its mean to Carta's top decile. "
            "Interpret what the lift means in practical terms — does this deal "
            "diversify the portfolio, add return potential, or both? Discuss the deal parameters "
            "used in the simulation: conditional MOIC, survival probability, exit year range. "
            "Flag if the baseline is a running portfolio (with committed deals) vs. a simulated "
            "benchmark portfolio.\n\n"
            "PART 2 — CHECK SIZE OPTIMIZATION (RVM Section 8): "
            "Explain the fund-performance optimizer methodology: a $250K-increment sweep across "
            "the fund-constrained check size range, with VCSimulator run at each level to measure "
            "percentage change in fund TVPI at P10, P50, and P90. Describe the composite score "
            "(weighted sum of normalized Δ%P10/P50/P90 with stage-calibrated weights). "
            "State the fund-optimized check size and implied ownership. "
            "Compare to Kelly Criterion (full and half-Kelly) as a benchmark. "
            "Discuss fund concentration constraints. "
            "Reference the Fund TVPI Impact and Composite Score charts injected below. "
            "Conclude with a recommendation on check sizing given VoLo's portfolio construction "
            "goals and the deal terms."
        ),
        "report_fields": ["portfolio_impact", "check_optimization", "simulation"],
    },
    {
        "key": "portfolio_tracking_scorecard",
        "title": "Portfolio Tracking Scorecard",
        "is_synthesis": False,
        "followon_only": True,
        "guidance": (
            "This section is ONLY included for follow-on investments. It evaluates how the company has "
            "performed against the goals, milestones, and commitments made in the original investment memo "
            "and subsequent board reporting.\n\n"
            "STRUCTURE:\n\n"
            "1. SCORECARD TABLE: Build a structured scorecard with these columns:\n"
            "   - Commitment/Goal (what was promised in the original IC memo)\n"
            "   - Status: one of MET, PARTIALLY MET, DELAYED, CHANGED, MISSED\n"
            "   - Evidence (specific data points from board reports)\n"
            "   - Commentary (brief IC-level interpretation)\n\n"
            "Categories to evaluate (group rows by category):\n"
            "   a) FINANCIAL MILESTONES — Revenue targets, burn rate, runway, unit economics, fundraising progress\n"
            "   b) PRODUCT & TECHNOLOGY — Product roadmap milestones, TRL progression, IP/patent filings, key technical achievements\n"
            "   c) COMMERCIAL TRACTION — Customer acquisition, pipeline, LOIs, partnerships, market penetration\n"
            "   d) TEAM & ORGANIZATION — Key hires, org build-out, board composition, advisory additions\n"
            "   e) CARBON IMPACT — tCO2 milestones, pilot deployments, measurement/verification progress\n"
            "   f) USE OF FUNDS — How capital was actually deployed vs. stated plan\n\n"
            "2. OVERALL SCORE: Provide a high-level summary score:\n"
            "   - Calculate percentage MET or PARTIALLY MET vs total commitments\n"
            "   - Assign an overall grade: STRONG EXECUTION (>80%%), ON TRACK (60-80%%), CAUTION (<60%%)\n"
            "   - Highlight the most important positive and negative deviations\n\n"
            "3. TRAJECTORY NARRATIVE: 2-3 paragraphs analyzing:\n"
            "   - Pattern of execution: Are delays systemic or isolated?\n"
            "   - Pivot vs. plan: If goals changed, was the pivot strategic and well-reasoned or reactive?\n"
            "   - Velocity trend: Is the company accelerating, steady, or decelerating?\n"
            "   - Management credibility: How reliable have management projections been?\n\n"
            "4. IMPLICATIONS FOR FOLLOW-ON:\n"
            "   - Does the execution track record support the follow-on valuation?\n"
            "   - What new milestones should VoLo require before the next investment decision?\n"
            "   - Are there governance or reporting gaps to address?\n\n"
            "IMPORTANT: Extract specific commitments from the PRIOR IC MEMO documents (category: prior_ic_memo). "
            "Cross-reference against BOARD REPORT documents (category: board_report). "
            "If either source is missing, note the gap explicitly. Never fabricate milestones."
        ),
        "report_fields": ["inputs", "simulation"],
    },
    {
        "key": "recommendation",
        "title": "Investment Recommendation",
        "is_synthesis": True,
        "guidance": (
            "Clear recommendation: Invest / Pass / Conditional. Summarize the bull case and bear case. "
            "State proposed terms (check size, board seat, key conditions). Quantitative highlights: "
            "MOIC, IRR, P(>3x), carbon t/$, portfolio impact. Identify follow-up diligence items. "
            "This section is written LAST after all other sections."
        ),
    },
]

# Section keys for non-synthesis (written in pass 2) vs synthesis (written in pass 3)
_DATA_SECTIONS = [s for s in MEMO_SECTIONS if not s["is_synthesis"]]
_SYNTHESIS_SECTIONS = [s for s in MEMO_SECTIONS if s["is_synthesis"]]


def _get_active_sections(investment_type: str = "first"):
    """Return (data_sections, synthesis_sections) filtered by investment type."""
    if investment_type == "followon":
        return _DATA_SECTIONS, _SYNTHESIS_SECTIONS
    # For first investments, exclude follow-on only sections
    data = [s for s in _DATA_SECTIONS if not s.get("followon_only")]
    return data, _SYNTHESIS_SECTIONS


# ═══════════════════════════════════════════════════════════════════════════════
#  WRITING STYLE GUIDE — derived from high-quality reference memos
# ═══════════════════════════════════════════════════════════════════════════════

_STYLE_GUIDE = """
WRITING STYLE — emulate the style of a top-tier venture capital diligence memo:

VOICE & TONE:
- Write with conviction and intellectual authority, but remain grounded in evidence
- Use a confident, declarative voice — avoid hedging language like "it seems", "perhaps", "it could be argued"
- Be direct about both strengths AND weaknesses; credibility comes from balanced honesty, not cheerleading
- Maintain a tone that is professional yet engaging — this should read like a compelling narrative, not a dry report
- Write for a sophisticated but diverse audience: IC members may include deep technologists, financial experts, and generalist partners

NARRATIVE STRUCTURE:
- Open each section with a strong orienting statement that frames WHY this topic matters for the investment thesis
- Build from context to specifics: set the macro landscape first, then zoom into the company's specific position
- Connect every technical or market fact back to its investment implication — never leave data uninterpreted
- Use transitional sentences that connect ideas across paragraphs; the memo should flow as a coherent story
- Close sections with a forward-looking statement about what this means for the deal or remaining diligence

TECHNICAL DEPTH WITH ACCESSIBILITY:
- Explain complex technical concepts clearly enough for a non-specialist, then immediately layer on the expert-level detail
- Use parenthetical clarifications to define jargon inline: "quasi-axisymmetric (QA) stellarator", "technology readiness level (TRL)"
- When describing technology, always connect the mechanism to its commercial advantage: "planar magnets... making the system modular, resilient, and far easier to scale"
- Quantify everything possible: use specific numbers, percentages, dollar amounts, timelines, and comparisons
- Use analogies sparingly but effectively to make complex concepts intuitive

DATA-DRIVEN ARGUMENTATION:
- Lead with the strongest quantitative evidence — dollar amounts, growth rates, probability metrics, multiples
- Cite specific data points with attribution: "Princeton ZERO Lab analysis suggests...", "According to the company's financial model..."
- Frame quantitative findings in context: not just "MOIC of 4.2x" but what that means relative to the fund's targets
- Present ranges and scenarios rather than single-point estimates where uncertainty exists
- Include third-party validation wherever available: peer-reviewed research, government data, independent analyses

COMPETITIVE & MARKET FRAMING:
- Position the company within its competitive landscape explicitly — don't just describe in isolation
- Use direct, specific comparisons: "Unlike CFS, which requires complex 3D magnets, Thea uses planar coils..."
- Acknowledge the competitive strengths of alternatives before explaining why the company's approach is differentiated
- Frame market dynamics as investment tailwinds or headwinds with supporting evidence

RISK & OBJECTIVITY:
- Identify specific, concrete risks — avoid vague "the market could change" statements
- For each risk, assess severity and proximity, then propose a specific mitigant or monitoring plan
- Use language like "Gaps / Concerns" and "While [strength], [counterpoint]" to signal balanced analysis
- Clearly distinguish between risks that are company-specific vs. industry-wide
- Note what the company CANNOT yet prove and what additional diligence would be needed

FORMATTING DISCIPLINE:
- Use prose paragraphs for analytical narrative; use bullet lists only for catalogs of discrete items (customer types, milestones, risk factors)
- Bold key terms and company names for scannability but don't overdo it
- Keep paragraphs focused on a single idea — typically 3-5 sentences
- Subheadings should be descriptive and create a scannable outline
"""


# ═══════════════════════════════════════════════════════════════════════════════
#  PASS 1 — DOCUMENT EXTRACTION: read each doc and allocate facts to sections
# ═══════════════════════════════════════════════════════════════════════════════

_EXTRACTION_SYSTEM = """You are a data room analyst for VoLo Earth Ventures. Your job is to read a document
and extract the key facts, figures, and insights that are relevant to an investment memo.

You will be given:
1. A document from the data room
2. A list of memo sections

For EACH memo section, extract the relevant information from this document using the format below.
If the document has nothing relevant to a section, write "No relevant information in this document."

EXTRACTION GUIDELINES — what to pull and how to tag it:

1. QUANTITATIVE DATA (always include):
   - Specific numbers: revenue ($12.3M ARR), growth rates (127% YoY), margins (62% gross)
   - Dates and timelines: "first commercial plant targeted for 2030"
   - Comparative metrics: "3x faster than incumbent approach", "$250M NOAK vs $1.24B FOAK"
   - Financial terms: round size, valuation, ownership, check size, liquidation preferences

2. ATTRIBUTED CLAIMS (tag the source):
   - Who says it and how credible: "Per Princeton ZERO Lab techno-economic analysis..."
   - Third-party validation vs company claims: mark "[company claim]" vs "[independent]" vs "[peer-reviewed]"
   - Customer/partner evidence: "LOI from [utility name]", "pilot with [partner]"

3. INVESTMENT IMPLICATIONS (brief annotation):
   - After key facts, add a short "→" annotation on what it means for the deal
   - Example: "Magnet production at 1/day, line of sight to 1/hour → path to mass manufacturing validated"
   - Example: "D-D fuel for pilot, D-T for commercial → fuel switch remains an undeRisked engineering step"

4. COMPETITIVE CONTEXT (when present):
   - Direct comparisons to named competitors or alternative approaches
   - What makes this approach better/worse/different and why
   - Industry benchmarks or standard performance metrics

5. RISKS & GAPS (flag explicitly):
   - Anything the document acknowledges as unproven, untested, or uncertain
   - Missing data that would be expected (e.g. no unit economics provided)
   - Assumptions that underpin projections (tag as "[assumption]")

Output format — use exactly these section headers:

[SECTION_KEY]
- Fact or data point with specifics [source tag if applicable]
  → Investment implication (one line)
- Another fact with numbers, names, dates
- Risk/gap: what's missing or unproven [flag]

CRITICAL — ACCURACY:
- Extract ONLY facts, names, numbers, and claims that are EXPLICITLY stated in the document.
- NEVER infer, assume, or fabricate details that are not present in the text.
- If a facility, location, partner, or customer is named in the document, include it exactly as written.
- If the document does not name something specifically, do NOT guess or fill in a name.
- Quote specific language from the document when possible to preserve accuracy.

Repeat for each section. Be thorough but avoid repetition. Extract the RAW MATERIAL that a
skilled memo writer needs to construct a compelling, data-driven narrative."""

_EXTRACTION_SECTIONS_LIST = "\n".join(
    f"- [{s['key']}] {s['title']}: {s['guidance'][:200]}"
    for s in _DATA_SECTIONS
)


def _is_fatal_llm_error(e: Exception) -> bool:
    """Return True for errors that will affect every subsequent call (credit/auth issues)."""
    msg = str(e).lower()
    return (
        "credit balance exhausted" in msg
        or "credit_balance_exhausted" in msg
        or "402" in str(e)
        or "invalid_api_key" in msg
        or "authentication_error" in msg
        or "401" in str(e)
    )


def _pass1_extract_document(client, model: str, doc: dict, max_chars: int = 60_000) -> dict:
    """
    Pass 1: Send a single document to the LLM and get back per-section fact extractions.
    Returns {section_key: "extracted bullet points..."}.
    """
    text = doc["extracted_text"][:max_chars]
    cat = doc["doc_category"].replace('_', ' ').title()
    fname = doc["file_name"]

    user_msg = f"""# DOCUMENT: [{cat}] {fname}
(File size: {len(doc['extracted_text'])} chars)

{text}

---

# MEMO SECTIONS TO EXTRACT INTO:
{_EXTRACTION_SECTIONS_LIST}

Extract all relevant information from this document into the appropriate sections.
Include quantitative data, attributed claims, investment implications, competitive context, and risk flags."""

    try:
        response = client.messages.create(
            model=model,
            max_tokens=6000,
            system=_EXTRACTION_SYSTEM,
            messages=[{"role": "user", "content": user_msg}],
        )
        reply = "".join(b.text for b in response.content if b.type == "text")
        tokens_in = response.usage.input_tokens if response.usage else 0
        tokens_out = response.usage.output_tokens if response.usage else 0
    except Exception as e:
        logger.error(f"Pass 1 extraction failed for {fname}: {e}")
        if _is_fatal_llm_error(e):
            raise
        return {"_error": str(e), "_tokens_in": 0, "_tokens_out": 0}

    # Parse the reply into section buckets
    sections = {}
    current_key = None
    current_lines = []
    for line in reply.split('\n'):
        stripped = line.strip()
        # Check for section header like [company_overview] or [COMPANY_OVERVIEW]
        bracket_match = re.match(r'^\[(\w+)\]', stripped)
        if bracket_match:
            if current_key and current_lines:
                sections[current_key] = '\n'.join(current_lines)
            current_key = bracket_match.group(1).lower()
            rest = stripped[bracket_match.end():].strip()
            current_lines = [rest] if rest else []
        elif current_key:
            current_lines.append(line)

    if current_key and current_lines:
        sections[current_key] = '\n'.join(current_lines)

    sections["_tokens_in"] = tokens_in
    sections["_tokens_out"] = tokens_out
    sections["_source"] = fname
    return sections


def _build_citation_index(raw_docs: list) -> dict:
    """
    Build a numbered citation index from uploaded documents.
    Returns {filename: citation_number} and a formatted legend string.
    """
    index = {}
    legend_parts = []
    for i, doc in enumerate(raw_docs, 1):
        fname = doc["file_name"]
        cat = doc["doc_category"].replace('_', ' ').title()
        index[fname] = i
        legend_parts.append(f"[{i}] [{cat}] {fname}")
    legend = "\n".join(legend_parts)
    return {"map": index, "legend": legend, "count": len(raw_docs)}


def _aggregate_section_briefs(all_extractions: list, section_key: str,
                               citation_index: dict = None) -> str:
    """Combine extractions from all documents for a single section, with citation numbers."""
    parts = []
    cite_map = citation_index.get("map", {}) if citation_index else {}
    for ext in all_extractions:
        source = ext.get("_source", "Unknown")
        text = ext.get(section_key, "").strip()
        if text and "no relevant information" not in text.lower():
            cite_num = cite_map.get(source, "?")
            parts.append(f"### Source [{cite_num}]: {source}\n{text}")
    return "\n\n".join(parts) if parts else ""


# ═══════════════════════════════════════════════════════════════════════════════
#  PASS 2 — SECTION WRITING: write each section from its aggregated brief
# ═══════════════════════════════════════════════════════════════════════════════

_SECTION_WRITER_SYSTEM = f"""You are VoLo Earth Ventures' Investment Committee memo writer.
You are writing ONE section of an investment memorandum.

CRITICAL — FACTUAL ACCURACY:
- ONLY state facts, names, numbers, dates, and claims that appear explicitly in the provided source documents or report data.
- NEVER fabricate PERSON names (founders, team members, advisors, customers), COMPANY names (competitors, partners, acquirers, investors, comparable startups), facility names, locations, dollar amounts, percentages, timelines, or partnerships. If a specific person or company is not named in the sources for the subject company, do not name one — write "the founding team" / "incumbent competitors" generically instead.
- Every named person and company must trace to the SUBJECT company's own source documents (the company named in the user message). If you find yourself introducing a person or organization that is not in those sources, STOP — that is a hallucination; delete it.
- If you cannot find a specific detail in the sources, either omit it or flag it as "not confirmed in data room materials" or a diligence gap.
- Every quantitative claim (dollar amount, percentage, date, capacity figure) MUST have a citation [n] or [RVM]. If you cannot cite it, do not write it.
- When in doubt, be less specific rather than risk inventing details. "The company's manufacturing facility" is better than fabricating a facility name.

CLAIM DISCIPLINE (source vs. verification):
- CONSERVATIVE ATTRIBUTION — present UNVERIFIED technical performance, customer activity, partnerships, pilots, or commercial traction as REPORTED, not as established fact: "the company reports", "management claims", "company materials indicate", "reported pre-pilot", "planned evaluation". State it as fact only when independent evidence in the sources supports it.
- SOURCE vs MODEL — keep company-reported, independently-verified, extracted-financial-model, and simulation/model-generated figures distinct. NEVER describe a simulation output or a model estimate as a founder projection or a company-reported number.
- COMMERCIAL-STAGE LADDER — use these exact terms and never upgrade a lower stage into a higher one (do NOT call "pilot planning", "pre-pilot", "technical discussion", or "evaluation" a "commercial pilot"): discussion < evaluation < pre-pilot < scoped pilot < paid pilot < qualification < commercial agreement < production deployment. Use the HIGHEST stage the evidence actually supports, and use that same term consistently.
- NO UNSUPPORTED PROMOTION — avoid conclusory promotional phrasing ("the physics works", "first commercially viable application", "eliminates adoption friction", "transformational", "the question is not whether the technology is real") unless verified evidence supports the conclusion.
- NO FALSE GAPS — do not write "not provided" / "no data available" / "data room did not contain" for a datum that is present elsewhere in your material; if a topic is only partly covered, name the SPECIFIC remaining gap (e.g. pricing validation, unit volumes, mfg CAPEX, working-capital, gross-margin support, financing cadence) instead.
- SAY IT ONCE — do not re-derive the same explanation (technology, TAM, customer traction, drop-in integration, binary risk, founder projections) that belongs in another section; state it in its home section and reference it briefly elsewhere.

Rules:
1. Write in professional, data-driven prose — cite specific numbers, percentages, and dollar amounts FROM THE SOURCE DOCUMENTS ONLY
2. Be thorough but avoid padding — every sentence should add value
3. Balance the bull case and bear case — credibility comes from honest assessment, not advocacy
4. Use Markdown formatting: ### for sub-sections, **bold** for emphasis, bullet lists only for catalogs of discrete items
5. Do NOT include the section title as a header — it will be added automatically
6. Length must match the grounded material available. 400-800 words is a CEILING for well-documented sections, never a target to pad toward. When the sources are thin, write a SHORT section (a few honest sentences) that names the diligence gaps — a 3-sentence grounded section is far better than 500 words of invented specifics. Do NOT invent people, companies, or facts to reach a word count.
7. Do NOT reference other sections or say "as discussed in..." — each section stands alone
8. If information is missing or insufficient, explicitly note it as a diligence gap requiring follow-up — do NOT fill gaps with invented details
9. Open with a strong orienting statement that frames why this topic matters for the investment thesis
10. Connect every fact back to its investment implication — never leave data uninterpreted

{_STYLE_GUIDE}"""


def _get_report_fields_for_section(report_context: str, section: dict) -> str:
    """Extract the relevant portions of the report context for a given section."""
    if not report_context:
        return ""

    # The report context has ## headers like "## SIMULATION RESULTS", "## CARBON IMPACT", etc.
    field_mapping = {
        "inputs": ["INPUT PARAMETERS", "DEAL REPORT:", "THIS ROUND"],
        "simulation": ["SIMULATION RESULTS"],
        "adoption": ["MARKET ADOPTION"],
        "carbon": ["CARBON IMPACT"],
        "portfolio_impact": ["PORTFOLIO IMPACT"],
        "sensitivity": ["SENSITIVITY ANALYSIS"],
        "valuation_comps": ["VALUATION CONTEXT"],
        "check_optimization": ["CHECK SIZE OPTIMIZATION"],
        "extraction": ["FINANCIAL MODEL EXTRACTION"],
    }

    wanted_headers = set()
    for rf in section.get("report_fields", []):
        for hdr in field_mapping.get(rf, []):
            wanted_headers.add(hdr)

    if not wanted_headers:
        return ""

    # Parse report_context into sections by ## headers
    parts = []
    current_header = ""
    current_lines = []
    for line in report_context.split('\n'):
        if line.startswith('## '):
            if current_header and current_lines:
                if any(h in current_header for h in wanted_headers):
                    parts.append('\n'.join(current_lines))
            current_header = line
            current_lines = [line]
        elif line.startswith('# '):
            # Top-level header — always include if "DEAL REPORT:" is wanted
            if current_header and current_lines:
                if any(h in current_header for h in wanted_headers):
                    parts.append('\n'.join(current_lines))
            current_header = line
            current_lines = [line]
            if any(h in line for h in wanted_headers):
                pass  # will be included
        else:
            current_lines.append(line)

    if current_header and current_lines:
        if any(h in current_header for h in wanted_headers):
            parts.append('\n'.join(current_lines))

    return '\n\n'.join(parts)


def _pass2_write_section(client, model: str, section: dict, brief: str,
                         report_slice: str, template_guidance: str,
                         company_name: str, links: list,
                         citation_legend: str = "") -> dict:
    """Pass 2: Write a single memo section from its aggregated brief + report data."""

    # No grounded material for this section — do NOT hand the model an empty
    # brief and tell it to "be thorough" (that is precisely when it confabulates
    # people and competitor companies from other deals). Emit an explicit
    # diligence-gap stub and skip the LLM call entirely.
    if not (brief or "").strip() and not (report_slice or "").strip():
        return {
            "text": (f"*The data room did not contain material specific to this "
                     f"section for {company_name or 'this company'}. Treat this as a "
                     f"diligence gap to close before IC — request the supporting "
                     f"documentation from the company.*"),
            "tokens_in": 0, "tokens_out": 0,
        }

    user_parts = [f"# SECTION: {section['title']}"]
    user_parts.append(f"## Section Purpose\n{section['guidance']}")

    if template_guidance:
        user_parts.append(f"## Template Guidance for This Section\n{template_guidance}")

    if citation_legend:
        user_parts.append(f"## Source Document Index\n{citation_legend}\n\nWhen citing information from these documents, use inline citations like [1], [2], etc. matching the numbers above. Place citations at the end of the sentence or clause they support. Use [RVM] when citing quantitative report data. You may combine citations: [1][3] or [1, RVM].")

    if report_slice:
        user_parts.append(f"## Quantitative Report Data [RVM]\n{report_slice}")

    if brief:
        user_parts.append(f"## Extracted Data Room Facts\n{brief}")
    else:
        user_parts.append("## Extracted Data Room Facts\nNo data room documents contained information for this section.")

    if links:
        user_parts.append("## Reference Links\n" + "\n".join(f"- {l}" for l in links))

    user_parts.append(f"\nWrite the '{section['title']}' section of the investment memo for {company_name or 'this company'}. Use ONLY the facts, people, and companies named above — do not introduce any person or organization that is not in this material. Match the section length to the grounded material: a thin section should be short and flag diligence gaps, never padded with invented specifics. Cite your sources using [n] notation.")

    user_msg = "\n\n".join(user_parts)

    try:
        response = client.messages.create(
            model=model,
            max_tokens=3000,
            system=_SECTION_WRITER_SYSTEM,
            messages=[{"role": "user", "content": user_msg}],
        )
        text = "".join(b.text for b in response.content if b.type == "text")
        tokens_in = response.usage.input_tokens if response.usage else 0
        tokens_out = response.usage.output_tokens if response.usage else 0
        return {"text": text.strip(), "tokens_in": tokens_in, "tokens_out": tokens_out}
    except Exception as e:
        logger.error(f"Pass 2 section write failed for {section['key']}: {e}")
        if _is_fatal_llm_error(e):
            raise
        return {"text": f"*[Generation failed for this section: {str(e)[:200]}]*", "tokens_in": 0, "tokens_out": 0}


# ═══════════════════════════════════════════════════════════════════════════════
#  PASS 3 — SYNTHESIS: Exec Summary & Recommendation from all sections
# ═══════════════════════════════════════════════════════════════════════════════

_SYNTHESIS_SYSTEM = f"""You are VoLo Earth Ventures' Investment Committee memo writer.
You are writing a SYNTHESIS section that draws from the entire investment memo.

CRITICAL — FACTUAL ACCURACY:
- ONLY reference facts, names, numbers, and claims that appear in the section texts provided to you.
- NEVER fabricate facility names, locations, dollar amounts, partnerships, customer names, or any specifics not present in the source material.
- Every quantitative claim must be traceable to the section texts or report data provided. If a number does not appear in your inputs, do not invent it.
- When synthesizing, use the same level of specificity as the source sections — do not add details that are not there.

CLAIM DISCIPLINE:
- CONSERVATIVE ATTRIBUTION — carry through the section texts' attribution: unverified technical, customer, pilot, partnership, or traction claims stay "the company reports" / "management claims", not established fact.
- SOURCE vs MODEL — never present a simulation output or a model estimate as a founder projection or a company-reported figure; keep them distinct.
- COMMERCIAL-STAGE LADDER — use the section texts' stage terms consistently (discussion < evaluation < pre-pilot < scoped pilot < paid pilot < qualification < commercial agreement < production deployment); never upgrade a lower stage into a higher one.
- SAY IT ONCE — synthesize; do not restate at length the same technology/TAM/traction/risk explanation the data sections already cover.

Rules:
1. Synthesize across all sections — do not just summarize one part
2. For the Investment Overview: write a compelling one-liner, populate the deal terms table using the CURRENT round's terms (pre-/post-money, check size, ownership) from the "Deal Terms — CURRENT ROUND" block, and list portfolio themes as concise bullets
3. For High Level Opportunities / Risks: write 3-6 specific, evidence-backed bullet points — each should be punchy and data-driven, drawn from the section texts
4. For the Investment Recommendation: give a clear verdict (Invest / Pass / Conditional) with the bull and bear case
5. Include the most important quantitative highlights: MOIC, IRR, P(>3x), carbon t/$, portfolio impact — ONLY if present in the provided data
6. Do NOT include the section title as a header — it will be added automatically
7. Write in clear, authoritative prose grounded in the evidence — conviction comes from specific, sourced facts, not adjectives or hyperbole
8. State a clear verdict (Invest / Pass / Conditional) and EXPLAIN the specific rationale behind it. A Pass must read as a reasoned pass (say WHY it is a pass); an Invest must be justified by the evidence. Do NOT inflate either into promotional language or an unexplained endorsement.
9. FOLLOW-ON deals (when the Deal Terms block says "Follow-on"): VoLo is ALREADY an investor. Every relevant section must address BOTH the new round — the decision at hand, using the current round's terms — AND VoLo's prior investment (combined ownership, total invested across rounds, and how this check builds on the existing position). Keep the two distinct; never use a prior round's valuation as the current headline terms, and never omit the prior investment from the memo.

{_STYLE_GUIDE}"""


def _pass3_synthesize(client, model: str, section: dict, all_section_texts: dict,
                      report_context: str, company_name: str,
                      additional_instructions: str) -> dict:
    """Pass 3: Write synthesis sections (Exec Summary, Recommendation) from all completed sections."""

    # Build a condensed version of all sections
    digest_parts = []
    for s in MEMO_SECTIONS:
        if s["is_synthesis"]:
            continue
        text = all_section_texts.get(s["key"], "")
        if text:
            # Truncate each section to ~800 chars for the digest
            truncated = text[:800] + ("..." if len(text) > 800 else "")
            digest_parts.append(f"### {s['title']}\n{truncated}")

    digest = "\n\n".join(digest_parts)

    # Also include key quantitative highlights from the report
    quant_summary = ""
    if report_context:
        # Pull the hero metrics and headline numbers
        for line in report_context.split('\n'):
            if any(k in line for k in ['Expected MOIC', 'P(>3x)', 'Expected IRR', 'Survival Rate',
                                        'Total Lifecycle tCO2', 'Risk-Adjusted', 'TVPI',
                                        'Optimal Check', 'Fund P50 Impact']):
                quant_summary += line + '\n'

    # The quant filter above keeps only metric lines, so the current round's deal
    # terms would otherwise never reach the synthesis writer — which is how the
    # Investment Overview / Recommendation ended up citing a prior (follow-on) round.
    # Lift the authoritative block out of report_context and pass it through verbatim.
    deal_terms_block = _extract_context_section(report_context, "## THIS ROUND")

    user_parts = [f"# SYNTHESIS SECTION: {section['title']}"]
    user_parts.append(f"## Purpose\n{section['guidance']}")

    if deal_terms_block:
        user_parts.append(
            f"## Deal Terms — CURRENT ROUND (authoritative; use these for the deal terms table)\n{deal_terms_block}"
        )

    if quant_summary:
        user_parts.append(f"## Key Quantitative Highlights\n{quant_summary}")

    user_parts.append(f"## Memo Sections Written So Far\n{digest}")

    if additional_instructions:
        user_parts.append(f"## Additional Instructions\n{additional_instructions}")

    user_parts.append(f"\nWrite the '{section['title']}' section for {company_name or 'this company'}.")

    user_msg = "\n\n".join(user_parts)

    try:
        response = client.messages.create(
            model=model,
            max_tokens=2000,
            system=_SYNTHESIS_SYSTEM,
            messages=[{"role": "user", "content": user_msg}],
        )
        text = "".join(b.text for b in response.content if b.type == "text")
        tokens_in = response.usage.input_tokens if response.usage else 0
        tokens_out = response.usage.output_tokens if response.usage else 0
        return {"text": text.strip(), "tokens_in": tokens_in, "tokens_out": tokens_out}
    except Exception as e:
        logger.error(f"Pass 3 synthesis failed for {section['key']}: {e}")
        if _is_fatal_llm_error(e):
            raise
        return {"text": f"*[Synthesis failed: {str(e)[:200]}]*", "tokens_in": 0, "tokens_out": 0}


# ═══════════════════════════════════════════════════════════════════════════════
#  TEMPLATE PARSING — extract per-section guidance from user template
# ═══════════════════════════════════════════════════════════════════════════════

def _parse_template_sections(template_text: str) -> dict:
    """
    Parse user template into per-section guidance.
    Returns {section_key: "template guidance text"}.
    Uses fuzzy matching on section titles with keyword fallback.
    """
    if not template_text:
        return {}

    # Build mapping from lowercase title fragments to section keys
    title_map = {}
    for s in MEMO_SECTIONS:
        title_lower = s["title"].lower()
        title_map[title_lower] = s["key"]

    # Additional keyword → section key mappings for the VoLo template
    _KEYWORD_MAP = {
        "investment overview": "investment_overview",
        "portfolio themes": "investment_overview",
        "high level opportunities": "high_level_opportunities",
        "opportunities": "high_level_opportunities",
        "high level risks": "high_level_risks",
        "company overview": "company_overview",
        "company": "company_overview",
        "market": "market",
        "business model": "business_model",
        "unit economics": "business_model",
        "team": "team",
        "traction": "traction",
        "competitive position": "competitive_position",
        "competitive": "competitive_position",
        "competition": "competitive_position",
        "carbon impact": "carbon_impact",
        "carbon": "carbon_impact",
        "technology": "technology_ip_moat",
        "ip and moat": "technology_ip_moat",
        "moat": "technology_ip_moat",
        "financing overview": "financing_overview",
        "financing": "financing_overview",
        "round and exit": "financing_overview",
        "exit planning": "financing_overview",
        "deal structure": "financing_overview",
        "round pricing": "financing_overview",
        "capital position": "financing_overview",
        "recommendation": "recommendation",
        "investment recommendation": "recommendation",
    }

    def _match_header(header_text: str) -> Optional[str]:
        ht = header_text.lower().strip()
        # Exact title match
        if ht in title_map:
            return title_map[ht]
        # Keyword map match
        for phrase, key in _KEYWORD_MAP.items():
            if phrase in ht or ht in phrase:
                return key
        # Fallback: single long word match
        for s in MEMO_SECTIONS:
            for word in s["title"].lower().split():
                if len(word) > 4 and word in ht:
                    return s["key"]
        return None

    # Split template by markdown headers (# ## ### or underline-style)
    sections = {}
    current_key = None
    current_lines = []

    for line in template_text.split('\n'):
        # Match markdown headers or pandoc-style underline headers
        header_match = re.match(r'^#{1,3}\s+(.+)', line)
        if not header_match:
            # Also catch pandoc's underline-style: "Title\n====" or "Title\n----"
            # (already converted by pandoc — the title becomes "# Title")
            pass
        if header_match:
            if current_key and current_lines:
                sections[current_key] = '\n'.join(current_lines).strip()
            header_text = header_match.group(1).strip()
            current_key = _match_header(header_text)
            current_lines = []
        elif current_key is not None:
            current_lines.append(line)

    if current_key and current_lines:
        sections[current_key] = '\n'.join(current_lines).strip()

    return sections


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN GENERATE ENDPOINT — orchestrates the 3-pass pipeline
#
#  Memo generation runs as a background job (like DDR). The HTTP endpoint
#  starts a worker thread and returns a job_id immediately so the request
#  cannot be killed by a proxy/CDN HTTP timeout (~2–3 min) while the pipeline
#  is still making its 15–20 sequential LLM calls. The frontend polls
#  /api/memo/generate/status/{job_id} for real progress.
# ═══════════════════════════════════════════════════════════════════════════════


def _run_memo_background(job_id: str, req: "MemoGenerateRequest",
                         user_id: int, user_username: str):
    """Background worker: run _build_memo_payload and stash the result/error."""
    def _update(**kwargs):
        with _MEMO_LOCK:
            job = _MEMO_JOBS.get(job_id)
            if job is None:
                return
            job.update(kwargs)

    def _progress_cb(pct: int, msg: str):
        _update(progress_pct=pct, progress_msg=msg)

    try:
        _update(status="running", progress_pct=1, progress_msg="Starting memo generation...")
        result = _build_memo_payload(req, user_id, progress_cb=_progress_cb)
        _update(
            status="complete",
            progress_pct=100,
            progress_msg="Memo ready.",
            result=result,
            finished_at=datetime.now().isoformat(),
        )
        logger.info(f"[MEMO] Job {job_id} complete for user={user_username}")
    except Exception as e:
        err_str = str(e)[:500] or type(e).__name__
        logger.exception(f"[MEMO] Job {job_id} failed")
        _update(
            status="error",
            progress_pct=0,
            progress_msg="Generation failed.",
            error=err_str,
            finished_at=datetime.now().isoformat(),
        )


@router.post("/generate")
async def generate_memo_start(
    req: MemoGenerateRequest,
    user: CurrentUser = Depends(get_current_user),
):
    """Start a memo-generation background job. Returns a job_id for polling.

    The heavy work (15–20 sequential LLM calls, ~3–5 min) runs in a background
    thread so the HTTP response returns immediately and the connection can't
    be cut by an intermediary proxy timeout.
    """
    _prune_old_memo_jobs()

    # Pre-flight: fail fast on obvious misconfiguration before enqueuing.
    prefs = get_model_preferences(user.id)
    model = req.model_override or prefs.get("memo_generation", MODEL_DEFAULTS["memo_generation"])
    if model.startswith("qwen"):
        if not os.environ.get("REFIANT_API_KEY", ""):
            raise HTTPException(500, "REFIANT_API_KEY not configured. Set it in .env or environment.")
    else:
        if not os.environ.get("ANTHROPIC_API_KEY", ""):
            raise HTTPException(500, "ANTHROPIC_API_KEY not configured. Set it in .env or environment.")

    # Minimal input validation: must have *something* to work with.
    has_report = bool(req.report_id)
    has_session = bool(req.session_id)
    has_library = bool(req.library_id)
    has_template = bool((req.template_text or "").strip() or req.template_id)
    if not (has_report or has_session or has_library or has_template):
        raise HTTPException(400, "Select a deal report, upload documents, or provide a template.")

    job_id = str(uuid.uuid4())[:12]
    with _MEMO_LOCK:
        _MEMO_JOBS[job_id] = {
            "status": "queued",
            "progress_pct": 0,
            "progress_msg": "Queued...",
            "error": None,
            "result": None,
            "user_id": user.id,
            "started_at": datetime.now().isoformat(),
            "finished_at": None,
            "_created_mono": time.time(),
        }

    t = threading.Thread(
        target=_run_memo_background,
        args=(job_id, req, user.id, user.username),
        daemon=True,
    )
    t.start()
    return {"job_id": job_id, "status": "queued"}


@router.get("/generate/status/{job_id}")
async def generate_memo_status(
    job_id: str,
    user: CurrentUser = Depends(get_current_user),
):
    """Poll memo-generation status. Returns progress while running, the full
    response payload once complete, or the error message on failure."""
    with _MEMO_LOCK:
        job = _MEMO_JOBS.get(job_id)
        snapshot = dict(job) if job else None

    if not snapshot:
        raise HTTPException(404, "Memo job not found or expired")
    if snapshot.get("user_id") != user.id:
        raise HTTPException(403, "Not your job")

    resp = {
        "job_id": job_id,
        "status": snapshot["status"],
        "progress_pct": snapshot["progress_pct"],
        "progress_msg": snapshot["progress_msg"],
        "error": snapshot["error"],
        "started_at": snapshot["started_at"],
        "finished_at": snapshot["finished_at"],
    }
    # Inline the full payload on completion so the client gets the memo in one
    # poll — same shape the old synchronous /generate endpoint returned.
    if snapshot["status"] == "complete" and snapshot.get("result"):
        resp["result"] = snapshot["result"]
    return resp


def _build_memo_payload(
    req: MemoGenerateRequest,
    user_id: int,
    progress_cb=None,
):
    """Run the full memo-generation pipeline synchronously. Returns the final
    response payload dict, or raises RuntimeError with a user-friendly message.

    Designed to be invoked from a background worker thread — no HTTPException,
    no async/await. Emits real progress via `progress_cb(pct, msg)`.
    """
    from ..engine.llm_utils import make_llm_client
    start_time = time.time()

    def _progress(pct: int, msg: str):
        if progress_cb:
            try: progress_cb(pct, msg)
            except Exception: pass

    _progress(2, "Resolving model and loading inputs...")

    # ── Resolve model ──
    prefs = get_model_preferences(user_id)
    model = req.model_override or prefs.get("memo_generation", MODEL_DEFAULTS["memo_generation"])
    is_refiant = model.startswith("qwen")

    if is_refiant:
        api_key = os.environ.get("REFIANT_API_KEY", "")
        if not api_key:
            raise RuntimeError("REFIANT_API_KEY not configured. Set it in .env or environment.")
    else:
        api_key = os.environ.get("ANTHROPIC_API_KEY", "")
        if not api_key:
            raise RuntimeError("ANTHROPIC_API_KEY not configured. Set it in .env or environment.")
    client = make_llm_client(is_refiant, api_key)

    # ── Resolve active sections based on investment type ──
    active_data, active_synth = _get_active_sections(req.investment_type)
    is_followon = req.investment_type == "followon"

    # ── Gather raw inputs ──
    report_data_json = None  # Raw report JSON for chart embedding
    report_inputs_json = None  # Raw inputs JSON (for the QA consistency backstop)
    conn = get_db()
    try:
        # Report
        report_context = ""
        company_name = req.company_name
        report_id = req.report_id
        if report_id:
            # Reports are shared firm-wide — any signed-in user can pull any
            # report into their memo (no owner_id filter).
            row = conn.execute(
                "SELECT * FROM deal_reports WHERE id=?",
                (report_id,),
            ).fetchone()
            if row:
                report_context = _build_report_context(row)
                report_data_json = row["report_json"]
                report_inputs_json = row["inputs_json"]
                if not company_name:
                    company_name = row["company_name"]

        # Template
        template_text = req.template_text or ""
        template_id = req.template_id
        if not template_text and template_id:
            tpl = conn.execute(
                "SELECT content FROM memo_templates WHERE id=? AND owner_id=?",
                (template_id, user_id),
            ).fetchone()
            if tpl:
                template_text = tpl["content"]
    finally:
        conn.close()

    # Load raw documents — from session uploads and/or deal library
    raw_docs = []
    if req.session_id:
        raw_docs = _load_raw_documents(req.session_id, user_id)
    if req.library_id:
        from .drive import load_library_documents
        library_docs = load_library_documents(req.library_id, user_id)
        raw_docs.extend(library_docs)

    links = req.links or []
    template_sections = _parse_template_sections(template_text)
    total_tokens_in = 0
    total_tokens_out = 0
    pass_log = []

    # ── Build citation index ──
    citation_index = _build_citation_index(raw_docs) if raw_docs else {"map": {}, "legend": "", "count": 0}

    # ═════════════════════════════════════════════════════════════════════════
    #  Run all LLM passes. We emit real progress via `progress_cb` before
    #  each LLM call so the frontend can show actual status (not an estimate).
    #  Budget: 2% → 95% spans the LLM pipeline; the remaining 5% covers
    #  markdown assembly, HTML conversion, and the DB insert.
    # ═════════════════════════════════════════════════════════════════════════

    _locked_preview = req.locked_sections or {}
    _n_docs = len(raw_docs)
    _n_data_steps = sum(1 for s in active_data if s["key"] not in _locked_preview)
    _n_synth_steps = sum(1 for s in active_synth if s["key"] not in _locked_preview)
    _total_llm_steps = max(1, _n_docs + _n_data_steps + _n_synth_steps)
    _pct_span = 93  # 2% → 95%

    def _emit_step_progress(step_idx: int, msg: str):
        pct = 2 + int((step_idx / _total_llm_steps) * _pct_span)
        _progress(min(pct, 95), msg)

    def _strip_leading_heading(text: str, section_title: str) -> str:
        """Remove any leading markdown heading the LLM added despite instructions.
        Strips the first line if it is a #/##/### heading that matches (or approximates)
        the section title, then strips any leading blank lines left behind."""
        if not text:
            return text
        lines = text.split("\n")
        first = lines[0].lstrip()
        if first.startswith("#"):
            # Remove all leading # characters and compare loosely
            heading_text = first.lstrip("#").strip().lower()
            title_lower  = section_title.lower()
            # Strip if the heading is the same as or a prefix/suffix of the title
            if (heading_text in title_lower or title_lower in heading_text
                    or heading_text.replace(" ", "") == title_lower.replace(" ", "")):
                lines = lines[1:]
                # Drop any blank lines that were sitting right after the heading
                while lines and not lines[0].strip():
                    lines = lines[1:]
        return "\n".join(lines)

    def _run_llm_pipeline():
        _total_in = 0
        _total_out = 0
        _pass_log = []
        _step = 0
        _locked = req.locked_sections or {}
        section_texts = {}

        # ── ENGINE DISPATCH ──────────────────────────────────────────────
        # v2: manifest pass + targeted deep reads + cached section writing
        # v1: legacy per-doc Pass 1 → per-section Pass 2
        # In both cases, Pass 3 (synthesis) runs the same way after.
        engine = (req.engine_version or "v1").lower()

        if engine == "v2":
            from ..engine.memo_v2 import run_memo_v2_pipeline
            _emit_step_progress(0, "Reading through the data room...")
            v2_section_texts, v2_in, v2_out, v2_log, manifest_meta = run_memo_v2_pipeline(
                client=client,
                model=model,
                raw_docs=raw_docs,
                report_context=report_context,
                memo_sections=active_data,
                template_sections=template_sections,
                company_name=company_name,
                links=links,
                locked_sections=_locked,
                additional_instructions=req.additional_instructions,
                progress_cb=_progress,
            )
            # Strip any leading heading the model added despite instructions
            for sk, txt in v2_section_texts.items():
                section_def = next((s for s in active_data if s["key"] == sk), None)
                if section_def:
                    v2_section_texts[sk] = _strip_leading_heading(txt, section_def["title"])
            section_texts.update(v2_section_texts)
            _total_in += v2_in
            _total_out += v2_out
            _pass_log.extend(v2_log)
            _pass_log.append({"stage": "manifest_meta", "meta": manifest_meta})
            # v2 emits its own absolute-percent progress (8% → ~90%) via the
            # progress_cb. The shared Pass 3 below uses _emit_step_progress
            # which calculates percent from _step / _total_llm_steps. Advance
            # _step so the synthesis bar continues from where v2 left off
            # instead of jumping backwards.
            _step = int(_total_llm_steps * 0.94)
        else:
            # ── v1 PASS 1: Extract facts from each document ──────────────
            all_extractions = []
            if raw_docs:
                for i, doc in enumerate(raw_docs):
                    _emit_step_progress(_step, f"Extracting facts from {doc['file_name']} ({i+1}/{_n_docs})")
                    extraction = _pass1_extract_document(client, model, doc)
                    all_extractions.append(extraction)
                    _total_in += extraction.get("_tokens_in", 0)
                    _total_out += extraction.get("_tokens_out", 0)
                    _pass_log.append({"pass": 1, "doc": doc["file_name"],
                                      "tokens": extraction.get("_tokens_in", 0) + extraction.get("_tokens_out", 0)})
                    _step += 1

            # ── v1 PASS 2: Write each data section ───────────────────────
            _data_total = len([s for s in active_data if s["key"] not in _locked])
            _data_done = 0
            for section in active_data:
                sk = section["key"]
                if sk in _locked:
                    # Preserve the user's saved edits exactly — skip LLM
                    section_texts[sk] = _locked[sk]
                    _pass_log.append({"pass": 2, "section": sk, "tokens": 0, "locked": True})
                    continue
                _data_done += 1
                _emit_step_progress(_step, f"Writing section: {section['title']} ({_data_done}/{_data_total})")
                brief = _aggregate_section_briefs(all_extractions, sk, citation_index)
                report_slice = _get_report_fields_for_section(report_context, section)
                tpl_guidance = template_sections.get(sk, "")

                result = _pass2_write_section(
                    client, model, section, brief, report_slice,
                    tpl_guidance, company_name, links,
                    citation_legend=citation_index["legend"]
                )
                section_texts[sk] = _strip_leading_heading(result["text"], section["title"])
                _total_in += result["tokens_in"]
                _total_out += result["tokens_out"]
                _pass_log.append({"pass": 2, "section": sk,
                                  "tokens": result["tokens_in"] + result["tokens_out"]})
                _step += 1

        # ── PASS 3 (shared): Synthesize cross-cutting sections ───────────
        _synth_total = len([s for s in active_synth if s["key"] not in _locked])
        _synth_done = 0
        for section in active_synth:
            sk = section["key"]
            if sk in _locked:
                section_texts[sk] = _locked[sk]
                _pass_log.append({"pass": 3, "section": sk, "tokens": 0, "locked": True})
                continue
            _synth_done += 1
            _emit_step_progress(_step, f"Synthesizing: {section['title']} ({_synth_done}/{_synth_total})")
            result = _pass3_synthesize(
                client, model, section, section_texts,
                report_context, company_name,
                req.additional_instructions
            )
            section_texts[sk] = _strip_leading_heading(result["text"], section["title"])
            _total_in += result["tokens_in"]
            _total_out += result["tokens_out"]
            _pass_log.append({"pass": 3, "section": sk,
                              "tokens": result["tokens_in"] + result["tokens_out"]})
            _step += 1

        return section_texts, _total_in, _total_out, _pass_log

    try:
        pipeline_result = _run_llm_pipeline()
    except Exception as pipeline_err:
        err_str = str(pipeline_err)
        # Classify the error into a human-readable message
        err_lower = err_str.lower()
        if "credit balance exhausted" in err_lower or "credit_balance_exhausted" in err_lower:
            detail = (
                "Credit balance exhausted — the API account has no remaining credits. "
                "Top up at console.anthropic.com (Anthropic) or refiant.ai (Refiant/Qwen) "
                "then try again."
            )
        elif "overloaded" in err_lower or "529" in err_str:
            detail = "The AI provider is currently overloaded. Please wait a moment and try again."
        elif "rate limit" in err_lower or "rate_limit" in err_lower or "429" in err_str:
            detail = "Rate limit reached. Please wait a moment and try again."
        elif "invalid_api_key" in err_lower or "authentication" in err_lower or "401" in err_str:
            detail = "Invalid API key. Check your ANTHROPIC_API_KEY or REFIANT_API_KEY in .env."
        elif "context_length" in err_lower or "context length" in err_lower:
            detail = "Input is too long for the selected model. Try reducing the number of documents."
        else:
            detail = f"Memo generation failed: {err_str[:300]}"
        raise RuntimeError(detail)
    section_texts, total_tokens_in, total_tokens_out, pass_log = pipeline_result

    # Prepend two deterministic, clearly-separated blocks to the Financing
    # Overview section, built straight from the report so the numbers can't
    # drift: (1) Financing Summary — prior financing vs the current round
    # (terms); (2) Return Profile — This Check (the decision) vs Total Position
    # (blended prior + new, flagged as context only). These are explicit
    # separations on top of the financing detail woven through the memo.
    try:
        _fin_key = "financing_overview"
        if _fin_key in section_texts and _fin_key not in (req.locked_sections or {}):
            _rp = json.loads(report_data_json) if isinstance(report_data_json, str) else (report_data_json or {})
            _rp = _rp or {}
            _blocks = [
                _financing_summary_block(_rp.get("deal_overview", {})),
                _return_profile_block(_rp),
            ]
            _prefix = "\n".join(b for b in _blocks if b)
            if _prefix:
                section_texts[_fin_key] = _prefix + "\n" + (section_texts.get(_fin_key) or "")
    except Exception:
        pass

    _progress(95, "Assembling memo and saving to database...")

    # ═════════════════════════════════════════════════════════════════════════
    #  ASSEMBLE final memo in section order
    # ═════════════════════════════════════════════════════════════════════════
    memo_parts = [f"# Investment Memorandum: {company_name or 'Deal Analysis'}\n"]
    memo_parts.append(f"*VoLo Earth Ventures — Confidential*\n")

    all_active_sections = active_data + active_synth
    # Reorder: synthesis sections that come first in MEMO_SECTIONS, then data, then remaining synthesis
    ordered_sections = [s for s in MEMO_SECTIONS if s in all_active_sections]
    for section in ordered_sections:
        text = section_texts.get(section["key"], "")
        if text:
            memo_parts.append(f"## {section['title']}\n\n{text}\n")

    # Appendix: data sources with citation numbers
    memo_parts.append("## Appendix: Sources\n")

    # 1. Data room documents reviewed
    if raw_docs:
        memo_parts.append("### I. Data Room Documents Reviewed")
        cite_map = citation_index.get("map", {})
        for doc in raw_docs:
            cat = doc["doc_category"].replace('_', ' ').title()
            num = cite_map.get(doc["file_name"], "?")
            memo_parts.append(f"- **[{num}]** [{cat}] {doc['file_name']}")
        memo_parts.append("")

    # 2. RVM deal report
    if report_context:
        memo_parts.append("### II. Quantitative Model (RVM)")
        memo_parts.append("- **[RVM]** VoLo Return Validation Model — Deal Report. Proprietary three-scenario DCF engine with Monte Carlo simulation across revenue, margin, exit multiple, and timing uncertainty. Produces MOIC, IRR, P(>3×), and carbon impact distributions.\n")

    # 3. Reference links from memo request
    if links:
        memo_parts.append("### III. Reference Links")
        for l in links:
            memo_parts.append(f"- {l}")
        memo_parts.append("")

    # 4. Underlying analytical data sources used by the VoLo RVM
    memo_parts.append("### IV. Underlying Analytical Data Sources")
    memo_parts.append(
        "_The following datasets and methodologies underpin the quantitative assumptions "
        "embedded in the VoLo Return Validation Model (RVM) and deal pipeline analysis._\n"
    )
    _PIPELINE_SOURCES = [
        {
            "id": "carta_rounds",
            "name": "Carta Insights — Fund Forecasting Profiles",
            "provider": "Carta",
            "category": "Financing & Valuation",
            "description": "Round sizing by stage/sector (p10–p90 percentiles), pre/post-money valuations, ESOP metrics, graduation rates, and time-to-graduation. Used to calibrate stage-appropriate check size ranges, valuation entry points, and dilution assumptions.",
            "url": "https://carta.com/blog/startup-financing-data/",
            "used_for": "Check size defaults, valuation benchmarks, dilution modeling",
        },
        {
            "id": "carta_benchmarks",
            "name": "Carta TVPI Fund Benchmarks",
            "provider": "Carta",
            "category": "Fund Performance",
            "description": "Fund TVPI percentiles (p10, p50, p75, p90) by fund age. Used for portfolio-level performance overlay and convergence-driven position sizing.",
            "url": "https://carta.com/blog/startup-financing-data/",
            "used_for": "Fund TVPI benchmarking, portfolio construction optimization",
        },
        {
            "id": "nrel_atb",
            "name": "NREL Annual Technology Baseline 2024 v3",
            "provider": "National Renewable Energy Laboratory (NREL)",
            "category": "Technology Cost",
            "description": "LCOE projections by technology and cost case, deployment cost benchmarks, representative technology classes, capacity factors, and CAPEX curves. Primary source for Bass diffusion parameter calibration across energy technology archetypes.",
            "url": "https://atb.nrel.gov/",
            "used_for": "Technology cost benchmarks, Bass S-curve calibration, adoption timelines",
        },
        {
            "id": "lazard_lcoe",
            "name": "Lazard Levelized Cost of Energy+ (LCOE+)",
            "provider": "Lazard",
            "category": "Technology Cost",
            "description": "Energy cost benchmarks (solar, wind, geothermal, nuclear, gas, coal, battery storage) in $/MWh ranges. Used to benchmark cost competitiveness of portfolio company technologies against incumbent energy sources.",
            "url": "https://www.lazard.com/research-insights/levelized-cost-of-energyplus/",
            "used_for": "Technology cost competitiveness, incumbent displacement analysis",
        },
        {
            "id": "doe_electrification",
            "name": "DOE Electrification Pathways Data Appendix",
            "provider": "U.S. Department of Energy",
            "category": "Technology Deployment",
            "description": "Electrification pathways for EV and industrial technology deployment scenarios. Informs market adoption rate calibration and TAM trajectory for electrification-adjacent technologies.",
            "url": "https://www.energy.gov/eere/analysis/electrification-futures-study",
            "used_for": "EV and industrial electrification TAM projections",
        },
        {
            "id": "damodaran_comps",
            "name": "Damodaran EV/EBITDA Public Comps",
            "provider": "Aswath Damodaran, Stern School of Business, NYU",
            "category": "Valuation Multiples",
            "description": "EV/EBITDA multiples for 97 US industries updated annually. Applied to terminal year EBITDA estimates with a 20% private-company discount (IPO haircut) to derive acquisition exit multiples.",
            "url": "https://pages.stern.nyu.edu/~adamodar/New_Home_Page/data.html",
            "used_for": "Exit multiple ranges, sector-specific valuation benchmarks",
        },
        {
            "id": "ebitda_margins",
            "name": "EBITDA Margin Ramp Model by TRL Level",
            "provider": "SaaS Capital, Bessemer Cloud Index, Battery Ventures, NREL/DOE (synthesized)",
            "category": "Financial Modeling",
            "description": "TRL-dependent EBITDA margin start/end/ramp parameters. Early-stage (TRL 1–3): −20% to +20% over 10 years; commercial stage (TRL 7–9): 18% to 32% over 2 years. Calibrated from SaaS and deep-tech operating benchmarks.",
            "url": None,
            "used_for": "Gross margin and EBITDA ramp assumptions by technology maturity stage",
        },
        {
            "id": "bass_diffusion",
            "name": "Bass Diffusion Technology Adoption Parameters",
            "provider": "NREL ATB + Historical Market Data (calibrated by VoLo)",
            "category": "Market Adoption",
            "description": "Innovation coefficient (p) and imitation coefficient (q) for 12 technology archetypes, inflection years, and maturity stages. Drives the S-curve adoption model and revenue cone simulation. Calibrated from historical deployment data including utility solar (1.3% → 13% in 10 yr), EV adoption, and enterprise SaaS penetration curves.",
            "url": "https://atb.nrel.gov/",
            "used_for": "Technology adoption S-curve, revenue projection fan chart, market timing",
        },
        {
            "id": "market_sizing",
            "name": "TAM/SAM/SOM Defaults by Technology Archetype",
            "provider": "BloombergNEF, Wood Mackenzie, IEA, Rystad Energy, McKinsey Global Institute, Gartner, Grand View Research",
            "category": "Market Sizing",
            "description": "Total addressable market sizing for 12 technology archetypes (e.g., Utility Solar $120B, EV Electrification $500B, AI/ML $300B, Industrial Decarbonization $180B). TAM estimates are cross-referenced across multiple provider forecasts and represent 2030 horizon figures.",
            "url": None,
            "used_for": "Base TAM/SAM/SOM defaults, market penetration ceiling, deal sizing context",
        },
        {
            "id": "carbon_intensity",
            "name": "Carbon Intensity & Avoided Emissions Model",
            "provider": "EPA eGRID, IPCC AR6, EIA AEO, VoLo Earth Proprietary (RVM 1.19)",
            "category": "Carbon Impact",
            "description": "Carbon intensity (tCO2/unit) by displaced resource using EPA eGRID regional emissions factors, IPCC AR6 lifecycle analysis, and EIA Annual Energy Outlook. Includes 40% methane leakage premium on natural gas per EPA 2014 methodology. TRL-to-risk divisor mapping for early-stage carbon accounting uncertainty.",
            "url": "https://www.epa.gov/egrid",
            "used_for": "Avoided emissions quantification, carbon cost of capital adjustment, ESG impact reporting",
        },
        {
            "id": "private_discount",
            "name": "Private Company Acquisition Discount (DLOM)",
            "provider": "Koeplin, Sarin & Shapiro (2000); Officer (2007)",
            "category": "Valuation Multiples",
            "description": "20% discount for lack of marketability (DLOM) applied to IPO/public comps when modeling acquisition exit paths. Academic range is 15–30%; 20% is the midpoint used per Koeplin et al. (2000) Journal of Financial Economics and Officer (2007) Journal of Finance.",
            "url": None,
            "used_for": "M&A exit multiple haircut, acquisition path valuation",
        },
        {
            "id": "cambridge_exits",
            "name": "Venture Exit Year Distribution",
            "provider": "Cambridge Associates Venture Benchmarks",
            "category": "Exit Modeling",
            "description": "Exit year probability weighting calibrated from Cambridge Associates venture fund data. Venture exits cluster in years 4–7 from entry, with tails to year 12+. Applied as a prior in Monte Carlo exit timing draws.",
            "url": "https://www.cambridgeassociates.com/",
            "used_for": "Exit timing distribution, hold period assumptions, IRR discounting",
        },
    ]
    for src in _PIPELINE_SOURCES:
        url_str = f" — [{src['url']}]({src['url']})" if src.get("url") else ""
        memo_parts.append(
            f"- **{src['name']}** *(_{src['provider']}_)*{url_str}  \n"
            f"  *Category: {src['category']} · Used for: {src['used_for']}*  \n"
            f"  {src['description']}\n"
        )
    memo_parts.append("")

    memo_md = "\n".join(memo_parts)

    # ── Consistency / QA backstop (fixes 1/4/13/14/33 + soft gating) ──────────
    # Prevention is upstream (canonical facts + prompt discipline). This SOFT
    # pass only appends a clearly-labeled review of anything that slipped
    # through; it never blocks rendering.
    try:
        _qa_rp = json.loads(report_data_json) if isinstance(report_data_json, str) else (report_data_json or {})
        _qa_inp = json.loads(report_inputs_json) if isinstance(report_inputs_json, str) else (report_inputs_json or {})
        _qa_block = _consistency_qa_block(_qa_rp or {}, _qa_inp or {}, section_texts, memo_md)
        if _qa_block:
            memo_md = memo_md + "\n\n" + _qa_block
    except Exception as _qa_err:
        logger.warning(f"QA consistency backstop skipped: {_qa_err}")

    elapsed = time.time() - start_time

    # ── Convert to HTML ──
    memo_html = _markdown_to_html(memo_md)

    # ── Save to DB ──
    # Store section-level texts for per-section editing
    sections_json_str = json.dumps(section_texts)

    conn = get_db()
    try:
        cur = conn.execute(
            """INSERT INTO generated_memos
               (owner_id, report_id, template_id, company_name, memo_markdown, memo_html,
                model_used, input_token_count, output_token_count, generation_time_s, sections_json, memo_session_id)
               VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
            (user_id, report_id, template_id, company_name, memo_md, memo_html,
             model, total_tokens_in, total_tokens_out, elapsed, sections_json_str, req.session_id or ""),
        )
        conn.commit()
        memo_id = cur.lastrowid
    finally:
        conn.close()

    # Build citation metadata for frontend popovers — rich page-level excerpts
    citations_meta = {}
    if raw_docs:
        cite_map = citation_index.get("map", {})
        for doc in raw_docs:
            num = cite_map.get(doc["file_name"], 0)
            cat = doc["doc_category"].replace('_', ' ').title()
            full_text = doc.get("extracted_text") or ""
            doc_id = doc.get("id")
            file_type = Path(doc["file_name"]).suffix.lower()

            # Build page-level excerpts for PDFs (split on [Page N] markers)
            pages = []
            if "[Page " in full_text:
                import re as _re
                page_splits = _re.split(r'\[Page (\d+)\]\n?', full_text)
                # page_splits: ['', '1', 'text...', '2', 'text...', ...]
                for i in range(1, len(page_splits) - 1, 2):
                    page_num = int(page_splits[i])
                    page_text = page_splits[i + 1].strip()
                    if page_text:
                        pages.append({
                            "page": page_num,
                            "text": page_text[:2000],
                            "truncated": len(page_text) > 2000,
                        })
            else:
                # Non-PDF or no page markers — treat entire text as one block
                if full_text.strip():
                    pages.append({
                        "page": None,
                        "text": full_text[:3000],
                        "truncated": len(full_text) > 3000,
                    })

            # Key facts summary (first 300 chars of non-empty content)
            summary = ""
            for p in pages:
                if p["text"].strip():
                    summary = p["text"][:300].strip()
                    if len(p["text"]) > 300:
                        summary += "..."
                    break

            citations_meta[str(num)] = {
                "number": num,
                "file_name": doc["file_name"],
                "category": cat,
                "file_type": file_type,
                "doc_id": doc_id,
                "total_chars": len(full_text),
                "total_pages": len(pages),
                "summary": summary,
                "pages": pages[:20],  # Cap at 20 pages for frontend
            }
    if report_context:
        # Build structured RVM citation with key metrics
        rvm_sections = []
        for line in report_context.split('\n'):
            line = line.strip()
            if line and not line.startswith('---'):
                rvm_sections.append(line)
        citations_meta["RVM"] = {
            "number": "RVM",
            "file_name": "VoLo Return Validation Model",
            "category": "Deal Report",
            "file_type": "rvm",
            "total_chars": len(report_context),
            "total_pages": 1,
            "summary": "\n".join(rvm_sections[:10]),
            "pages": [{"page": None, "text": report_context[:5000], "truncated": len(report_context) > 5000}],
        }

    # Parse report data for frontend chart embedding
    report_data_parsed = None
    if report_data_json:
        try:
            report_data_parsed = json.loads(report_data_json) if isinstance(report_data_json, str) else report_data_json
        except (json.JSONDecodeError, TypeError):
            report_data_parsed = None

    # Build image metadata for inline embedding — auto-match images to sections by category.
    # raw_docs already includes image files (see _load_raw_documents), so we just filter here.
    # We use file_type (stored as the extension at upload time) to avoid re-parsing file_name.
    image_docs = []
    for doc in raw_docs:
        ext = (doc.get("file_type") or Path(doc["file_name"]).suffix).lower()
        if ext in _IMAGE_EXTS and doc.get("id"):
            image_docs.append({
                "doc_id": doc["id"],
                "file_name": doc["file_name"],
                "category": doc.get("doc_category", "other"),
                "file_type": ext,
                "section_targets": _image_category_to_sections(doc.get("doc_category", "other"), doc.get("file_name", "")),
            })

    return {
        "id": memo_id,
        "company_name": company_name,
        "memo_markdown": memo_md,
        "memo_html": memo_html,
        "model_used": model,
        "input_tokens": total_tokens_in,
        "output_tokens": total_tokens_out,
        "generation_time_s": round(elapsed, 2),
        "citations": citations_meta,
        "report_data": report_data_parsed,
        "image_docs": image_docs,
        "sections": section_texts,
        "pipeline": {
            "documents_processed": len(raw_docs),
            "sections_written": len(section_texts),
            "passes": pass_log,
            "total_llm_calls": len(pass_log),
        },
    }


def _image_category_to_sections(category: str, file_name: str = "") -> list:
    """Map a document category (and filename) to memo sections for image embedding.
    Filename keywords take priority over category when they match — this handles
    cases like board_report PDFs whose extracted images have descriptive filenames.
    """
    CATEGORY_MAP = {
        "pitch_deck": ["company_overview", "market", "business_model", "competitive_position"],
        "financial_model": ["financials", "business_model", "financing_overview"],
        "term_sheet": ["financing_overview"],
        "cap_table": ["financing_overview"],
        "legal": ["financing_overview"],
        "ip_patent": ["technology_ip_moat"],
        "customer_reference": ["traction", "competitive_position"],
        "market_research": ["market", "company_overview"],
        "technical_diligence": ["technology_ip_moat", "company_overview"],
        "team_bios": ["team"],
        "board_materials": ["traction", "financing_overview"],
        "board_report": ["technology_ip_moat", "competitive_position", "market", "traction", "company_overview"],
        "screenshots": ["company_overview", "technology_ip_moat", "market", "traction", "business_model", "financials"],
        "other": ["company_overview"],
    }
    # Filename-based routing — override category when filename signals a specific section
    fname = file_name.lower()
    FILENAME_SIGNALS = [
        (["competitive", "competition", "landscape", "competitor"], ["competitive_position", "market"]),
        (["technology", "tech", "platform", "architecture", "system", "resource_evaluation", "evaluation", "technical"], ["technology_ip_moat", "company_overview"]),
        (["market", "tam", "industry", "sector", "sizing"], ["market", "company_overview"]),
        (["team", "management", "leadership", "founder", "bio"], ["team"]),
        (["financial", "revenue", "forecast", "projection", "model", "p&l", "income"], ["financials", "business_model"]),
        (["traction", "customer", "pipeline", "milestone", "kpi", "growth"], ["traction", "competitive_position"]),
        (["carbon", "emission", "climate", "sustainability", "ghg"], ["carbon_impact"]),
        (["patent", "ip", "intellectual"], ["technology_ip_moat"]),
        (["cap_table", "captable", "ownership", "equity"], ["financing_overview"]),
    ]
    for keywords, sections in FILENAME_SIGNALS:
        if any(kw in fname for kw in keywords):
            return sections
    return CATEGORY_MAP.get(category, ["company_overview"])


def _markdown_to_html(md: str) -> str:
    """Lightweight markdown-to-HTML conversion."""
    html = md

    # Headers
    html = re.sub(r'^#### (.+)$', r'<h4>\1</h4>', html, flags=re.MULTILINE)
    html = re.sub(r'^### (.+)$', r'<h3>\1</h3>', html, flags=re.MULTILINE)
    html = re.sub(r'^## (.+)$', r'<h2>\1</h2>', html, flags=re.MULTILINE)
    html = re.sub(r'^# (.+)$', r'<h1>\1</h1>', html, flags=re.MULTILINE)

    # Bold and italic
    html = re.sub(r'\*\*\*(.+?)\*\*\*', r'<strong><em>\1</em></strong>', html)
    html = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', html)
    html = re.sub(r'\*(.+?)\*', r'<em>\1</em>', html)

    # Lists
    html = re.sub(r'^- (.+)$', r'<li>\1</li>', html, flags=re.MULTILINE)
    html = re.sub(r'(<li>.*</li>\n?)+', lambda m: '<ul>' + m.group(0) + '</ul>', html)

    # Numbered lists
    html = re.sub(r'^\d+\. (.+)$', r'<li>\1</li>', html, flags=re.MULTILINE)

    # Tables (basic pipe tables)
    lines = html.split('\n')
    in_table = False
    table_lines = []
    result = []
    for line in lines:
        if '|' in line and line.strip().startswith('|'):
            if not in_table:
                in_table = True
                table_lines = []
            table_lines.append(line)
        else:
            if in_table:
                result.append(_convert_table(table_lines))
                in_table = False
                table_lines = []
            result.append(line)
    if in_table:
        result.append(_convert_table(table_lines))
    html = '\n'.join(result)

    # Paragraphs — wrap non-tag text blocks
    html = re.sub(r'^(?!<[hul1-9ol]|<li|<table|<tr|<th|<td)(.+)$', r'<p>\1</p>', html, flags=re.MULTILINE)

    # Clean up empty paragraphs
    html = re.sub(r'<p>\s*</p>', '', html)

    # ── Inline citations: convert [1], [2], [RVM], [1, 3], etc. into clickable spans ──
    def _cite_replace(m):
        raw = m.group(1)  # e.g. "1" or "RVM" or "1, 3"
        parts = [p.strip() for p in raw.split(',')]
        spans = []
        for p in parts:
            css_class = "memo-cite"
            data_cite = p
            if p.upper() == "RVM":
                css_class += " memo-cite-rvm"
            spans.append(f'<span class="{css_class}" data-cite="{data_cite}" title="Source [{p}] — click to view">[{p}]</span>')
        return ''.join(spans)

    # Match [n], [RVM], [n, m], but NOT markdown links [text](url) or bold **[text]**
    html = re.sub(r'(?<!\()\[(\d+(?:\s*,\s*\d+)*(?:\s*,\s*RVM)?|RVM(?:\s*,\s*\d+)*)\](?!\()', _cite_replace, html)

    return html


def _convert_table(lines):
    """Convert markdown pipe table to HTML."""
    if len(lines) < 2:
        return '\n'.join(lines)

    html = '<table class="memo-table">'

    # Header row
    cells = [c.strip() for c in lines[0].strip('|').split('|')]
    html += '<thead><tr>' + ''.join(f'<th>{c}</th>' for c in cells) + '</tr></thead>'

    # Body rows (skip separator line)
    html += '<tbody>'
    for line in lines[2:]:
        cells = [c.strip() for c in line.strip('|').split('|')]
        html += '<tr>' + ''.join(f'<td>{c}</td>' for c in cells) + '</tr>'
    html += '</tbody></table>'

    return html


# ═══════════════════════════════════════════════════════════════════════════════
#  MEMO HISTORY
# ═══════════════════════════════════════════════════════════════════════════════

@router.get("/history")
async def list_memos(user: CurrentUser = Depends(get_current_user)):
    conn = get_db()
    try:
        rows = conn.execute(
            """SELECT id, company_name, model_used, input_token_count, output_token_count,
                      generation_time_s, status, created_at
               FROM generated_memos WHERE owner_id=? ORDER BY created_at DESC LIMIT 50""",
            (user.id,),
        ).fetchall()
        return [dict(r) for r in rows]
    finally:
        conn.close()


@router.get("/history/{memo_id}")
async def get_memo(memo_id: int, user: CurrentUser = Depends(get_current_user)):
    """Fetch a single memo. Visible to any authenticated VoLo team member —
    the team library is shared so analysts can review each other's memos."""
    conn = get_db()
    try:
        row = conn.execute(
            "SELECT * FROM generated_memos WHERE id=?",
            (memo_id,),
        ).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Memo not found")
        result = dict(row)
        # Parse sections_json so frontend can use it directly
        try:
            result["sections"] = json.loads(row["sections_json"] or "{}")
        except (json.JSONDecodeError, TypeError):
            result["sections"] = {}
        # Include report_data for chart injection
        if row["report_id"]:
            rpt_row = conn.execute(
                "SELECT report_json FROM deal_reports WHERE id=?",
                (row["report_id"],),
            ).fetchone()
            if rpt_row:
                try:
                    result["report_data"] = json.loads(rpt_row["report_json"] or "{}")
                except (json.JSONDecodeError, TypeError):
                    result["report_data"] = None
        # Include image_docs from the memo's session so frontend can inject them
        image_docs = []
        session_id = (result.get("memo_session_id") or "")
        if session_id:
            img_rows = conn.execute(
                """SELECT id, file_name, doc_category, file_type, file_path
                   FROM memo_documents
                   WHERE memo_session_id=? AND file_type IN ('.png','.jpg','.jpeg','.gif','.webp')""",
                (session_id,),
            ).fetchall()
            for ir in img_rows:
                fpath = Path(ir["file_path"]) if ir["file_path"] else None
                if fpath and fpath.exists():
                    image_docs.append({
                        "doc_id": ir["id"],
                        "file_name": ir["file_name"],
                        "category": ir["doc_category"],
                        "file_type": ir["file_type"],
                        "section_targets": _image_category_to_sections(ir["doc_category"], ir["file_name"]),
                    })
        result["image_docs"] = image_docs
        return result
    finally:
        conn.close()


@router.delete("/history/{memo_id}")
async def delete_memo(memo_id: int, user: CurrentUser = Depends(get_current_user)):
    conn = get_db()
    try:
        # Shared team library: any member may delete (per team decision).
        conn.execute("DELETE FROM generated_memos WHERE id=?", (memo_id,))
        conn.commit()
        return {"ok": True}
    finally:
        conn.close()


class _MemoRenameRequest(BaseModel):
    title: str = ""


@router.patch("/history/{memo_id}/title")
async def rename_memo(
    memo_id: int,
    body: _MemoRenameRequest,
    user: CurrentUser = Depends(get_current_user),
):
    """Rename an IC memo (any teammate can rename — memos are shared). Pass
    an empty title to reset to the default 'Investment Memo' label."""
    new_title = (body.title or "").strip()[:200]
    conn = get_db()
    try:
        cur = conn.execute(
            "UPDATE generated_memos SET custom_title=? WHERE id=?",
            (new_title, memo_id),
        )
        conn.commit()
        if (cur.rowcount or 0) == 0:
            raise HTTPException(404, "Memo not found.")
    finally:
        conn.close()
    return {"ok": True, "id": memo_id, "custom_title": new_title}


# ═══════════════════════════════════════════════════════════════════════════════
#  DOCX EXPORT
# ═══════════════════════════════════════════════════════════════════════════════

@router.get("/history/{memo_id}/docx")
async def export_memo_docx(memo_id: int, token: Optional[str] = Query(None)):
    """Export memo as .docx with embedded data room images. Supports token via query param."""
    if not token:
        raise HTTPException(status_code=401, detail="Token required")
    payload = decode_token(token)
    if not payload:
        raise HTTPException(status_code=401, detail="Invalid or expired token")
    user = CurrentUser(uid=int(payload["sub"]), username=payload["user"], role=payload["role"])
    conn = get_db()
    try:
        # Shared team library: any signed-in member can export a teammate's memo.
        row = conn.execute(
            "SELECT company_name, memo_markdown, report_id, memo_session_id FROM generated_memos WHERE id=?",
            (memo_id,),
        ).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Memo not found")

        # Fetch image documents scoped to THIS memo's session only
        image_docs = []
        memo_session_id = row["memo_session_id"] or ""
        if memo_session_id:
            # Images belong to the memo's own session — scope by session, not the
            # exporter, so a teammate's docx export still embeds the right images.
            session_rows = conn.execute(
                """SELECT id, file_name, file_path, file_type, doc_category
                   FROM memo_documents
                   WHERE memo_session_id=? AND file_type IN ('.png','.jpg','.jpeg','.gif','.webp')
                   ORDER BY uploaded_at""",
                (memo_session_id,),
            ).fetchall()
        else:
            # Fallback for older memos without a session_id: use most recent images for this user
            session_rows = conn.execute(
                """SELECT id, file_name, file_path, file_type, doc_category
                   FROM memo_documents
                   WHERE owner_id=? AND file_type IN ('.png','.jpg','.jpeg','.gif','.webp')
                   ORDER BY uploaded_at DESC LIMIT 20""",
                (user.id,),
            ).fetchall()
        for sr in session_rows:
            fpath = Path(sr["file_path"]) if sr["file_path"] else None
            if fpath and fpath.exists():
                image_docs.append({
                    "file_path": str(fpath),
                    "file_name": sr["file_name"],
                    "category": sr["doc_category"],
                    "section_targets": _image_category_to_sections(sr["doc_category"], sr["file_name"]),
                })
    finally:
        conn.close()

    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        raise HTTPException(status_code=500, detail="python-docx not installed")

    doc = Document()

    # Style setup
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)

    # Title
    title = doc.add_heading(f"Investment Memorandum: {row['company_name']}", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Subtitle
    sub = doc.add_paragraph("VoLo Earth Ventures — Confidential")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.runs[0].font.color.rgb = RGBColor(0x58, 0x60, 0x69)
    sub.runs[0].font.size = Pt(12)

    doc.add_paragraph("")  # spacer

    # Build section → image mapping (one image per section max)
    section_images = {}  # section_title_keyword → file_path
    used_sections = set()
    SECTION_KEYWORDS = {
        'company_overview': 'company overview',
        'market': 'market',
        'business_model': 'business model',
        'team': 'team',
        'traction': 'traction',
        'competitive_position': 'competitive position',
        'carbon_impact': 'carbon impact',
        'technology_ip_moat': 'technology',
        'financing_overview': 'financing overview',
    }
    for img in image_docs:
        for sec_key in img.get("section_targets", []):
            kw = SECTION_KEYWORDS.get(sec_key)
            if kw and kw not in used_sections:
                section_images[kw] = img
                used_sections.add(kw)
                break

    # Parse markdown into doc, injecting images after section headings
    lines = row["memo_markdown"].split('\n')
    current_section = None
    image_injected_for_section = set()
    para_count_in_section = 0

    for line in lines:
        line = line.rstrip()
        if line.startswith('#### '):
            doc.add_heading(line[5:], level=4)
        elif line.startswith('### '):
            doc.add_heading(line[4:], level=3)
        elif line.startswith('## '):
            heading_text = line[3:]
            doc.add_heading(heading_text, level=2)
            current_section = heading_text.lower()
            para_count_in_section = 0
        elif line.startswith('# '):
            heading_text = line[2:]
            doc.add_heading(heading_text, level=1)
            current_section = heading_text.lower()
            para_count_in_section = 0
        elif line.startswith('- '):
            doc.add_paragraph(line[2:], style='List Bullet')
        elif re.match(r'^\d+\. ', line):
            doc.add_paragraph(re.sub(r'^\d+\. ', '', line), style='List Number')
        elif line.strip():
            p = doc.add_paragraph()
            _add_formatted_text(p, line)
            para_count_in_section += 1

            # Inject image after the first substantial paragraph of matching sections
            if current_section and para_count_in_section == 1:
                for kw, img in section_images.items():
                    if kw in current_section and kw not in image_injected_for_section:
                        try:
                            img_path = Path(img["file_path"])
                            if img_path.exists():
                                doc.add_picture(str(img_path), width=Inches(5.5))
                                last_para = doc.paragraphs[-1]
                                last_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                # Caption
                                cap_name = img["file_name"].rsplit('.', 1)[0].replace('_', ' ').replace('-', ' ').title()
                                cap_cat = img["category"].replace('_', ' ').title()
                                cap = doc.add_paragraph(f"{cap_cat}: {cap_name}")
                                cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                cap_run = cap.runs[0] if cap.runs else cap.add_run("")
                                cap_run.font.size = Pt(9)
                                cap_run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
                                cap_run.font.italic = True
                                image_injected_for_section.add(kw)
                        except Exception as img_err:
                            logger.warning(f"Could not embed image {img['file_name']}: {img_err}")
                        break

    # Save to temp file
    tmp = tempfile.NamedTemporaryFile(suffix='.docx', delete=False)
    doc.save(tmp.name)
    tmp.close()

    safe_name = re.sub(r'[^\w\-.]', '_', row["company_name"] or "memo")
    return FileResponse(
        tmp.name,
        media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        filename=f"Investment_Memo_{safe_name}.docx",
    )


def _add_formatted_text(paragraph, text):
    """Parse bold/italic markdown inline and add runs to paragraph."""
    # Split on bold markers
    parts = re.split(r'(\*\*\*.+?\*\*\*|\*\*.+?\*\*|\*.+?\*)', text)
    for part in parts:
        if part.startswith('***') and part.endswith('***'):
            run = paragraph.add_run(part[3:-3])
            run.bold = True
            run.italic = True
        elif part.startswith('**') and part.endswith('**'):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        elif part.startswith('*') and part.endswith('*'):
            run = paragraph.add_run(part[1:-1])
            run.italic = True
        else:
            paragraph.add_run(part)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION-LEVEL REVIEW & REVISION
# ═══════════════════════════════════════════════════════════════════════════════

class SectionReviseRequest(BaseModel):
    memo_id: int
    section_key: str
    instructions: str  # e.g. "make the tone more positive", "add more detail on competitors"
    model: Optional[str] = None
    additional_doc_text: Optional[str] = None  # optional new document content to incorporate


class SectionDirectEditRequest(BaseModel):
    memo_id: int
    section_key: str
    new_text: str  # the manually edited markdown text


@router.post("/revise-section")
async def revise_section(
    req: SectionReviseRequest,
    user: CurrentUser = Depends(get_current_user),
):
    """
    Revise a single memo section using LLM based on reviewer instructions.

    Takes the current section text + reviewer instructions and regenerates
    just that section, preserving all other sections. Records the revision
    in memo_revisions for change tracking.
    """
    conn = get_db()
    try:
        # Memos are a shared team library: any signed-in VoLo member can edit a
        # teammate's memo (matches the firm-wide GET). Attribution is preserved
        # via memo_revisions.revised_by. (Owner-scoping here caused a spurious
        # "Memo not found" when a non-owner tried to save edits.)
        row = conn.execute(
            "SELECT * FROM generated_memos WHERE id=?",
            (req.memo_id,),
        ).fetchone()
        if not row:
            raise HTTPException(404, "Memo not found")

        # Load section texts
        try:
            sections = json.loads(row["sections_json"] or "{}")
        except (json.JSONDecodeError, TypeError):
            sections = {}

        if req.section_key not in sections:
            raise HTTPException(404, f"Section '{req.section_key}' not found in this memo")

        old_text = sections[req.section_key]

        # Find the section definition for title and guidance
        section_def = None
        for s in MEMO_SECTIONS:
            if s["key"] == req.section_key:
                section_def = s
                break
        if not section_def:
            raise HTTPException(400, f"Unknown section key: {req.section_key}")

        # Get report context if available
        report_context = ""
        if row["report_id"]:
            # Reports are shared firm-wide — drop owner scoping so a memo
            # built off a teammate's report can still be revised by section.
            rpt_row = conn.execute(
                "SELECT report_json FROM deal_reports WHERE id=?",
                (row["report_id"],),
            ).fetchone()
            if rpt_row:
                try:
                    rpt = json.loads(rpt_row["report_json"]) if isinstance(rpt_row["report_json"], str) else rpt_row["report_json"]
                    report_context = _build_report_context_from_parsed(rpt, section_def)
                except Exception:
                    pass

        # Resolve model
        model_prefs = get_model_preferences(user.id)
        model = req.model or model_prefs.get("memo", "claude-sonnet-4-20250514")

        # Build the revision prompt
        def _run_revision():
            import anthropic
            api_key = os.environ.get("ANTHROPIC_API_KEY", "")
            if not api_key:
                raise ValueError("ANTHROPIC_API_KEY not set")
            client = anthropic.Anthropic(api_key=api_key)

            system_prompt = f"""{_SECTION_WRITER_SYSTEM}

REVISION MODE: You are revising an existing section of an investment memo.
You will receive the current section text and specific revision instructions from a reviewer.
Apply the requested changes while maintaining consistency with the rest of the memo.
Preserve citations [n] and [RVM] references. Keep the same general structure unless told otherwise.
Do NOT include the section title as a header — it will be added automatically."""

            user_parts = [
                f"# SECTION: {section_def['title']}",
                f"## Section Purpose\n{section_def['guidance']}",
                f"## Current Section Text\n{old_text}",
                f"## Reviewer Instructions\n{req.instructions}",
            ]
            if report_context:
                user_parts.append(f"## Quantitative Report Data [RVM]\n{report_context}")
            if req.additional_doc_text:
                user_parts.append(f"## Additional Document Content\n{req.additional_doc_text}")

            user_parts.append(f"\nRevise the '{section_def['title']}' section according to the reviewer's instructions. Maintain IC-quality writing.")

            response = client.messages.create(
                model=model,
                max_tokens=3000,
                system=system_prompt,
                messages=[{"role": "user", "content": "\n\n".join(user_parts)}],
            )
            text = "".join(b.text for b in response.content if b.type == "text")
            tokens_in = response.usage.input_tokens if response.usage else 0
            tokens_out = response.usage.output_tokens if response.usage else 0
            return text.strip(), tokens_in, tokens_out

        new_text, tokens_in, tokens_out = await asyncio.to_thread(_run_revision)

        # Update sections_json
        sections[req.section_key] = new_text
        sections_json_str = json.dumps(sections)

        # Reassemble full memo markdown and HTML
        memo_md, memo_html = _reassemble_memo(sections, row["company_name"])

        # Save updated memo
        conn.execute(
            """UPDATE generated_memos
               SET sections_json=?, memo_markdown=?, memo_html=?
               WHERE id=?""",
            (sections_json_str, memo_md, memo_html, req.memo_id),
        )

        # Record revision
        conn.execute(
            """INSERT INTO memo_revisions
               (memo_id, section_key, revision_type, old_text, new_text,
                instructions, revised_by, model_used, tokens_in, tokens_out)
               VALUES (?, ?, 'llm', ?, ?, ?, ?, ?, ?, ?)""",
            (req.memo_id, req.section_key, old_text, new_text,
             req.instructions, user.username, model, tokens_in, tokens_out),
        )
        conn.commit()

        # Convert just the new section to HTML for immediate frontend update
        section_html = _markdown_to_html(f"## {section_def['title']}\n\n{new_text}\n")

        return {
            "section_key": req.section_key,
            "new_text": new_text,
            "section_html": section_html,
            "model_used": model,
            "tokens_in": tokens_in,
            "tokens_out": tokens_out,
            "memo_html": memo_html,
        }
    finally:
        conn.close()


@router.post("/edit-section")
async def edit_section_direct(
    req: SectionDirectEditRequest,
    user: CurrentUser = Depends(get_current_user),
):
    """
    Save a direct manual edit to a memo section.
    Records the change in memo_revisions for tracking.
    """
    conn = get_db()
    try:
        # Memos are a shared team library: any signed-in VoLo member can edit a
        # teammate's memo (matches the firm-wide GET). Attribution is preserved
        # via memo_revisions.revised_by. (Owner-scoping here caused a spurious
        # "Memo not found" when a non-owner tried to save edits.)
        row = conn.execute(
            "SELECT * FROM generated_memos WHERE id=?",
            (req.memo_id,),
        ).fetchone()
        if not row:
            raise HTTPException(404, "Memo not found")

        try:
            sections = json.loads(row["sections_json"] or "{}")
        except (json.JSONDecodeError, TypeError):
            sections = {}

        if req.section_key not in sections:
            raise HTTPException(404, f"Section '{req.section_key}' not found")

        old_text = sections[req.section_key]
        sections[req.section_key] = req.new_text
        sections_json_str = json.dumps(sections)

        memo_md, memo_html = _reassemble_memo(sections, row["company_name"])

        conn.execute(
            """UPDATE generated_memos
               SET sections_json=?, memo_markdown=?, memo_html=?
               WHERE id=?""",
            (sections_json_str, memo_md, memo_html, req.memo_id),
        )

        conn.execute(
            """INSERT INTO memo_revisions
               (memo_id, section_key, revision_type, old_text, new_text,
                instructions, revised_by)
               VALUES (?, ?, 'manual', ?, ?, '', ?)""",
            (req.memo_id, req.section_key, old_text, req.new_text, user.username),
        )
        conn.commit()

        section_def = next((s for s in MEMO_SECTIONS if s["key"] == req.section_key), None)
        title = section_def["title"] if section_def else req.section_key
        section_html = _markdown_to_html(f"## {title}\n\n{req.new_text}\n")

        return {
            "section_key": req.section_key,
            "new_text": req.new_text,
            "section_html": section_html,
            "memo_html": memo_html,
        }
    finally:
        conn.close()


@router.get("/revisions/{memo_id}")
async def get_revisions(memo_id: int, user: CurrentUser = Depends(get_current_user)):
    """Get revision history for a memo, grouped by section."""
    conn = get_db()
    try:
        rows = conn.execute(
            """SELECT r.* FROM memo_revisions r
               WHERE r.memo_id=?
               ORDER BY r.created_at DESC""",
            (memo_id,),
        ).fetchall()
        return [dict(r) for r in rows]
    finally:
        conn.close()


@router.get("/revisions/{memo_id}/{section_key}")
async def get_section_revisions(
    memo_id: int, section_key: str,
    user: CurrentUser = Depends(get_current_user),
):
    """Get revision history for a specific section."""
    conn = get_db()
    try:
        rows = conn.execute(
            """SELECT r.* FROM memo_revisions r
               WHERE r.memo_id=? AND r.section_key=?
               ORDER BY r.created_at DESC""",
            (memo_id, section_key),
        ).fetchall()
        return [dict(r) for r in rows]
    finally:
        conn.close()


def _strip_section_heading(text: str, section_title: str) -> str:
    """Remove any leading markdown heading the LLM added to a section body."""
    if not text:
        return text
    lines = text.split("\n")
    first = lines[0].lstrip()
    if first.startswith("#"):
        heading_text = first.lstrip("#").strip().lower()
        title_lower  = section_title.lower()
        if (heading_text in title_lower or title_lower in heading_text
                or heading_text.replace(" ", "") == title_lower.replace(" ", "")):
            lines = lines[1:]
            while lines and not lines[0].strip():
                lines = lines[1:]
    return "\n".join(lines)


def _reassemble_memo(sections: dict, company_name: str) -> tuple:
    """Reassemble full memo markdown + HTML from section texts dict."""
    memo_parts = [f"# Investment Memorandum: {company_name or 'Deal Analysis'}\n"]
    memo_parts.append(f"*VoLo Earth Ventures — Confidential*\n")

    for section in MEMO_SECTIONS:
        text = sections.get(section["key"], "")
        if text:
            text = _strip_section_heading(text, section["title"])
            memo_parts.append(f"## {section['title']}\n\n{text}\n")

    memo_md = "\n".join(memo_parts)
    memo_html = _markdown_to_html(memo_md)
    return memo_md, memo_html


def _build_report_context_from_parsed(rpt: dict, section_def: dict) -> str:
    """Build report context string from parsed report JSON for a specific section."""
    parts = []
    report_fields = section_def.get("report_fields", [])

    # Always surface the current round's deal terms first so any section being
    # revised — especially the synthesis sections (Investment Overview,
    # Recommendation) — anchors on the CURRENT round, not a prior one.
    parts.extend(_format_round_terms(rpt.get("deal_overview", {})))

    if "simulation" in report_fields or "inputs" in report_fields:
        sim = rpt.get("simulation", {})
        if sim:
            parts.append(f"Monte Carlo simulation ({sim.get('n_simulations', 'N/A')} paths):")
            moic = sim.get("moic_unconditional", {})
            if moic:
                parts.append(f"  MOIC: mean={moic.get('mean', 'N/A')}x, median={moic.get('median', 'N/A')}x")
            prob = sim.get("probability", {})
            if prob:
                parts.append(f"  P(>3x)={prob.get('gt_3x', 'N/A')}, P(loss)={prob.get('total_loss', 'N/A')}")

    if "carbon" in report_fields:
        carbon = rpt.get("carbon_impact", {})
        outputs = carbon.get("outputs", {})
        intermediates = carbon.get("intermediates", {})
        ci = carbon.get("carbon_inputs", {})
        ov = rpt.get("deal_overview", {})
        risk_div = carbon.get("risk_divisor_used", "N/A")
        if outputs:
            parts.append("=== CARBON IMPACT (RVM Section 4) ===")
            # Inputs
            if ci:
                parts.append(f"Displacement Chain Inputs:")
                parts.append(f"  Unit definition: {ci.get('unit_definition', 'N/A')}")
                parts.append(f"  Unit service life: {ci.get('unit_service_life_yrs', 'N/A')} years")
                parts.append(f"  Displaced resource: {ci.get('displaced_resource', 'N/A')}")
                parts.append(f"  Baseline lifetime production: {ci.get('baseline_lifetime_prod', 'N/A')} {ci.get('specific_production_units', '')}")
                parts.append(f"  Range improvement (fraction displaced): {ci.get('range_improvement', 'N/A')}")
                parts.append(f"  Commercial launch year: {ci.get('commercial_launch_yr', 'N/A')}")
                vols = ci.get('year_volumes', [])
                if vols:
                    parts.append(f"  10-year volume forecast (units): {', '.join(str(round(v,1)) for v in vols[:10])}")
                if ci.get('emb_displaced_resource'):
                    parts.append(f"  Embodied carbon resource: {ci.get('emb_displaced_resource')} "
                                 f"(baseline: {ci.get('emb_baseline_production', 0)}, "
                                 f"range_improvement: {ci.get('emb_range_improvement', 0)})")
            # Intermediates
            if intermediates:
                parts.append(f"Calculated Intermediates:")
                parts.append(f"  JD (displaced vol per unit): {intermediates.get('jd', 'N/A')}")
                ci_ser = intermediates.get('operating_ci_series', [])
                if ci_ser:
                    parts.append(f"  Operating CI series (tCO2/unit-prod, Y1-Y10): {', '.join(f'{v:.4f}' for v in ci_ser[:10])}")
                ann_op = intermediates.get('annual_operating', [])
                if ann_op:
                    parts.append(f"  Annual operating impact (tCO2): {', '.join(f'{v:.1f}' for v in ann_op[:10])}")
                ann_emb = intermediates.get('annual_embodied', [])
                if any(v and v > 0 for v in (ann_emb or [])):
                    parts.append(f"  Annual embodied impact (tCO2): {', '.join(f'{v:.1f}' for v in ann_emb[:10])}")
                ann_lc = intermediates.get('annual_lifecycle', [])
                if ann_lc:
                    parts.append(f"  Annual lifecycle impact (tCO2): {', '.join(f'{v:.1f}' for v in ann_lc[:10])}")
                parts.append(f"  Total operating: {intermediates.get('total_operating', 'N/A')} tCO2")
                parts.append(f"  Total embodied: {intermediates.get('total_embodied', 'N/A')} tCO2")
                parts.append(f"  Total lifecycle: {intermediates.get('total_lifecycle', 'N/A')} tCO2")
            # Outputs
            parts.append(f"VoLo-Level Outputs:")
            parts.append(f"  Total lifecycle tCO2 (company): {outputs.get('company_tonnes', 'N/A')}")
            parts.append(f"  VoLo pro-rata tCO2 ({ov.get('entry_ownership_pct', '?')}% ownership): {outputs.get('volo_prorata', 'N/A')}")
            parts.append(f"  Risk divisor: {risk_div} ({carbon.get('risk_divisor_source', '')})")
            parts.append(f"  Risk-adjusted tCO2: {outputs.get('volo_risk_adj', 'N/A')}")
            parts.append(f"  t/$ (unadjusted): {outputs.get('tonnes_per_dollar', 'N/A')}")
            parts.append(f"  t/$ (risk-adjusted): {outputs.get('risk_adj_tpd', 'N/A')}")

    if "check_optimization" in report_fields or "portfolio_impact" in report_fields:
        ps = rpt.get("position_sizing", {})
        gso = (ps.get("grid_search") or {}).get("optimal") or {}
        constraints = ps.get("fund_constraints") or {}
        kelly = ps.get("kelly_reference") or {}
        if ps.get("has_data"):
            parts.append(f"Check size optimization:")
            parts.append(f"  Fund-optimized check: ${gso.get('check_m', 'N/A')}M (ownership: {gso.get('ownership_pct', 'N/A')}%)")
            parts.append(f"  Fund P50 impact: {gso.get('fund_p50_pct_chg', 'N/A')}")
            parts.append(f"  Fund constraints: min=${constraints.get('min_check_m', 'N/A')}M, max=${constraints.get('max_check_m', 'N/A')}M, fund_size=${constraints.get('fund_size_m', 'N/A')}M")
            parts.append(f"  Kelly full=${kelly.get('optimal_check_m', 'N/A')}M, half-Kelly=${kelly.get('half_kelly_check_m', 'N/A')}M")
        fops = rpt.get("followon_position_sizing", {})
        if fops.get("has_data"):
            _pr = fops.get("prorata", {}) or {}
            _comb = fops.get("combined", {}) or {}
            _bl = _comb.get("blended_stats") or {}
            _combown = _comb.get("combined_ownership_pct", _comb.get("total_ownership_pct_approx", "N/A"))
            parts.append(f"Follow-on sizing (#{fops.get('followon_number', '?')} follow-on; basis: {fops.get('sizing_method', 'N/A')}):")
            parts.append(f"  Prior deployed: ${_comb.get('total_prior_m', 'N/A')}M; recommended follow-on: ${fops.get('recommended_followon_check_m', 'N/A')}M; combined invested: ${_comb.get('total_invested_m', 'N/A')}M")
            parts.append(f"  Diluted ownership entering: {_pr.get('current_diluted_ownership_pct', 'N/A')}% -> combined post-round: {_combown}%; pro-rata check: ${_pr.get('pro_rata_check_m', 'N/A')}M ({_pr.get('recommended_vs_pro_rata', 'N/A')})")
            if _bl.get("blended_moic_p50") is not None:
                parts.append(f"  Blended MOIC: P10 {_bl.get('blended_moic_p10')}x / P50 {_bl.get('blended_moic_p50')}x / P90 {_bl.get('blended_moic_p90')}x")
        pi = rpt.get("portfolio_impact", {})
        if pi.get("has_data"):
            parts.append(f"Portfolio impact (Section 3):")
            parts.append(f"  TVPI base={pi.get('tvpi_base_mean', 'N/A')}x -> with deal={pi.get('tvpi_new_mean', 'N/A')}x (lift={pi.get('tvpi_mean_lift', 'N/A')}x)")
            parts.append(f"  TVPI With Deal by percentile: P50 {pi.get('tvpi_new_p50', 'N/A')}x, P75 {pi.get('tvpi_new_p75', 'N/A')}x, TOP DECILE P90 {pi.get('tvpi_new_p90', 'N/A')}x")
            parts.append("  Carta benchmark (mature): P50 ~1.65x, P75 ~2.2x, TOP DECILE P90 ~4.0x. The TVPI Mean is the expected case (below the top decile); the fund's P90 is the like-for-like top-decile comparison — never compare the mean to Carta's top decile.")
            parts.append(f"  IRR base={pi.get('irr_base_mean', 'N/A')} -> with deal={pi.get('irr_new_mean', 'N/A')} (lift={pi.get('irr_mean_lift', 'N/A')})")
            parts.append(f"  N committed deals: {pi.get('n_committed_deals', 0)}, narrative: {pi.get('narrative', '')}")

    if "valuation_comps" in report_fields:
        vc = rpt.get("valuation_context", {})
        if vc:
            parts.append(f"Valuation comps: IPO EV/EBITDA mean={vc.get('ipo_ev_ebitda_mean', 'N/A')}x")

    if "adoption" in report_fields:
        adopt = rpt.get("adoption_analysis", {})
        scurve = adopt.get("scurve", {})
        if scurve:
            parts.append(f"S-curve: bass_p={scurve.get('bass_p_mean', 'N/A')}, bass_q={scurve.get('bass_q_mean', 'N/A')}")

    return "\n".join(parts) if parts else ""
