"""
Google Drive integration — per-user OAuth.

Each user connects their own Google account once via the OAuth flow; the app
then reads Drive on their behalf, inheriting whatever per-folder permissions
they already have. Refresh tokens are encrypted at rest with a Fernet key
from GOOGLE_TOKEN_ENCRYPTION_KEY.

OAuth flow:
  GET  /api/drive/oauth/authorize        — redirects user to Google consent screen
  GET  /api/drive/oauth/callback         — Google redirects back here with a code
  POST /api/drive/disconnect             — revoke + delete this user's credentials
  GET  /api/drive/connection-status      — { connected, google_email, connected_at }

Library / sync:
  POST   /api/drive/libraries              — create / link a deal library to a Drive folder
  GET    /api/drive/libraries              — list all libraries
  GET    /api/drive/libraries/{id}         — get library details + documents
  DELETE /api/drive/libraries/{id}         — remove a library
  POST   /api/drive/libraries/{id}/sync    — sync from Google Drive (pull new/changed files)
  GET    /api/drive/libraries/{id}/documents — list documents in a library
  PUT    /api/drive/documents/{id}/category — update document category
"""

import base64
import hashlib
import io
import json
import logging
import os
import re
import tempfile
import time
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Optional

# Google's OAuth library can return scopes in a different order than requested
# (e.g. it normalizes "openid" / "email"), which causes oauthlib to raise
# "Scope has changed" on token exchange. This env flag tells oauthlib to
# accept the normalized scopes. Must be set BEFORE oauthlib is imported, and
# we use direct assignment (not setdefault) so it can't be silently overridden
# by an incomplete env config.
os.environ["OAUTHLIB_RELAX_TOKEN_SCOPE"] = "1"

import jwt
from fastapi import APIRouter, Depends, HTTPException, Query, Request
from fastapi.responses import JSONResponse, RedirectResponse
from pydantic import BaseModel

from ..auth import CurrentUser, SECRET_KEY, decode_token, get_current_user, get_optional_user
from ..database import get_db

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/drive", tags=["drive"])

# ── OAuth configuration ───────────────────────────────────────────────────────
# Scopes requested at consent time. drive.readonly + email + openid lets us
# (a) read files the user can see, (b) record which Google account connected.
_OAUTH_SCOPES = [
    "https://www.googleapis.com/auth/drive.readonly",
    "https://www.googleapis.com/auth/userinfo.email",
    "openid",
]

# Google OAuth requires returned scopes to match exactly. The userinfo and
# openid scopes can come back with full URL prefixes; we accept both forms.
_REQUIRED_SCOPE = "https://www.googleapis.com/auth/drive.readonly"

# State token TTL — covers the round-trip to Google + user clicking through
_OAUTH_STATE_TTL_SECONDS = 600


def _oauth_redirect_uri(request: Optional[Request] = None) -> str:
    """The OAuth redirect URI must EXACTLY match what's registered in Google
    Cloud Console. We prefer the explicit env var; if it's not set we build
    one from the incoming request.

    Reverse proxies (Replit, Cloudflare, etc.) sometimes chain x-forwarded-*
    headers, producing values like "https,https" or hosts with ":443" port
    suffixes. We normalize aggressively so the URL we send on /authorize is
    byte-identical to the URL we send on /callback's token exchange."""
    explicit = os.environ.get("GOOGLE_OAUTH_REDIRECT_URI", "").strip()
    if explicit:
        return explicit
    if request is not None:
        proto_raw = request.headers.get("x-forwarded-proto") or request.url.scheme
        # Chained proxies can produce "https,https"; first value is canonical
        scheme = proto_raw.split(",")[0].strip().lower() or "https"
        host_raw = request.headers.get("x-forwarded-host") or request.url.netloc
        host = host_raw.split(",")[0].strip()
        # Strip default ports — Google compares strings, ":443" vs none differs
        if scheme == "https" and host.endswith(":443"):
            host = host[:-4]
        if scheme == "http" and host.endswith(":80"):
            host = host[:-3]
        return f"{scheme}://{host}/api/drive/oauth/callback"
    raise HTTPException(
        status_code=500,
        detail="GOOGLE_OAUTH_REDIRECT_URI not configured and no request context.",
    )


def _oauth_client_config(request: Optional[Request] = None) -> dict:
    """Build the client_config dict that google-auth-oauthlib's Flow expects."""
    client_id = os.environ.get("GOOGLE_OAUTH_CLIENT_ID", "").strip()
    client_secret = os.environ.get("GOOGLE_OAUTH_CLIENT_SECRET", "").strip()
    if not client_id or not client_secret:
        raise HTTPException(
            status_code=500,
            detail=(
                "Google Drive integration is not configured. "
                "Set GOOGLE_OAUTH_CLIENT_ID and GOOGLE_OAUTH_CLIENT_SECRET in Replit Secrets."
            ),
        )
    return {
        "web": {
            "client_id": client_id,
            "client_secret": client_secret,
            "auth_uri": "https://accounts.google.com/o/oauth2/auth",
            "token_uri": "https://oauth2.googleapis.com/token",
            "redirect_uris": [_oauth_redirect_uri(request)],
        }
    }


# ── Refresh-token encryption (Fernet, key from env) ───────────────────────────
_FERNET_INSTANCE = None


def _get_fernet():
    """Lazily build a Fernet cipher from GOOGLE_TOKEN_ENCRYPTION_KEY.
    The key must be a 32-byte url-safe base64 string (Fernet.generate_key())."""
    global _FERNET_INSTANCE
    if _FERNET_INSTANCE is not None:
        return _FERNET_INSTANCE
    try:
        from cryptography.fernet import Fernet
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="cryptography package not installed. Run: pip install cryptography",
        )
    key = os.environ.get("GOOGLE_TOKEN_ENCRYPTION_KEY", "").strip()
    if not key:
        raise HTTPException(
            status_code=500,
            detail=(
                "GOOGLE_TOKEN_ENCRYPTION_KEY is not set. Generate one with: "
                "python -c \"from cryptography.fernet import Fernet; "
                "print(Fernet.generate_key().decode())\" — then add it to Replit Secrets."
            ),
        )
    try:
        _FERNET_INSTANCE = Fernet(key.encode() if isinstance(key, str) else key)
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"GOOGLE_TOKEN_ENCRYPTION_KEY is invalid: {e}",
        )
    return _FERNET_INSTANCE


def _encrypt_token(plaintext: str) -> str:
    """Encrypt a refresh token for at-rest storage in the DB."""
    f = _get_fernet()
    return f.encrypt(plaintext.encode("utf-8")).decode("ascii")


def _decrypt_token(ciphertext: str) -> str:
    """Reverse of _encrypt_token."""
    f = _get_fernet()
    return f.decrypt(ciphertext.encode("ascii")).decode("utf-8")


# ── Per-user credential storage ───────────────────────────────────────────────

def _save_user_credentials(user_id: int, refresh_token: str,
                           google_email: str, scopes: list) -> None:
    """Upsert encrypted refresh token + metadata for a user."""
    enc = _encrypt_token(refresh_token)
    scopes_str = " ".join(scopes) if scopes else ""
    conn = get_db()
    try:
        existing = conn.execute(
            "SELECT id FROM user_drive_credentials WHERE user_id=?", (user_id,)
        ).fetchone()
        if existing:
            conn.execute(
                """UPDATE user_drive_credentials
                   SET google_email=?, refresh_token_enc=?, scopes=?,
                       connected_at=datetime('now')
                   WHERE user_id=?""",
                (google_email, enc, scopes_str, user_id),
            )
        else:
            conn.execute(
                """INSERT INTO user_drive_credentials
                   (user_id, google_email, refresh_token_enc, scopes)
                   VALUES (?, ?, ?, ?)""",
                (user_id, google_email, enc, scopes_str),
            )
        conn.commit()
    finally:
        conn.close()


def _delete_user_credentials(user_id: int) -> bool:
    """Remove a user's stored Drive credentials. Returns True if a row was deleted."""
    conn = get_db()
    try:
        cur = conn.execute(
            "DELETE FROM user_drive_credentials WHERE user_id=?", (user_id,)
        )
        conn.commit()
        return (cur.rowcount or 0) > 0
    finally:
        conn.close()


def _get_user_drive_status(user_id: int) -> dict:
    """Return connection status for a user (no secrets)."""
    conn = get_db()
    try:
        row = conn.execute(
            """SELECT google_email, connected_at, last_used_at, scopes
               FROM user_drive_credentials WHERE user_id=?""",
            (user_id,),
        ).fetchone()
        if not row:
            return {"connected": False}
        return {
            "connected": True,
            "google_email": row["google_email"],
            "connected_at": row["connected_at"],
            "last_used_at": row["last_used_at"],
            "scopes": row["scopes"],
        }
    finally:
        conn.close()


def _load_user_oauth_credentials(user_id: int):
    """Build a google.oauth2.credentials.Credentials object for this user.
    Returns None if the user has not connected Drive."""
    try:
        from google.oauth2.credentials import Credentials
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="google-auth not installed. Run: pip install google-auth google-auth-oauthlib google-api-python-client",
        )

    conn = get_db()
    try:
        row = conn.execute(
            "SELECT refresh_token_enc, scopes FROM user_drive_credentials WHERE user_id=?",
            (user_id,),
        ).fetchone()
    finally:
        conn.close()
    if not row:
        return None

    try:
        refresh_token = _decrypt_token(row["refresh_token_enc"])
    except Exception as e:
        logger.error(f"Failed to decrypt refresh token for user {user_id}: {e}")
        return None

    scopes = (row["scopes"] or "").split() or _OAUTH_SCOPES
    client_id = os.environ.get("GOOGLE_OAUTH_CLIENT_ID", "").strip()
    client_secret = os.environ.get("GOOGLE_OAUTH_CLIENT_SECRET", "").strip()

    return Credentials(
        token=None,                       # no current access token; will refresh
        refresh_token=refresh_token,
        token_uri="https://oauth2.googleapis.com/token",
        client_id=client_id,
        client_secret=client_secret,
        scopes=scopes,
    )


def _touch_last_used(user_id: int) -> None:
    """Update last_used_at after a successful Drive call."""
    conn = get_db()
    try:
        conn.execute(
            "UPDATE user_drive_credentials SET last_used_at=datetime('now') WHERE user_id=?",
            (user_id,),
        )
        conn.commit()
    finally:
        conn.close()


# ── Signed state parameter (CSRF protection + PKCE verifier transport) ────────

def _make_oauth_state(user_id: int, code_verifier: Optional[str] = None) -> str:
    """Sign a short-lived JWT carrying the user_id and (if PKCE is in use)
    the code_verifier. The OAuth callback recovers both from this token so
    it can identify the user AND complete the PKCE handshake without trusting
    query-string input."""
    payload = {
        "uid": user_id,
        "kind": "drive_oauth",
        "exp": datetime.now(timezone.utc) + timedelta(seconds=_OAUTH_STATE_TTL_SECONDS),
    }
    if code_verifier:
        payload["cv"] = code_verifier
    return jwt.encode(payload, SECRET_KEY, algorithm="HS256")


def _verify_oauth_state(state: str) -> tuple[int, Optional[str]]:
    """Validate the state token and return (user_id, code_verifier_or_None).
    Raises HTTPException on invalid/expired state."""
    try:
        payload = jwt.decode(state, SECRET_KEY, algorithms=["HS256"])
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=400, detail="OAuth flow expired. Please try again.")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=400, detail="Invalid OAuth state.")
    if payload.get("kind") != "drive_oauth":
        raise HTTPException(status_code=400, detail="Invalid OAuth state.")
    return int(payload["uid"]), payload.get("cv")


def _generate_pkce_verifier() -> str:
    """RFC 7636 PKCE code verifier — 43-128 chars from the unreserved set.
    We use 96 chars (~512 bits of entropy) which fits comfortably and is
    accepted by every server. Generated with secrets.token_urlsafe to ensure
    cryptographic randomness."""
    import secrets as _secrets
    return _secrets.token_urlsafe(72)[:96]

# ── Temp directory for downloaded files during extraction ─────────────────────
_DRIVE_CACHE_DIR = Path(__file__).resolve().parent.parent.parent / "data" / "drive_cache"
_DRIVE_CACHE_DIR.mkdir(parents=True, exist_ok=True)

# Supported file types for text extraction
_EXTRACTABLE_MIMES = {
    'application/pdf': '.pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
    'text/plain': '.txt',
    'text/csv': '.csv',
    'text/markdown': '.md',
    'application/json': '.json',
    'text/html': '.html',
}

# Google Docs/Sheets/Slides need to be exported
_GOOGLE_EXPORT_MIMES = {
    'application/vnd.google-apps.document': ('application/vnd.openxmlformats-officedocument.wordprocessingml.document', '.docx'),
    'application/vnd.google-apps.spreadsheet': ('application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', '.xlsx'),
    'application/vnd.google-apps.presentation': ('application/vnd.openxmlformats-officedocument.presentationml.presentation', '.pptx'),
}

# Category inference from filename / mime
_CATEGORY_PATTERNS = {
    'financial_model': ['financial', 'model', 'proforma', 'pro forma', 'forecast', 'projections', 'budget'],
    'pitch_deck': ['pitch', 'deck', 'presentation', 'investor update'],
    'term_sheet': ['term sheet', 'terms', 'loi', 'letter of intent'],
    'cap_table': ['cap table', 'captable', 'ownership'],
    'legal': ['legal', 'contract', 'agreement', 'nda', 'msa', 'incorporation', 'bylaws'],
    'ip_patent': ['patent', 'ip', 'intellectual property', 'trademark'],
    'customer_reference': ['customer', 'reference', 'testimonial', 'case study'],
    'market_research': ['market', 'research', 'analysis', 'tam', 'landscape', 'industry'],
    'technical_diligence': ['technical', 'diligence', 'architecture', 'engineering', 'tech'],
    'team_bios': ['team', 'bio', 'leadership', 'management', 'founders'],
    'board_materials': ['board', 'minutes', 'governance'],
}


def _infer_category(file_name: str) -> str:
    """Guess document category from filename."""
    name_lower = file_name.lower()
    for category, keywords in _CATEGORY_PATTERNS.items():
        if any(kw in name_lower for kw in keywords):
            return category
    return 'other'


def _parse_drive_folder_id(url_or_id: str) -> str:
    """Extract Google Drive folder ID from a URL or raw ID."""
    url_or_id = url_or_id.strip()
    # Direct ID (no slashes)
    if re.match(r'^[a-zA-Z0-9_-]{10,}$', url_or_id):
        return url_or_id
    # URL patterns
    m = re.search(r'/folders/([a-zA-Z0-9_-]+)', url_or_id)
    if m:
        return m.group(1)
    m = re.search(r'[?&]id=([a-zA-Z0-9_-]+)', url_or_id)
    if m:
        return m.group(1)
    raise ValueError(f"Could not parse Drive folder ID from: {url_or_id[:100]}")


def _get_drive_service(user_id: int):
    """Build an authenticated Drive v3 service for a specific user.
    Reads the user's stored refresh token, exchanges it for a fresh access
    token, and returns a service. Raises 401 with `code=drive_not_connected`
    if the user has never connected, or if the token has been revoked."""
    try:
        from googleapiclient.discovery import build
        from google.auth.transport.requests import Request as GoogleAuthRequest
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="Google API libraries not installed. Run: pip install google-api-python-client google-auth google-auth-oauthlib",
        )

    creds = _load_user_oauth_credentials(user_id)
    if creds is None:
        raise HTTPException(
            status_code=401,
            detail={
                "code": "drive_not_connected",
                "message": "Google Drive is not connected for this user. Connect it from the IC Memo tab.",
            },
        )

    # Refresh the access token. If the refresh itself fails (revoked / expired),
    # purge the stored credentials and tell the user to reconnect.
    try:
        creds.refresh(GoogleAuthRequest())
    except Exception as e:
        logger.warning(f"Drive token refresh failed for user {user_id}: {e}")
        _delete_user_credentials(user_id)
        raise HTTPException(
            status_code=401,
            detail={
                "code": "drive_token_revoked",
                "message": "Your Google Drive connection has expired or been revoked. Please reconnect.",
            },
        )

    _touch_last_used(user_id)
    return build("drive", "v3", credentials=creds, cache_discovery=False)


def _list_files_recursive(service, folder_id: str, path_prefix: str = "") -> list:
    """Recursively list all files in a Drive folder and subfolders.

    Important: by default the Drive API silently returns 0 results for files
    that live inside a Shared Drive (formerly Team Drive) or for some
    Workspace-shared folders. We explicitly enable shared-drive support on
    every list call so the same code works for personal folders, "Shared
    with me" items, and Shared Drive content. This is the standard fix
    for the "I see files in Drive but the app shows 0 docs" issue."""
    all_files = []
    page_token = None

    while True:
        # NOTE: do NOT pass corpora="allDrives" here. That parameter is for
        # cross-drive SEARCH queries; combined with `'folderId' in parents` it
        # over-restricts and filters out "Shared with me" content (which is
        # technically owned by another user's My Drive, not a Shared Drive
        # workspace). The correct combination for a folder-children query is
        # supportsAllDrives + includeItemsFromAllDrives WITHOUT a corpora
        # override — Drive auto-detects where the folder lives.
        resp = service.files().list(
            q=f"'{folder_id}' in parents and trashed=false",
            fields="nextPageToken, files(id, name, mimeType, size, modifiedTime)",
            pageSize=200,
            pageToken=page_token,
            supportsAllDrives=True,
            includeItemsFromAllDrives=True,
        ).execute()

        for f in resp.get("files", []):
            if f["mimeType"] == "application/vnd.google-apps.folder":
                # Recurse into subfolder
                subfolder_path = f"{path_prefix}{f['name']}/"
                all_files.extend(_list_files_recursive(service, f["id"], subfolder_path))
            else:
                f["subfolder_path"] = path_prefix
                all_files.append(f)

        page_token = resp.get("nextPageToken")
        if not page_token:
            break

    return all_files


def _download_file(service, file_info: dict) -> bytes:
    """Download a file from Drive. Handles Google Docs export."""
    mime = file_info["mimeType"]
    file_id = file_info["id"]

    # supportsAllDrives is needed when the file lives in a Shared Drive,
    # otherwise download silently fails. Harmless on personal-drive files.
    if mime in _GOOGLE_EXPORT_MIMES:
        export_mime, _ = _GOOGLE_EXPORT_MIMES[mime]
        request = service.files().export_media(fileId=file_id, mimeType=export_mime)
    else:
        request = service.files().get_media(fileId=file_id, supportsAllDrives=True)

    from googleapiclient.http import MediaIoBaseDownload
    buffer = io.BytesIO()
    downloader = MediaIoBaseDownload(buffer, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return buffer.getvalue()


def _extract_text(file_bytes: bytes, file_name: str, mime_type: str) -> str:
    """Extract text from file bytes. Reuses the same logic as memo.py."""
    # Determine the effective extension
    ext = Path(file_name).suffix.lower()
    if mime_type in _GOOGLE_EXPORT_MIMES:
        _, ext = _GOOGLE_EXPORT_MIMES[mime_type]

    text = ""
    try:
        if ext in ('.txt', '.md', '.csv', '.json', '.html'):
            text = file_bytes.decode('utf-8', errors='replace')[:100_000]

        elif ext == '.pdf':
            try:
                import fitz
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                pages = [page.get_text() for page in doc]
                text = "\n\n".join(pages)[:200_000]
                doc.close()
            except ImportError:
                text = "[PDF text extraction requires PyMuPDF]"

        elif ext == '.docx':
            try:
                from docx import Document
                doc = Document(io.BytesIO(file_bytes))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                text = "\n".join(paragraphs)[:200_000]
            except ImportError:
                text = "[DOCX text extraction requires python-docx]"

        elif ext in ('.xlsx', '.xls'):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
                parts = []
                for ws in wb.worksheets[:10]:
                    parts.append(f"=== Sheet: {ws.title} ===")
                    row_count = 0
                    for row in ws.iter_rows(values_only=True):
                        vals = [str(c) if c is not None else '' for c in row]
                        if any(v.strip() for v in vals):
                            parts.append('\t'.join(vals))
                            row_count += 1
                            if row_count > 300:
                                parts.append(f"... (truncated, {ws.max_row} total rows)")
                                break
                wb.close()
                text = "\n".join(parts)[:200_000]
            except ImportError:
                text = "[Excel text extraction requires openpyxl]"

        elif ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                parts = []
                for i, slide in enumerate(prs.slides):
                    parts.append(f"=== Slide {i+1} ===")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    parts.append(para.text.strip())
                text = "\n".join(parts)[:200_000]
            except ImportError:
                text = "[PPTX text extraction requires python-pptx]"

    except Exception as e:
        text = f"[Extraction error: {str(e)[:200]}]"

    return text


# ═══════════════════════════════════════════════════════════════════════════════
#  OAUTH FLOW — per-user Google Drive sign-in
# ═══════════════════════════════════════════════════════════════════════════════

@router.get("/oauth/authorize")
async def oauth_authorize(
    request: Request,
    token: Optional[str] = Query(None),
    user: Optional[CurrentUser] = Depends(get_optional_user),
):
    """Kick off the OAuth flow. Redirects the browser to Google's consent screen.

    This is a top-level browser navigation (not an XHR), so the browser does
    NOT send our Authorization header. We accept a `?token=` query param as
    a fallback — same pattern as /memo/history/{id}/docx and document view."""
    effective_user = user
    if not effective_user and token:
        payload = decode_token(token)
        if payload:
            effective_user = CurrentUser(uid=int(payload["sub"]), username=payload["user"], role=payload["role"])
    if not effective_user:
        raise HTTPException(status_code=401, detail="Authentication required")

    try:
        from google_auth_oauthlib.flow import Flow
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="google-auth-oauthlib not installed. Run: pip install google-auth-oauthlib",
        )

    redirect_uri = _oauth_redirect_uri(request)
    flow = Flow.from_client_config(
        _oauth_client_config(request),
        scopes=_OAUTH_SCOPES,
        redirect_uri=redirect_uri,
    )

    # PKCE: pre-generate the verifier ourselves and inject it into the flow
    # BEFORE authorization_url() runs. Google's OAuth library will hash it,
    # send the challenge to Google, and Google will then require the same
    # verifier on the token exchange. Because /authorize and /callback are
    # separate request handlers (each with its own Flow instance), we must
    # carry the verifier across the redirect — we do that by signing it into
    # the state JWT. autogenerate_code_verifier=False stops the lib from
    # generating its own and clobbering ours.
    code_verifier = _generate_pkce_verifier()
    flow.code_verifier = code_verifier
    flow.autogenerate_code_verifier = False

    state = _make_oauth_state(effective_user.id, code_verifier=code_verifier)
    auth_url, _ = flow.authorization_url(
        access_type="offline",          # required to get a refresh_token
        prompt="consent",                # force refresh_token issuance even if user re-authorizes
        include_granted_scopes="true",
        state=state,
    )
    return RedirectResponse(url=auth_url)


@router.get("/oauth/callback")
async def oauth_callback(
    request: Request,
    code: Optional[str] = Query(None),
    state: Optional[str] = Query(None),
    error: Optional[str] = Query(None),
):
    """Google redirects here after the user clicks Allow/Deny on the consent
    screen. Exchanges the authorization code for tokens, stores the refresh
    token (encrypted), and bounces the user back to the IC Memo tab."""
    if error:
        return RedirectResponse(
            url=f"/?drive_oauth=error&reason={error}#memo",
            status_code=302,
        )
    if not code or not state:
        raise HTTPException(status_code=400, detail="Missing code or state in OAuth callback.")

    user_id, code_verifier = _verify_oauth_state(state)

    try:
        from google_auth_oauthlib.flow import Flow
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="google-auth-oauthlib not installed. Run: pip install google-auth-oauthlib",
        )

    redirect_uri = _oauth_redirect_uri(request)
    flow = Flow.from_client_config(
        _oauth_client_config(request),
        scopes=_OAUTH_SCOPES,
        redirect_uri=redirect_uri,
    )

    # Restore the PKCE verifier we stored in the state JWT during /authorize.
    # Without this, fetch_token() would have no verifier to send and Google
    # would respond "(invalid_grant) Missing code verifier".
    if code_verifier:
        flow.code_verifier = code_verifier
        flow.autogenerate_code_verifier = False

    try:
        # Pass the full callback URL so the library can extract the code itself
        flow.fetch_token(code=code)
    except Exception as e:
        # Most token-exchange failures come from a redirect_uri mismatch
        # between what we sent on /authorize and what we send on /callback,
        # or from an authorized-redirect-URI list in Google Cloud that
        # doesn't include the URL we constructed from forwarded headers.
        # Log everything we can, and surface a short error message so the
        # user can pass it back without having to scrape Replit logs.
        from urllib.parse import quote
        err_text = str(e)[:300]
        logger.error(
            "OAuth token exchange failed for user %s. "
            "redirect_uri sent=%s, error=%s",
            user_id, redirect_uri, err_text,
        )
        return RedirectResponse(
            url=(
                "/?drive_oauth=error"
                f"&reason=token_exchange_failed"
                f"&detail={quote(err_text)}"
                f"&redirect_uri_used={quote(redirect_uri)}"
                "#memo"
            ),
            status_code=302,
        )

    creds = flow.credentials
    if not creds.refresh_token:
        # Should not happen with prompt=consent + access_type=offline, but be defensive
        return RedirectResponse(
            url="/?drive_oauth=error&reason=no_refresh_token#memo",
            status_code=302,
        )

    # Verify the granted scopes include Drive read access
    granted_scopes = set(creds.scopes or [])
    if _REQUIRED_SCOPE not in granted_scopes:
        return RedirectResponse(
            url="/?drive_oauth=error&reason=missing_drive_scope#memo",
            status_code=302,
        )

    # Resolve the user's Google email from the userinfo endpoint
    google_email = ""
    try:
        from googleapiclient.discovery import build
        oauth2_service = build("oauth2", "v2", credentials=creds, cache_discovery=False)
        info = oauth2_service.userinfo().get().execute()
        google_email = info.get("email", "")
    except Exception as e:
        logger.warning(f"Could not resolve Google email for user {user_id}: {e}")

    _save_user_credentials(
        user_id=user_id,
        refresh_token=creds.refresh_token,
        google_email=google_email,
        scopes=list(creds.scopes or _OAUTH_SCOPES),
    )

    return RedirectResponse(url="/?drive_oauth=success#memo", status_code=302)


@router.post("/disconnect")
async def oauth_disconnect(user: CurrentUser = Depends(get_current_user)):
    """Revoke the user's Drive access by deleting their stored credentials.
    The user can also revoke from their Google account settings; this just
    drops our copy of the refresh token. (We don't proactively call Google's
    /revoke endpoint because the token is already encrypted at rest and
    deleting it is sufficient to stop us from using it.)"""
    deleted = _delete_user_credentials(user.id)
    return {"ok": True, "was_connected": deleted}


@router.get("/connection-status")
async def oauth_connection_status(user: CurrentUser = Depends(get_current_user)):
    """Report whether this user has connected Google Drive."""
    return _get_user_drive_status(user.id)


@router.get("/oauth/debug")
async def oauth_debug(request: Request, user: CurrentUser = Depends(get_current_user)):
    """Diagnostic — shows what redirect URI the app would send to Google,
    which secrets are configured, and which forwarded headers Replit is
    sending. Use this to verify config WITHOUT triggering an OAuth flow.

    Returns no secret values — just whether each secret is set and looks
    well-formed. Safe to share output with anyone debugging."""
    client_id = os.environ.get("GOOGLE_OAUTH_CLIENT_ID", "")
    client_secret = os.environ.get("GOOGLE_OAUTH_CLIENT_SECRET", "")
    enc_key = os.environ.get("GOOGLE_TOKEN_ENCRYPTION_KEY", "")
    explicit_redirect = os.environ.get("GOOGLE_OAUTH_REDIRECT_URI", "")

    # Compute the redirect URI the app would actually send
    try:
        computed_redirect = _oauth_redirect_uri(request)
    except Exception as e:
        computed_redirect = f"(could not compute: {e})"

    # Check Fernet key looks valid without revealing it
    fernet_ok = False
    fernet_error = ""
    if enc_key:
        try:
            from cryptography.fernet import Fernet
            Fernet(enc_key.encode() if isinstance(enc_key, str) else enc_key)
            fernet_ok = True
        except Exception as e:
            fernet_error = str(e)[:200]

    # Check Google libs are importable
    google_libs = {}
    for name in ("google.oauth2.credentials", "google_auth_oauthlib.flow",
                 "googleapiclient.discovery", "cryptography.fernet"):
        try:
            __import__(name)
            google_libs[name] = "ok"
        except ImportError as e:
            google_libs[name] = f"NOT INSTALLED: {e}"

    return {
        "redirect_uri_computed": computed_redirect,
        "redirect_uri_env_override": explicit_redirect or "(not set — using computed)",
        "redirect_uri_match_check": (
            "explicit env var matches computed" if explicit_redirect == computed_redirect
            else "explicit env var differs from computed (explicit wins)" if explicit_redirect
            else "(no explicit override; computed is used)"
        ),
        "client_id_set": bool(client_id),
        "client_id_length": len(client_id),
        "client_id_looks_valid": client_id.endswith(".apps.googleusercontent.com") if client_id else False,
        "client_secret_set": bool(client_secret),
        "client_secret_length": len(client_secret),
        "encryption_key_set": bool(enc_key),
        "encryption_key_valid": fernet_ok,
        "encryption_key_error": fernet_error,
        "scopes_requested": _OAUTH_SCOPES,
        "google_libraries": google_libs,
        "forwarded_headers": {
            "x-forwarded-proto": request.headers.get("x-forwarded-proto"),
            "x-forwarded-host": request.headers.get("x-forwarded-host"),
            "host": request.headers.get("host"),
            "url_scheme": request.url.scheme,
            "url_netloc": request.url.netloc,
        },
        "current_user_email": user.username,
        "instructions": (
            "Compare 'redirect_uri_computed' against the Authorized redirect URIs "
            "in your Google Cloud OAuth Client. They must match BYTE-FOR-BYTE — "
            "trailing slash, scheme, port, case all matter. If they differ, either "
            "(a) update the Google Cloud entry to match, or (b) set the "
            "GOOGLE_OAUTH_REDIRECT_URI Replit Secret to lock it down."
        ),
    }


# ═══════════════════════════════════════════════════════════════════════════════
#  LIBRARY CRUD
# ═══════════════════════════════════════════════════════════════════════════════

class LibraryCreate(BaseModel):
    company_name: str
    drive_folder_url: str


class LibraryUpdate(BaseModel):
    company_name: Optional[str] = None
    drive_folder_url: Optional[str] = None


@router.post("/libraries")
async def create_library(req: LibraryCreate, user: CurrentUser = Depends(get_current_user)):
    try:
        folder_id = _parse_drive_folder_id(req.drive_folder_url)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))

    conn = get_db()
    try:
        # Check if library already exists for this folder
        existing = conn.execute(
            "SELECT id FROM deal_document_libraries WHERE owner_id=? AND drive_folder_id=?",
            (user.id, folder_id),
        ).fetchone()
        if existing:
            return {"id": existing["id"], "already_exists": True}

        cur = conn.execute(
            """INSERT INTO deal_document_libraries (owner_id, company_name, drive_folder_id, drive_folder_url)
               VALUES (?, ?, ?, ?)""",
            (user.id, req.company_name, folder_id, req.drive_folder_url),
        )
        conn.commit()
        return {"id": cur.lastrowid, "folder_id": folder_id}
    finally:
        conn.close()


@router.get("/libraries")
async def list_libraries(user: CurrentUser = Depends(get_current_user)):
    conn = get_db()
    try:
        rows = conn.execute(
            """SELECT id, company_name, drive_folder_id, drive_folder_url, last_synced_at,
                      sync_status, doc_count, created_at
               FROM deal_document_libraries WHERE owner_id=? ORDER BY updated_at DESC""",
            (user.id,),
        ).fetchall()
        return [dict(r) for r in rows]
    finally:
        conn.close()


@router.get("/libraries/{lib_id}")
async def get_library(lib_id: int, user: CurrentUser = Depends(get_current_user)):
    conn = get_db()
    try:
        row = conn.execute(
            "SELECT * FROM deal_document_libraries WHERE id=? AND owner_id=?",
            (lib_id, user.id),
        ).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Library not found")
        lib = dict(row)

        # Include documents
        docs = conn.execute(
            """SELECT id, file_name, file_type, file_size, mime_type, subfolder_path,
                      doc_category, drive_modified, last_extracted,
                      LENGTH(extracted_text) as extracted_chars
               FROM deal_documents WHERE library_id=? AND owner_id=?
               ORDER BY subfolder_path, file_name""",
            (lib_id, user.id),
        ).fetchall()
        lib["documents"] = [dict(d) for d in docs]
        return lib
    finally:
        conn.close()


@router.delete("/libraries/{lib_id}")
async def delete_library(lib_id: int, user: CurrentUser = Depends(get_current_user)):
    conn = get_db()
    try:
        conn.execute("DELETE FROM deal_documents WHERE library_id=? AND owner_id=?", (lib_id, user.id))
        conn.execute("DELETE FROM deal_document_libraries WHERE id=? AND owner_id=?", (lib_id, user.id))
        conn.commit()
        return {"ok": True}
    finally:
        conn.close()


# ═══════════════════════════════════════════════════════════════════════════════
#  SYNC — pull files from Google Drive
# ═══════════════════════════════════════════════════════════════════════════════

@router.post("/libraries/{lib_id}/sync")
async def sync_library(lib_id: int, user: CurrentUser = Depends(get_current_user)):
    """
    Sync a library from Google Drive:
    1. List all files recursively in the Drive folder
    2. For new/changed files: download, extract text, store in DB
    3. Remove DB records for files deleted from Drive
    """
    conn = get_db()
    try:
        lib = conn.execute(
            "SELECT * FROM deal_document_libraries WHERE id=? AND owner_id=?",
            (lib_id, user.id),
        ).fetchone()
        if not lib:
            raise HTTPException(status_code=404, detail="Library not found")
        folder_id = lib["drive_folder_id"]
    finally:
        conn.close()

    # Update status to syncing
    conn = get_db()
    try:
        conn.execute(
            "UPDATE deal_document_libraries SET sync_status='syncing', updated_at=datetime('now') WHERE id=?",
            (lib_id,),
        )
        conn.commit()
    finally:
        conn.close()

    start_time = time.time()
    stats = {"new": 0, "updated": 0, "unchanged": 0, "removed": 0, "errors": 0, "skipped": 0}

    try:
        service = _get_drive_service(user.id)

        # List all files in the folder recursively
        drive_files = _list_files_recursive(service, folder_id)
        drive_file_ids = set()

        for f in drive_files:
            drive_file_ids.add(f["id"])
            mime = f.get("mimeType", "")
            name = f.get("name", "unknown")
            modified = f.get("modifiedTime", "")
            size = int(f.get("size", 0)) if f.get("size") else 0
            subfolder = f.get("subfolder_path", "")

            # Check if this file type is extractable
            is_extractable = mime in _EXTRACTABLE_MIMES or mime in _GOOGLE_EXPORT_MIMES
            if not is_extractable:
                # Check by extension
                ext = Path(name).suffix.lower()
                if ext not in ('.pdf', '.docx', '.xlsx', '.pptx', '.txt', '.csv', '.md', '.json', '.html'):
                    stats["skipped"] += 1
                    continue

            # Check if we already have this file
            conn = get_db()
            try:
                existing = conn.execute(
                    "SELECT id, drive_modified, extraction_hash FROM deal_documents WHERE library_id=? AND drive_file_id=?",
                    (lib_id, f["id"]),
                ).fetchone()
            finally:
                conn.close()

            if existing and existing["drive_modified"] == modified:
                # File unchanged
                stats["unchanged"] += 1
                continue

            # Download and extract
            try:
                file_bytes = _download_file(service, f)
                text = _extract_text(file_bytes, name, mime)
                text_hash = hashlib.md5(text.encode()).hexdigest()

                ext = Path(name).suffix.lower()
                if mime in _GOOGLE_EXPORT_MIMES:
                    _, ext = _GOOGLE_EXPORT_MIMES[mime]

                category = _infer_category(name)

                conn = get_db()
                try:
                    if existing:
                        # Update existing record
                        conn.execute(
                            """UPDATE deal_documents SET
                                file_name=?, file_type=?, file_size=?, mime_type=?,
                                subfolder_path=?, doc_category=?, extracted_text=?,
                                extraction_hash=?, drive_modified=?,
                                last_extracted=datetime('now'), updated_at=datetime('now')
                               WHERE id=?""",
                            (name, ext, size, mime, subfolder, category, text,
                             text_hash, modified, existing["id"]),
                        )
                        stats["updated"] += 1
                    else:
                        # Insert new record
                        conn.execute(
                            """INSERT INTO deal_documents
                               (library_id, owner_id, drive_file_id, file_name, file_type,
                                file_size, mime_type, subfolder_path, doc_category,
                                extracted_text, extraction_hash, drive_modified, last_extracted)
                               VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, datetime('now'))""",
                            (lib_id, user.id, f["id"], name, ext, size, mime,
                             subfolder, category, text, text_hash, modified),
                        )
                        stats["new"] += 1
                    conn.commit()
                finally:
                    conn.close()

            except Exception as e:
                logger.error(f"Drive sync: failed to process {name}: {e}")
                stats["errors"] += 1
                continue

        # Remove documents that are no longer in Drive
        conn = get_db()
        try:
            existing_docs = conn.execute(
                "SELECT id, drive_file_id FROM deal_documents WHERE library_id=?",
                (lib_id,),
            ).fetchall()
            for doc in existing_docs:
                if doc["drive_file_id"] not in drive_file_ids:
                    conn.execute("DELETE FROM deal_documents WHERE id=?", (doc["id"],))
                    stats["removed"] += 1
            conn.commit()
        finally:
            conn.close()

        # Update library metadata
        elapsed = round(time.time() - start_time, 2)
        conn = get_db()
        try:
            total_docs = conn.execute(
                "SELECT COUNT(*) as c FROM deal_documents WHERE library_id=?", (lib_id,),
            ).fetchone()["c"]
            conn.execute(
                """UPDATE deal_document_libraries SET
                    sync_status='synced', last_synced_at=datetime('now'),
                    doc_count=?, updated_at=datetime('now')
                   WHERE id=?""",
                (total_docs, lib_id),
            )
            conn.commit()
        finally:
            conn.close()

        return {
            "status": "synced",
            "elapsed_s": elapsed,
            "total_docs": total_docs,
            "stats": stats,
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Drive sync failed for library {lib_id}: {e}")
        conn = get_db()
        try:
            conn.execute(
                "UPDATE deal_document_libraries SET sync_status='error', updated_at=datetime('now') WHERE id=?",
                (lib_id,),
            )
            conn.commit()
        finally:
            conn.close()
        # Use a longer truncation so Google's "Drive API not enabled" error
        # — which includes the activation URL — is visible to the user.
        raise HTTPException(status_code=500, detail=f"Sync failed: {str(e)[:1200]}")


# ═══════════════════════════════════════════════════════════════════════════════
#  DOCUMENT MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════════════

@router.get("/libraries/{lib_id}/documents")
async def list_library_documents(lib_id: int, user: CurrentUser = Depends(get_current_user)):
    conn = get_db()
    try:
        docs = conn.execute(
            """SELECT id, file_name, file_type, file_size, mime_type, subfolder_path,
                      doc_category, drive_modified, last_extracted,
                      LENGTH(extracted_text) as extracted_chars
               FROM deal_documents WHERE library_id=? AND owner_id=?
               ORDER BY subfolder_path, file_name""",
            (lib_id, user.id),
        ).fetchall()
        return [dict(d) for d in docs]
    finally:
        conn.close()


class CategoryUpdate(BaseModel):
    category: str


@router.put("/documents/{doc_id}/category")
async def update_document_category(doc_id: int, req: CategoryUpdate, user: CurrentUser = Depends(get_current_user)):
    conn = get_db()
    try:
        conn.execute(
            "UPDATE deal_documents SET doc_category=?, updated_at=datetime('now') WHERE id=? AND owner_id=?",
            (req.category, doc_id, user.id),
        )
        conn.commit()
        return {"ok": True}
    finally:
        conn.close()


# ═══════════════════════════════════════════════════════════════════════════════
#  HELPER — Load documents from library for memo generation
# ═══════════════════════════════════════════════════════════════════════════════

def load_library_documents(library_id: int, owner_id: int) -> list:
    """Load all extracted documents from a deal library for use in memo generation.
    Returns list of dicts compatible with the memo pipeline's expected format."""
    conn = get_db()
    try:
        rows = conn.execute(
            """SELECT id, file_name, doc_category, extracted_text, subfolder_path
               FROM deal_documents WHERE library_id=? AND owner_id=?
               ORDER BY subfolder_path, doc_category, file_name""",
            (library_id, owner_id),
        ).fetchall()
        return [dict(r) for r in rows if (r["extracted_text"] or "").strip()]
    finally:
        conn.close()
