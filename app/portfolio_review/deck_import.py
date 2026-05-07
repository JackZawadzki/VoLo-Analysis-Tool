"""
Import a Monthly All-Team PortCo Updates deck into pr_traction_snapshots.

The deck format (one row per company on the "Where We Are Today" slides):
    [#] | [Company Name] | [Deal Lead] | [Status update narrative]

We parse those rows out of the slide tables, then send each narrative through
Claude to extract structured metrics (status, revenue, growth, milestones, etc.)
and persist a pr_traction_snapshot row.

Run:
    python -m app.portfolio_review.deck_import path/to/deck.pptx
"""

from __future__ import annotations

import json
import logging
import os
import re
import sqlite3
import sys
from datetime import datetime
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


# ── Parse the deck ────────────────────────────────────────────────────────────
def _walk_shape(shape, rows: list[list[str]]):
    """Recursive walk: collect table rows from any nested group/table."""
    if shape.has_table:
        for row in shape.table.rows:
            cells = [(cell.text_frame.text if cell.text_frame else "").strip() for cell in row.cells]
            if any(cells):
                rows.append(cells)
        return
    if shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            _walk_shape(sub, rows)


def parse_company_rows(pptx_path: Path) -> list[dict]:
    """Extract one record per company from the 'Where We Are Today' tables.

    The deck has a Fund I section (typically rows 1-26) and a Fund II section
    (rows 1-9). We don't try to distinguish funds here — we just match each
    name back to the pr_companies table which already has fund tagged.
    """
    from pptx import Presentation

    prs = Presentation(pptx_path)
    rows: list[list[str]] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            _walk_shape(shape, rows)

    companies: list[dict] = []
    seen_names = set()
    for r in rows:
        if len(r) < 4:
            continue
        idx_cell, name_cell, lead_cell, narrative_cell = r[0], r[1], r[2], r[3]
        if not re.match(r"^\d+$", idx_cell.strip()):
            continue
        name = name_cell.strip()
        if not name:
            continue
        narrative = narrative_cell.strip()
        if not narrative or narrative.lower().startswith("see fund"):
            continue

        key = name.lower()
        if key in seen_names:
            continue
        seen_names.add(key)
        companies.append({
            "name": name,
            "lead": lead_cell.strip(),
            "narrative": narrative,
            "deck_row_index": int(idx_cell.strip()),
        })
    return companies


def parse_fundraising_list(pptx_path: Path) -> dict[str, str]:
    """Parse the 'PortCos Actively Fundraising Now' slide.

    Returns {normalized_company_name: fundraising_description}.
    """
    from pptx import Presentation

    prs = Presentation(pptx_path)
    for slide in prs.slides:
        text_blocks = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_blocks.append(shape.text)
        all_text = "\n".join(text_blocks)
        if "Actively Fundraising" not in all_text:
            continue
        # Lines like "Blue Frontier - Nearing term sheet on $20M+ CN Bridge..."
        out: dict[str, str] = {}
        for line in all_text.splitlines():
            line = line.strip()
            if not line or "Actively Fundraising" in line or "Confidential" in line:
                continue
            m = re.match(r"^([^-]+?)\s*[-–—]\s*(.+)$", line)
            if not m:
                continue
            name, desc = m.group(1).strip(), m.group(2).strip()
            if name and desc:
                out[_normalize_name(name)] = desc
        return out
    return {}


# ── Match parsed names to pr_companies rows ───────────────────────────────────
def _normalize_name(name: str) -> str:
    s = (name or "").lower()
    s = re.sub(r"\b(inc|incorporated|corp|corporation|ltd|llc|co|company|public benefit corporation|pbc)\.?\b", "", s)
    s = re.sub(r"[^a-z0-9 ]", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def match_to_company(conn: sqlite3.Connection, deck_name: str) -> Optional[int]:
    """Return pr_companies.id for the given deck name, or None."""
    norm = _normalize_name(deck_name)
    rows = conn.execute("SELECT id, name FROM pr_companies").fetchall()
    by_norm = {_normalize_name(r["name"]): r["id"] for r in rows}
    if norm in by_norm:
        return by_norm[norm]
    # Partial match — deck name is "AICrete" but DB has "AICrete Corp."
    for cn, cid in by_norm.items():
        if not cn or not norm:
            continue
        if cn.startswith(norm) or norm.startswith(cn) or norm in cn or cn in norm:
            if min(len(cn), len(norm)) >= 3:
                return cid
    return None


# ── Anthropic extraction ──────────────────────────────────────────────────────
EXTRACTION_PROMPT = """Extract structured portfolio metrics from this status update.

Return ONLY a JSON object with this exact shape, no markdown or commentary:

{{
  "commercial_status": "Pre-Rev" | "Pilot" | "Commercial" | "Hyperscale",
  "revenue_current": <USD number or null>,
  "revenue_period": "<e.g. '2025' or 'Q4 2025'>",
  "arr_current": <USD number or null>,
  "fundraising_status": "<short label like 'Closed Series A', 'Raising Series B', 'Bridge round closing', or 'Not raising'>",
  "key_milestone": "<single most important recent milestone, 1 sentence>",
  "summary": "<concise 2-sentence summary of where this company is now>",
  "confidence": "low" | "medium" | "high"
}}

DEFINITIONS:
- "Pre-Rev"     — no paying customers
- "Pilot"       — pilots/LOIs only, no recurring revenue
- "Commercial"  — paying customers, recurring revenue under ~$5M ARR
- "Hyperscale"  — $5M+ ARR or major commercial traction with multiple Tier 1 customers

Numbers: if the update says "$10M+ revenue", set revenue_current=10000000. If it says "multi-million ARR", estimate arr_current=2000000 with confidence="low".

Company: {company_name}
Lead: {lead}
Status update: {narrative}
"""


def call_claude(company_name: str, lead: str, narrative: str) -> dict:
    from anthropic import Anthropic
    api_key = os.environ.get("ANTHROPIC_API_KEY", "").strip()
    if not api_key:
        raise RuntimeError("ANTHROPIC_API_KEY not set in env")
    client = Anthropic(api_key=api_key)
    prompt = EXTRACTION_PROMPT.format(
        company_name=company_name, lead=lead or "—", narrative=narrative,
    )
    msg = client.messages.create(
        model="claude-haiku-4-5-20251001",  # haiku is plenty for this — cheap + fast
        max_tokens=600,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = "".join(b.text for b in msg.content if hasattr(b, "text"))
    cleaned = raw.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
        cleaned = re.sub(r"\s*```\s*$", "", cleaned)
    try:
        data = json.loads(cleaned)
    except json.JSONDecodeError:
        m = re.search(r"\{.*\}", cleaned, re.DOTALL)
        if not m:
            raise ValueError(f"LLM did not return JSON: {cleaned[:200]}")
        data = json.loads(m.group())
    return {"data": data, "raw": raw, "model": "claude-haiku-4-5"}


# ── Main entry point ─────────────────────────────────────────────────────────
def import_deck(deck_path: Path, conn: sqlite3.Connection,
                source_name: Optional[str] = None) -> dict:
    """Parse the deck, extract per-company metrics with Claude, persist
    pr_traction_snapshots rows. Returns counts.

    Also parses the 'Actively Fundraising' slide and merges fundraising_status
    into each matching company's snapshot.
    """
    if not deck_path.exists():
        raise FileNotFoundError(deck_path)

    parsed = parse_company_rows(deck_path)
    fundraising_map = parse_fundraising_list(deck_path)
    print(f"Parsed {len(parsed)} company rows + {len(fundraising_map)} active fundraises")

    counts = {"matched": 0, "unmatched": 0, "snapshots_added": 0, "errors": 0}
    unmatched_names: list[str] = []

    for entry in parsed:
        name = entry["name"]
        cid = match_to_company(conn, name)
        if not cid:
            counts["unmatched"] += 1
            unmatched_names.append(name)
            continue
        counts["matched"] += 1

        # Look up fundraising status by normalized name
        fr_status = fundraising_map.get(_normalize_name(name), "")

        try:
            extraction = call_claude(name, entry["lead"], entry["narrative"])
            d = extraction["data"]
        except Exception as e:
            print(f"  [{name}] LLM failed: {e}")
            counts["errors"] += 1
            continue

        # If LLM produced its own fundraising_status and we don't have one
        # from the active list, prefer the LLM's narrative version.
        final_fr = fr_status or d.get("fundraising_status", "")

        source_files = [{
            "name": source_name or deck_path.name,
            "folder_type": "team_update",
            "lead": entry["lead"],
        }]

        conn.execute(
            """INSERT INTO pr_traction_snapshots
            (company_id, commercial_status, revenue_current, revenue_period,
             arr_current, notable_milestones, summary, source_files,
             model_used, confidence, raw_response,
             deal_lead, narrative_raw, fundraising_status, deck_row_index)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
            (
                cid,
                d.get("commercial_status", ""),
                d.get("revenue_current"),
                d.get("revenue_period", ""),
                d.get("arr_current"),
                d.get("key_milestone", ""),
                d.get("summary", ""),
                json.dumps(source_files),
                extraction["model"],
                d.get("confidence", "medium"),
                extraction["raw"],
                entry["lead"],
                entry["narrative"],
                final_fr,
                entry.get("deck_row_index"),
            ),
        )
        counts["snapshots_added"] += 1
        status = d.get("commercial_status", "?")
        rev = d.get("revenue_current")
        rev_str = f"${rev/1e6:.1f}M" if rev else "—"
        fr_marker = " 💰" if final_fr else ""
        print(f"  ✓ {name:<25}  {status:<12}  rev={rev_str:<8}  ({d.get('confidence','?')} conf){fr_marker}")

    conn.commit()
    counts["unmatched_names"] = unmatched_names
    counts["fundraising_count"] = sum(1 for v in fundraising_map.values() if v)
    return counts


if __name__ == "__main__":
    import argparse
    sys.path.insert(0, str(Path(__file__).resolve().parents[2]))
    from app.database import get_db, init_db
    from app.portfolio_review.schema import apply_schema

    parser = argparse.ArgumentParser(description="Import a portfolio team-update deck.")
    parser.add_argument("deck", help="Path to the .pptx file")
    args = parser.parse_args()

    init_db()
    with get_db() as conn:
        apply_schema(conn)
        result = import_deck(Path(args.deck), conn, source_name=Path(args.deck).name)
    print()
    print(json.dumps(result, indent=2, default=str))
