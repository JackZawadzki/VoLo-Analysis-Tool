"""
Traction & Status scanner — pulls portfolio company materials from Drive
and uses Claude to summarize commercial status and revenue growth.

Pipeline per company:
  1. Look up linked Drive folders (current + diligence)
  2. List files, pick the most relevant N (latest decks, financial models,
     board updates, investor letters)
  3. Download + extract text from each
  4. Send the extracted text to Claude with a structured prompt
  5. Parse the JSON response and persist a pr_traction_snapshot row

Folder discovery:
  Given a parent folder ID containing one subfolder per company, match
  subfolders to pr_companies.name (exact + fuzzy fallback).
"""

from __future__ import annotations

import io
import json
import logging
import os
import re
import sqlite3
from datetime import datetime
from pathlib import Path
from typing import Any, Optional

logger = logging.getLogger(__name__)

# How many files of each "kind" to send to the LLM per company per scan.
# More = richer context but more tokens. These limits keep cost predictable.
MAX_DECKS = 1
MAX_MODELS = 1
MAX_BOARD_UPDATES = 2
MAX_INVESTOR_UPDATES = 2

# Max characters of text per file to include in the prompt.
# Decks are usually short, models can be long — capping keeps latency sane.
MAX_CHARS_PER_FILE = 25_000

MIME_FOLDER = "application/vnd.google-apps.folder"
MIME_GSHEET = "application/vnd.google-apps.spreadsheet"
MIME_GDOC = "application/vnd.google-apps.document"
MIME_GSLIDES = "application/vnd.google-apps.presentation"

# File-type heuristics for picking the right docs to scan
DECK_PATTERNS = re.compile(r"deck|pitch|investor.?presentation", re.I)
MODEL_PATTERNS = re.compile(r"financial.?model|projection|forecast|p&l|operating", re.I)
BOARD_PATTERNS = re.compile(r"board|quarterly|board.?update|board.?pack", re.I)
INVESTOR_PATTERNS = re.compile(r"investor.?update|investor.?letter|monthly.?update|founder.?update", re.I)


# ── Folder discovery ─────────────────────────────────────────────────────────
def list_subfolders(service, parent_folder_id: str) -> list[dict]:
    """List immediate subfolders of a Drive folder."""
    q = f"'{parent_folder_id}' in parents and mimeType='{MIME_FOLDER}' and trashed=false"
    files = []
    page_token = None
    while True:
        resp = service.files().list(
            q=q, pageSize=200,
            fields="nextPageToken, files(id, name, modifiedTime)",
            supportsAllDrives=True, includeItemsFromAllDrives=True,
            pageToken=page_token,
        ).execute()
        files.extend(resp.get("files", []))
        page_token = resp.get("nextPageToken")
        if not page_token:
            break
    return files


def _normalize_name(name: str) -> str:
    """Lowercase, strip Inc/Corp/Ltd/etc., collapse whitespace and punctuation."""
    s = (name or "").lower()
    s = re.sub(r"\b(inc|incorporated|corp|corporation|ltd|llc|co|company|public benefit corporation|pbc)\.?\b", "", s)
    s = re.sub(r"[^a-z0-9 ]", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def match_folders_to_companies(
    conn: sqlite3.Connection, subfolders: list[dict],
    parent_folder_id: str, folder_type: str = "current",
) -> dict:
    """Match each subfolder to a pr_company by normalized-name. Persist the
    matches into pr_company_folders.

    Returns counts: {matched, unmatched, skipped_dupes}.
    """
    companies = conn.execute("SELECT id, name FROM pr_companies").fetchall()
    name_to_id = {_normalize_name(row["name"]): row["id"] for row in companies}

    matched, unmatched, skipped = 0, [], 0

    for f in subfolders:
        folder_name = f["name"]
        norm = _normalize_name(folder_name)
        cid = name_to_id.get(norm)
        if cid is None:
            # Fuzzy: try partial match (folder name contains company name or vice versa)
            for cn, cid_candidate in name_to_id.items():
                if cn and (cn in norm or norm in cn) and len(cn) >= 4:
                    cid = cid_candidate
                    break
        if cid is None:
            unmatched.append(folder_name)
            continue

        # Skip if this (company_id, folder_type) is already linked to a different folder
        existing = conn.execute(
            "SELECT drive_folder_id FROM pr_company_folders WHERE company_id=? AND folder_type=?",
            (cid, folder_type),
        ).fetchone()
        if existing and existing[0] != f["id"]:
            skipped += 1
            continue
        if existing and existing[0] == f["id"]:
            continue  # already linked to this folder

        conn.execute(
            """INSERT INTO pr_company_folders
               (company_id, folder_type, drive_folder_id, drive_folder_name, parent_folder_id, match_confidence)
               VALUES (?, ?, ?, ?, ?, ?)""",
            (cid, folder_type, f["id"], folder_name, parent_folder_id, "auto"),
        )
        matched += 1

    conn.commit()
    return {"matched": matched, "unmatched": unmatched, "skipped_existing": skipped}


# ── File picking + text extraction ────────────────────────────────────────────
def _list_files_recursive(service, folder_id: str, depth: int = 0, max_depth: int = 2) -> list[dict]:
    """List all non-folder files inside `folder_id`, recursing up to max_depth."""
    if depth > max_depth:
        return []
    out = []
    page_token = None
    q = f"'{folder_id}' in parents and trashed=false"
    while True:
        resp = service.files().list(
            q=q, pageSize=200,
            fields="nextPageToken, files(id, name, mimeType, modifiedTime, size)",
            supportsAllDrives=True, includeItemsFromAllDrives=True,
            pageToken=page_token,
        ).execute()
        for f in resp.get("files", []):
            if f["mimeType"] == MIME_FOLDER:
                out.extend(_list_files_recursive(service, f["id"], depth + 1, max_depth))
            else:
                out.append(f)
        page_token = resp.get("nextPageToken")
        if not page_token:
            break
    return out


def _classify_file(f: dict) -> Optional[str]:
    """Return 'deck' | 'model' | 'board' | 'investor' | None."""
    name = f.get("name", "")
    mime = f.get("mimeType", "")

    if BOARD_PATTERNS.search(name):
        return "board"
    if INVESTOR_PATTERNS.search(name):
        return "investor"
    if DECK_PATTERNS.search(name):
        return "deck"
    if MODEL_PATTERNS.search(name):
        return "model"

    # Fallback by mime type
    if mime == MIME_GSLIDES or mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return "deck"
    if mime == MIME_GSHEET or mime in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ):
        return "model"
    if mime == "application/pdf":
        return "investor"  # PDFs we can't classify by name → assume investor update
    return None


def _pick_files(files: list[dict]) -> list[dict]:
    """Pick the top-N files of each kind, ordered by modified date desc."""
    classified = [(f, _classify_file(f)) for f in files]
    classified = [(f, k) for f, k in classified if k is not None]

    by_kind: dict[str, list[dict]] = {"deck": [], "model": [], "board": [], "investor": []}
    for f, k in classified:
        by_kind[k].append(f)

    # Sort each list by modifiedTime desc and trim
    def sorted_by_mtime(lst):
        return sorted(lst, key=lambda x: x.get("modifiedTime", ""), reverse=True)

    picked = []
    picked.extend(sorted_by_mtime(by_kind["deck"])[:MAX_DECKS])
    picked.extend(sorted_by_mtime(by_kind["model"])[:MAX_MODELS])
    picked.extend(sorted_by_mtime(by_kind["board"])[:MAX_BOARD_UPDATES])
    picked.extend(sorted_by_mtime(by_kind["investor"])[:MAX_INVESTOR_UPDATES])
    return picked


def _download_and_extract_text(service, file_meta: dict) -> str:
    """Download a Drive file and extract plain text. Returns '' if extraction fails."""
    from googleapiclient.http import MediaIoBaseDownload
    mime = file_meta.get("mimeType", "")
    name = file_meta.get("name", "")

    # Google-native types: export to a text-friendly format
    export_mime = None
    if mime == MIME_GDOC:
        export_mime = "text/plain"
    elif mime == MIME_GSLIDES:
        export_mime = "text/plain"
    elif mime == MIME_GSHEET:
        export_mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

    try:
        if export_mime:
            request = service.files().export_media(fileId=file_meta["id"], mimeType=export_mime)
        else:
            request = service.files().get_media(fileId=file_meta["id"], supportsAllDrives=True)

        buf = io.BytesIO()
        downloader = MediaIoBaseDownload(buf, request)
        done = False
        while not done:
            _, done = downloader.next_chunk()
        raw = buf.getvalue()
    except Exception as e:
        logger.warning(f"Download failed for {name}: {e}")
        return ""

    # Decode by file type
    try:
        if mime == MIME_GDOC or mime == MIME_GSLIDES or mime.startswith("text/"):
            return raw.decode("utf-8", errors="ignore")[:MAX_CHARS_PER_FILE]
        if mime == "application/pdf":
            return _extract_pdf_text(raw)[:MAX_CHARS_PER_FILE]
        if mime in ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    MIME_GSHEET):
            return _extract_xlsx_text(raw)[:MAX_CHARS_PER_FILE]
        if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return _extract_pptx_text(raw)[:MAX_CHARS_PER_FILE]
        if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return _extract_docx_text(raw)[:MAX_CHARS_PER_FILE]
    except Exception as e:
        logger.warning(f"Text extract failed for {name} ({mime}): {e}")
    return ""


def _extract_pdf_text(raw: bytes) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(raw))
        return "\n".join(p.extract_text() or "" for p in reader.pages)
    except Exception:
        return ""


def _extract_xlsx_text(raw: bytes) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(raw), data_only=True)
        out: list[str] = []
        for sn in wb.sheetnames[:5]:  # first 5 sheets only
            ws = wb[sn]
            out.append(f"=== Sheet: {sn} ===")
            for row in ws.iter_rows(max_row=80, values_only=True):
                line = "\t".join("" if c is None else str(c) for c in row)
                if line.strip():
                    out.append(line)
        return "\n".join(out)
    except Exception:
        return ""


def _extract_pptx_text(raw: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        out: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    out.append(shape.text)
        return "\n".join(out)
    except Exception:
        return ""


def _extract_docx_text(raw: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(raw))
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception:
        return ""


# ── LLM extraction ────────────────────────────────────────────────────────────
EXTRACTION_PROMPT = """You are extracting structured portfolio metrics for a venture capital firm's quarterly
portfolio review. The firm wants commercial status and revenue growth for one company.

Return ONLY a JSON object — no markdown, no commentary — with this exact shape:

{{
  "commercial_status": "Pre-Rev" | "Pilot" | "Commercial" | "Hyperscale",
  "revenue_current": <number in USD or null>,
  "revenue_prior": <number in USD or null>,
  "revenue_period": "<e.g. 'FY2024 vs FY2023'>",
  "revenue_growth_pct": <decimal — 0.40 = 40% — or null>,
  "arr_current": <number in USD or null>,
  "customer_count": <integer or null>,
  "runway_months": <number or null>,
  "notable_milestones": "<1-3 sentences on key recent milestones>",
  "summary": "<2-4 sentences on overall traction and status>",
  "baseline_status": "<commercial status at investment time, from DD docs>",
  "baseline_revenue": <revenue at investment time or null>,
  "baseline_summary": "<1-2 sentences describing state at investment time>",
  "change_vs_baseline": "<1-2 sentences on what's changed since investment>",
  "confidence": "low" | "medium" | "high"
}}

DEFINITIONS:
- "Pre-Rev"      — no paying customers / pilot programs only
- "Pilot"        — early pilots or LOIs, not yet at scale
- "Commercial"   — multiple paying customers, recurring revenue under $5M ARR
- "Hyperscale"   — major commercial traction, $5M+ ARR or multiple Tier 1 customers

If a field is not directly stated, infer conservatively. Mark confidence="low" if you're guessing.

Company: {company_name}

The materials below are tagged by source folder ([CURRENT] = ongoing, [DILIGENCE] = investment-time baseline):

{file_blocks}
"""


def _call_anthropic(company_name: str, file_blocks: str) -> dict:
    """Call Claude to extract structured metrics. Returns the parsed JSON dict
    plus the raw response text for audit."""
    try:
        from anthropic import Anthropic
    except ImportError:
        raise RuntimeError("anthropic package not installed")

    api_key = os.environ.get("ANTHROPIC_API_KEY", "").strip()
    if not api_key:
        raise RuntimeError("ANTHROPIC_API_KEY not set in env")

    client = Anthropic(api_key=api_key)
    prompt = EXTRACTION_PROMPT.format(company_name=company_name, file_blocks=file_blocks)

    msg = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=2000,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = "".join(b.text for b in msg.content if hasattr(b, "text"))

    # Strip code fences if Claude added them despite instructions
    cleaned = raw.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
        cleaned = re.sub(r"\s*```\s*$", "", cleaned)

    try:
        data = json.loads(cleaned)
    except json.JSONDecodeError:
        # Try to extract the first JSON object from the text
        m = re.search(r"\{.*\}", cleaned, re.DOTALL)
        if m:
            data = json.loads(m.group())
        else:
            raise ValueError(f"LLM response was not valid JSON: {cleaned[:300]}")

    return {"data": data, "raw": raw, "model": "claude-opus-4-6"}


# ── Public entry point ───────────────────────────────────────────────────────
def scan_company(conn: sqlite3.Connection, company_id: int, user_id: int) -> dict:
    """Scan one company's folders and write a fresh pr_traction_snapshot row.

    Returns a status dict suitable for surfacing in the UI.
    """
    from ..routes.drive import _get_drive_service

    company = conn.execute(
        "SELECT id, name FROM pr_companies WHERE id=?", (company_id,)
    ).fetchone()
    if not company:
        raise ValueError(f"company {company_id} not found")

    folders = conn.execute(
        "SELECT * FROM pr_company_folders WHERE company_id=?", (company_id,)
    ).fetchall()
    if not folders:
        raise ValueError(
            f"No Drive folders linked for {company['name']}. "
            "Run the folder-discovery scan first."
        )

    service = _get_drive_service(user_id)

    # Collect picked files across all folder types, tagging each by folder type
    all_picked: list[tuple[dict, str]] = []  # (file_meta, folder_type)
    for f in folders:
        try:
            files = _list_files_recursive(service, f["drive_folder_id"])
            picked = _pick_files(files)
            for p in picked:
                all_picked.append((p, f["folder_type"]))
        except Exception as e:
            logger.warning(f"Listing failed for folder {f['drive_folder_name']}: {e}")
            continue

    if not all_picked:
        raise ValueError(f"No relevant files found in linked folders for {company['name']}")

    # Download + extract text from each
    file_blocks: list[str] = []
    source_files: list[dict] = []
    for file_meta, folder_type in all_picked:
        text = _download_and_extract_text(service, file_meta)
        if not text:
            continue
        tag = "[DILIGENCE]" if folder_type == "diligence" else "[CURRENT]"
        file_blocks.append(
            f"--- {tag} {file_meta['name']} (modified {file_meta.get('modifiedTime','?')}) ---\n{text}"
        )
        source_files.append({"name": file_meta["name"], "folder_type": folder_type,
                             "modified": file_meta.get("modifiedTime")})

    if not file_blocks:
        raise ValueError(f"Could not extract text from any file for {company['name']}")

    # Call Claude
    full_block = "\n\n".join(file_blocks)
    extraction = _call_anthropic(company["name"], full_block)
    d = extraction["data"]

    # Persist snapshot
    cur = conn.execute(
        """INSERT INTO pr_traction_snapshots
        (company_id, commercial_status, revenue_current, revenue_prior, revenue_period,
         revenue_growth_pct, arr_current, customer_count, runway_months,
         notable_milestones, summary, baseline_status, baseline_revenue, baseline_summary,
         change_vs_baseline, source_files, model_used, confidence, raw_response)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
        (
            company_id,
            d.get("commercial_status", ""),
            d.get("revenue_current"),
            d.get("revenue_prior"),
            d.get("revenue_period", ""),
            d.get("revenue_growth_pct"),
            d.get("arr_current"),
            d.get("customer_count"),
            d.get("runway_months"),
            d.get("notable_milestones", ""),
            d.get("summary", ""),
            d.get("baseline_status", ""),
            d.get("baseline_revenue"),
            d.get("baseline_summary", ""),
            d.get("change_vs_baseline", ""),
            json.dumps(source_files),
            extraction["model"],
            d.get("confidence", "medium"),
            extraction["raw"],
        ),
    )
    snapshot_id = cur.lastrowid

    # Update last_scanned_at on each folder
    for f in folders:
        conn.execute(
            "UPDATE pr_company_folders SET last_scanned_at = datetime('now') WHERE id = ?",
            (f["id"],),
        )
    conn.commit()

    return {
        "snapshot_id": snapshot_id,
        "company_id": company_id,
        "company_name": company["name"],
        "files_scanned": len(source_files),
        "model": extraction["model"],
        **{k: d.get(k) for k in (
            "commercial_status", "revenue_current", "revenue_growth_pct",
            "summary", "change_vs_baseline", "confidence",
        )},
    }
