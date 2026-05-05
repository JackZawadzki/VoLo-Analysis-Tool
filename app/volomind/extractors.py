"""File body extractors with optional dependencies.

Each extractor takes raw bytes and returns plain text or None. Optional
libraries (pypdf, python-docx, python-pptx) are imported lazily — their
absence surfaces as a one-time warning rather than crashing the sync.

Image-only PDFs fall back to Anthropic vision when ANTHROPIC_API_KEY is set
and pypdf returns sparse text.
"""

from __future__ import annotations

import io
import os
from typing import Callable, Optional


_warned: set[str] = set()


def _warn_once(key: str, message: str) -> None:
    if key in _warned:
        return
    _warned.add(key)
    print(f"[volomind/extractors] {message}", flush=True)


# --- PDF -------------------------------------------------------------------

_IMAGE_PDF_DENSITY_THRESHOLD = 100  # chars / page
_VISION_MAX_BYTES = 32 * 1024 * 1024
_VISION_MAX_PAGES = 100


def extract_pdf(data: bytes) -> Optional[str]:
    try:
        from pypdf import PdfReader
    except ImportError:
        _warn_once("pypdf", "pypdf not installed — PDF files will be skipped.")
        return None
    try:
        reader = PdfReader(io.BytesIO(data))
        page_count = len(reader.pages)
        parts = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                continue
        text = "\n\n".join(p for p in parts if p.strip())
    except Exception:
        return _extract_pdf_vision(data)

    density = (len(text) / page_count) if page_count else 0
    if density < _IMAGE_PDF_DENSITY_THRESHOLD:
        if page_count and page_count <= _VISION_MAX_PAGES and len(data) <= _VISION_MAX_BYTES:
            vision_text = _extract_pdf_vision(data)
            if vision_text:
                return vision_text
    return text or None


_VISION_PROMPT = (
    "Extract all text from this PDF, including text inside images, charts, "
    "diagrams, and slides. Preserve structure with markdown. For figures with "
    "no extractable text, give a one-line description in [brackets]. Return "
    "only extracted content, no preamble."
)


def _extract_pdf_vision(data: bytes) -> Optional[str]:
    api_key = os.environ.get("ANTHROPIC_API_KEY", "").strip()
    if not api_key:
        return None
    try:
        from anthropic import Anthropic
    except ImportError:
        _warn_once("anthropic", "anthropic SDK not installed — image-only PDF fallback disabled.")
        return None
    import base64
    pdf_b64 = base64.standard_b64encode(data).decode("utf-8")
    model = os.environ.get("ANTHROPIC_VISION_PDF_MODEL", "claude-haiku-4-5-20251001")
    try:
        client = Anthropic(api_key=api_key)
        msg = client.messages.create(
            model=model,
            max_tokens=8192,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "document", "source": {
                        "type": "base64",
                        "media_type": "application/pdf",
                        "data": pdf_b64,
                    }},
                    {"type": "text", "text": _VISION_PROMPT},
                ],
            }],
        )
        chunks = [b.text for b in msg.content if hasattr(b, "text")]
        text = "\n".join(c for c in chunks if c).strip()
        return text or None
    except Exception:
        return None


# --- DOCX ------------------------------------------------------------------

def _extract_docx_comments(data: bytes) -> list[str]:
    """Pull review comments from word/comments.xml inside the .docx zip."""
    import xml.etree.ElementTree as ET
    import zipfile
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            if "word/comments.xml" not in z.namelist():
                return []
            xml_bytes = z.read("word/comments.xml")
    except (zipfile.BadZipFile, KeyError):
        return []
    ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    try:
        root = ET.fromstring(xml_bytes)
    except ET.ParseError:
        return []
    out: list[str] = []
    for c in root.findall(f"{{{ns}}}comment"):
        author = c.get(f"{{{ns}}}author") or "Anonymous"
        parts = [t.text for t in c.iter(f"{{{ns}}}t") if t.text]
        text = " ".join(parts).strip()
        if text:
            out.append(f"[{author}] {text}")
    return out


def extract_docx(data: bytes) -> Optional[str]:
    try:
        import docx
    except ImportError:
        _warn_once("python-docx", "python-docx not installed — DOCX files will be skipped.")
        return None
    try:
        d = docx.Document(io.BytesIO(data))
        paras = [p.text for p in d.paragraphs if p.text.strip()]
        for table in d.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    paras.append(" | ".join(cells))
        body = "\n\n".join(paras) if paras else ""
        comments = _extract_docx_comments(data)
        sections: list[str] = []
        if body:
            sections.append(body)
        if comments:
            sections.append("## Comments\n" + "\n".join(comments))
        return "\n\n".join(sections) if sections else None
    except Exception:
        return None


# --- PPTX ------------------------------------------------------------------

def extract_pptx(data: bytes) -> Optional[str]:
    try:
        from pptx import Presentation
    except ImportError:
        _warn_once("python-pptx", "python-pptx not installed — PPTX files will be skipped.")
        return None
    try:
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            chunks = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                    chunks.append(shape.text_frame.text)
            notes = ""
            if getattr(slide, "has_notes_slide", False):
                try:
                    notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                except AttributeError:
                    notes = ""
            if not chunks and not notes:
                continue
            parts = [f"# Slide {i}"]
            if chunks:
                parts.append("\n".join(chunks))
            if notes:
                parts.append(f"## Speaker Notes\n{notes}")
            slides.append("\n".join(parts))
        return "\n\n".join(slides) if slides else None
    except Exception:
        return None


# --- Text ------------------------------------------------------------------

def _decode_text(data: bytes) -> Optional[str]:
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("utf-8", errors="replace")


# --- Dispatch --------------------------------------------------------------

_HANDLERS: dict[str, Callable[[bytes], Optional[str]]] = {
    "application/pdf": extract_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_pptx,
    "text/plain": _decode_text,
    "text/markdown": _decode_text,
}


def supported_mime(mime: str) -> bool:
    return mime in _HANDLERS


def extract(mime: str, data: bytes) -> Optional[str]:
    handler = _HANDLERS.get(mime)
    if handler is None:
        return None
    return handler(data)
