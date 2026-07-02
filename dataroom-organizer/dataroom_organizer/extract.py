"""
Content + number extraction per file type.

Two guarantees from the plan (Section 2d, Section 5):
  * the system reads the *contents* of a document, not just its filename;
  * numbers from spreadsheets and cap tables are read directly by code and
    carry their source cell reference -- the AI never re-types a figure, so
    figures cannot be invented.

Everything here is best-effort and defensive: an unreadable file yields empty
text plus a flag, never an exception that stops the crawl.
"""

from __future__ import annotations

import json
import re
from pathlib import Path

from . import config as C

TEXT_CAP = 20_000          # chars of text kept per file for classification
NUMBER_LABELS = [
    "revenue", "arr", "mrr", "gross margin", "ebitda", "net income", "cash",
    "burn", "runway", "valuation", "pre-money", "post-money", "raise",
    "fully diluted", "shares outstanding", "total shares", "ownership",
    "option pool", "headcount", "customers", "bookings", "tco2", "co2",
]


def _read_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="replace")[:TEXT_CAP]
    except Exception:
        return ""


def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except Exception:  # pragma: no cover
        from PyPDF2 import PdfReader  # type: ignore
    try:
        reader = PdfReader(str(path))
        out: list[str] = []
        for page in reader.pages[:15]:
            out.append(page.extract_text() or "")
            if sum(len(x) for x in out) > TEXT_CAP:
                break
        return "\n".join(out)[:TEXT_CAP]
    except Exception:
        return ""


def _read_docx(path: Path) -> str:
    try:
        import docx
        d = docx.Document(str(path))
        parts = [p.text for p in d.paragraphs]
        for tbl in d.tables:
            for row in tbl.rows:
                parts.append("\t".join(c.text for c in row.cells))
        return "\n".join(parts)[:TEXT_CAP]
    except Exception:
        return ""


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)[:TEXT_CAP]
    except Exception:
        return ""


def _read_xlsx(path: Path) -> tuple[str, list[dict]]:
    """Return (text_blob, extracted_numbers). Numbers carry sheet!cell refs."""
    try:
        from openpyxl import load_workbook
    except Exception:  # pragma: no cover
        return "", []
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception:
        return "", []

    text_parts: list[str] = []
    numbers: list[dict] = []
    try:
        for ws in wb.worksheets:
            text_parts.append(f"[sheet: {ws.title}]")
            for row in ws.iter_rows(max_row=400, max_col=40):
                label_cell = None
                for cell in row:
                    v = cell.value
                    if v is None:
                        continue
                    if isinstance(v, str):
                        s = v.strip()
                        if s:
                            text_parts.append(s)
                            low = s.lower()
                            if any(lbl in low for lbl in NUMBER_LABELS):
                                label_cell = (cell, s)
                    elif isinstance(v, (int, float)) and label_cell is not None:
                        # number sitting on a labelled row -> capture with its cell ref
                        lab_text = label_cell[1]
                        numbers.append({
                            "label": _clean_label(lab_text),
                            "value": v,
                            "ref": f"{ws.title}!{cell.coordinate}",
                        })
                        label_cell = None
                if sum(len(p) for p in text_parts) > TEXT_CAP:
                    break
            if sum(len(p) for p in text_parts) > TEXT_CAP:
                break
    finally:
        wb.close()
    # de-dupe numbers, keep first occurrence, cap
    seen = set()
    uniq = []
    for n in numbers:
        key = (n["label"], n["value"])
        if key in seen:
            continue
        seen.add(key)
        uniq.append(n)
        if len(uniq) >= 25:
            break
    return "\n".join(text_parts)[:TEXT_CAP], uniq


def _clean_label(s: str) -> str:
    s = re.sub(r"\s+", " ", s).strip(" :\t")
    return s[:60]


def _read_gnative(path: Path) -> dict:
    """A Google-native pointer (.gdoc/.gsheet/.gslides...). Read its doc_id; the
    real content lives online and cannot be read locally."""
    info = {"doc_id": "", "url": ""}
    try:
        data = json.loads(path.read_text(encoding="utf-8", errors="replace"))
        info["doc_id"] = data.get("doc_id", "")
        if info["doc_id"]:
            info["url"] = f"https://drive.google.com/open?id={info['doc_id']}"
    except Exception:
        pass
    return info


def _read_zip(path: Path) -> str:
    """List archive member names (central directory only) as a classification hint."""
    try:
        import zipfile
        with zipfile.ZipFile(path) as z:
            names = z.namelist()[:300]
        return "\n".join(Path(n).name for n in names)[:TEXT_CAP]
    except Exception:
        return ""


def _read_eml(path: Path) -> str:
    try:
        from email import policy
        from email.parser import BytesParser
        msg = BytesParser(policy=policy.default).parse(path.open("rb"))
        head = f"Subject: {msg['subject']}\nFrom: {msg['from']}\nTo: {msg['to']}\nDate: {msg['date']}\n"
        body = ""
        try:
            part = msg.get_body(preferencelist=("plain", "html"))
            if part:
                body = part.get_content()[:4000]
        except Exception:
            pass
        return (head + "\n" + body)[:TEXT_CAP]
    except Exception:
        return ""


IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp", ".heic"}
MACRO_EXTS = {".xlsm", ".xlsb", ".pptm", ".docm"}


FAST_CONTENT_EXTS = {".xlsx", ".xlsm", ".csv", ".tsv", ".txt", ".md", ".log", ".json"}


def extract(path: Path, *, content_mode: str = "full") -> dict:
    """Extract text + numbers + format flags for a single file.

    content_mode:
      "full" — read every readable type (slowest; parses PDFs/decks).
      "fast" — read only cheap types (spreadsheets/text) so code-read numbers and
               cap-table figures still populate, but slow PDF/deck parsing is skipped.
      "none" — metadata + format flags only (read-only preview; no content reads).
    Returns {text, numbers, flags, gnative}. Defensive: an unreadable file yields
    empty text plus a flag, never an exception.
    """
    ext = C.clean_ext(path.name) or path.suffix.lower()
    flags: list[str] = []
    text = ""
    numbers: list[dict] = []
    gnative: dict = {}

    try:
        size = path.stat().st_size
    except Exception:
        size = 0
    too_big = size > C.MAX_EXTRACT_MB * 1024 * 1024

    # ---- format-only signals (always cheap, no content read) ----
    if ext in C.GNATIVE_EXTS:
        flags += ["google-native", "content-not-read-locally"]
        gnative = _read_gnative(path)
    elif ext in IMAGE_EXTS:
        flags.append("image-only")
    elif ext in C.VIDEO_EXTS:
        flags.append("video")
    elif ext in C.LEGACY_BINARY_EXTS or ext in {".xls"}:
        flags.append("legacy-binary-format")
    elif ext in C.ARCHIVE_EXTS:
        flags.append("archive")
    if ext in MACRO_EXTS:
        flags.append("macro-enabled")
    if size > 25 * 1024 * 1024:
        flags.append("large-file")

    # ---- content read (skippable) ----
    if content_mode != "none" and not flags_block_content(flags):
        if too_big:
            flags.append("content-skipped-too-large")
        elif content_mode == "fast" and ext not in FAST_CONTENT_EXTS:
            flags.append("content-skipped-fast-mode")
        else:
            if ext in {".txt", ".md", ".csv", ".tsv", ".json", ".log"}:
                text = _read_text_file(path)
            elif ext == ".pdf":
                text = _read_pdf(path)
                if not text.strip():
                    flags.append("pdf-no-text-layer")
            elif ext == ".docx":
                text = _read_docx(path)
            elif ext == ".pptx":
                text = _read_pptx(path)
            elif ext in {".xlsx", ".xlsm"}:
                text, numbers = _read_xlsx(path)
            elif ext in C.ARCHIVE_EXTS:
                text = _read_zip(path)
            elif ext in C.EMAIL_EXTS:
                text = _read_eml(path) if ext == ".eml" else ""
                if not text:
                    flags.append("email-not-parsed")
            elif ext == "":
                flags.append("no-extension")
            elif ext not in C.GNATIVE_EXTS | IMAGE_EXTS | C.VIDEO_EXTS \
                    | C.LEGACY_BINARY_EXTS | {".xls"}:
                flags.append("unknown-format")

    return {"text": text or "", "numbers": numbers, "flags": flags, "gnative": gnative}


def flags_block_content(flags: list[str]) -> bool:
    """Formats whose 'content' is not text we read here."""
    return any(f in flags for f in
               ("google-native", "image-only", "video", "legacy-binary-format"))
