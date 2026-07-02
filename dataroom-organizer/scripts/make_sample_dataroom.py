"""
Generate a synthetic 'messy' data room that reproduces the problems the plan
describes, so the reorganizer can be tested end-to-end without the live drive.

The individual documents are deliberately *complex and realistic*: multi-section
IC memos, multi-sheet cap tables and financial models with real line items,
multi-page board decks / term sheets / diligence and impact reports, and investor
updates with full KPI tables -- so classification, code-read numbers, versioning
and the before/after comparison are exercised on lifelike content.

Pilot archetypes (plan, Section 4):
  * Cambium   — richly documented (every bucket, versions, duplicates, tricky names)
  * Reframe   — mixed (gaps, a cap-table *screenshot*)
  * Hestia    — nearly empty
  * NxLite    — older, still active (legacy naming, a monthly update *series*)
  * Veckta / Nth Cycle — scattered: their docs live in cross-cutting theme folders

Usage:  python scripts/make_sample_dataroom.py [OUT_DIR]
"""

from __future__ import annotations

import os
import shutil
import sys
import textwrap
import time
from datetime import datetime
from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Font
import docx
from docx.shared import Pt
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import LETTER

ROSTER = ["Cambium", "Reframe", "Hestia", "NxLite", "Veckta", "Nth Cycle"]


def _set_date(p: Path, iso: str) -> None:
    t = time.mktime(datetime.strptime(iso, "%Y-%m-%d").timetuple())
    os.utime(p, (t, t))


# --------------------------------------------------------------------------- #
# Format mechanics: rich PDF / DOCX / XLSX builders
# --------------------------------------------------------------------------- #

def pdf_doc(p: Path, title: str, subtitle: str, blocks: list[tuple], date: str):
    """blocks: list of ('H2'|'P'|'BULLET'|'TABLE', content)."""
    c = canvas.Canvas(str(p), pagesize=LETTER)
    width, height = LETTER
    margin = 60
    y = height - margin

    def newpage():
        nonlocal y
        c.showPage(); y = height - margin

    def write(text, font="Helvetica", size=10.5, gap=14, indent=0):
        nonlocal y
        if y < margin + 30:
            newpage()
        c.setFont(font, size)
        c.drawString(margin + indent, y, text)
        y -= gap

    c.setFont("Helvetica-Bold", 18); c.drawString(margin, y, title); y -= 24
    c.setFont("Helvetica-Oblique", 11); c.drawString(margin, y, subtitle); y -= 22
    for kind, content in blocks:
        if kind == "H2":
            y -= 6; write(content, "Helvetica-Bold", 13, 18)
        elif kind == "P":
            for w in textwrap.wrap(content, 98):
                write(w, "Helvetica", 10.5, 13.5)
            y -= 5
        elif kind == "BULLET":
            wrapped = textwrap.wrap(content, 92)
            for i, w in enumerate(wrapped):
                write(("• " if i == 0 else "  ") + w, "Helvetica", 10.5, 13.5, indent=12)
        elif kind == "TABLE":
            for row in content:
                write("  ".join(f"{str(x):<22}" for x in row), "Courier", 9, 12.5)
            y -= 6
    c.showPage(); c.save(); _set_date(p, date)


def docx_doc(p: Path, title: str, blocks: list[tuple], date: str):
    """blocks: ('H1'|'H2'|'P'|'BULLET'|'TABLE', content)."""
    d = docx.Document()
    style = d.styles["Normal"]; style.font.name = "Calibri"; style.font.size = Pt(11)
    d.add_heading(title, 0)
    for kind, content in blocks:
        if kind == "H1":
            d.add_heading(content, level=1)
        elif kind == "H2":
            d.add_heading(content, level=2)
        elif kind == "P":
            d.add_paragraph(content)
        elif kind == "BULLET":
            d.add_paragraph(content, style="List Bullet")
        elif kind == "TABLE":
            rows = content
            t = d.add_table(rows=len(rows), cols=len(rows[0]))
            t.style = "Light Grid Accent 1"
            for ri, row in enumerate(rows):
                for ci, val in enumerate(row):
                    t.cell(ri, ci).text = str(val)
    d.save(p); _set_date(p, date)


def xlsx_book(p: Path, sheets: dict[str, list[list]], date: str):
    wb = Workbook(); first = True
    for name, rows in sheets.items():
        ws = wb.active if first else wb.create_sheet()
        ws.title = name[:31]; first = False
        for r in rows:
            ws.append(r)
        for cell in ws[1]:
            cell.font = Font(bold=True)
    wb.save(p); _set_date(p, date)


def png_captable(p: Path, company: str, rows: list[list], date: str):
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (820, 420), "white"); dr = ImageDraw.Draw(img)
    dr.text((24, 18), f"{company} — CAPITALIZATION TABLE (exported screenshot)", fill="black")
    dr.line((24, 44, 796, 44), fill="black")
    for ri, row in enumerate(rows):
        y = 60 + ri * 30
        for ci, cell in enumerate(row):
            dr.text((28 + ci * 200, y), str(cell), fill="black")
        dr.line((24, y + 22, 796, y + 22), fill=(210, 210, 210))
    img.save(p); _set_date(p, date)


def pptx_deck(p: Path, company: str, sector: str, slides: list[tuple], date: str):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = f"{company}"
    s.placeholders[1].text = f"{sector} · Series A Investor Deck"
    for title, body in slides:
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = title
        sl.placeholders[1].text = body
    prs.save(p); _set_date(p, date)


def dup(src: Path, dst: Path):
    dst.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(src, dst)


# --------------------------------------------------------------------------- #
# Bespoke document content
# --------------------------------------------------------------------------- #

def ic_memo_blocks(company: str, sector: str, tag: str, nums: dict) -> list[tuple]:
    redacted = "redacted" in tag
    blocks = [
        ("P", f"Classification: {tag.upper()} — VoLo Earth Investment Committee. "
              f"Sector: {sector}. Stage: Series A. Round size: ${nums['round']/1e6:.1f}M."),
        ("H1", "1. Executive Summary"),
        ("P", f"{company} is a {sector.lower()} company addressing a large and growing "
              f"decarbonization market. We recommend a ${nums['check']/1e6:.1f}M investment "
              f"at a ${nums['pre']/1e6:.0f}M pre-money valuation (${nums['post']/1e6:.0f}M "
              f"post-money), for approximately {nums['own']:.1f}% ownership on a fully "
              f"diluted basis. The thesis rests on a differentiated technology, an "
              f"experienced team, and early commercial validation."),
        ("H1", "2. Company & Team"),
        ("P", f"Founded by repeat operators with deep domain expertise, {company} has "
              f"grown to {nums['headcount']} employees. The founding team previously "
              f"scaled hardware and software in adjacent climate markets."),
        ("H1", "3. Market & TAM"),
        ("P", f"We estimate a serviceable addressable market exceeding "
              f"${nums['tam']/1e9:.0f}B, driven by regulatory tailwinds, corporate net-zero "
              f"commitments, and falling technology costs. The beachhead segment alone "
              f"represents a multi-hundred-million-dollar opportunity."),
        ("H1", "4. Product & Technology"),
        ("P", "The core technology has been validated at pilot scale and is protected by "
              "a growing patent portfolio. Key technical risks relate to scale-up yield "
              "and unit economics at commercial volume, addressed in technical diligence."),
        ("H1", "5. Traction & KPIs"),
        ("TABLE", [["Metric", "FY22", "FY23", "FY24"],
                   ["Revenue ($M)", nums['rev'][0], nums['rev'][1], nums['rev'][2]],
                   ["ARR ($M)", round(nums['arr']*0.4, 1), round(nums['arr']*0.7, 1), nums['arr']],
                   ["Customers", nums['cust'][0], nums['cust'][1], nums['cust'][2]],
                   ["Gross margin %", 38, 47, 55]]),
        ("H1", "6. Business Model & Unit Economics"),
        ("P", "Revenue is a blend of product sales and recurring software/service "
              "contracts. Contribution margin improves materially with scale as fixed "
              "manufacturing and platform costs are absorbed."),
    ]
    if not redacted:
        blocks += [
            ("H1", "7. Financials & Projections"),
            ("P", f"FY24 revenue of ${nums['rev'][2]}M with a monthly net burn of "
                  f"~${nums['burn']/1e3:.0f}K and ${nums['cash']/1e6:.1f}M of cash, "
                  f"implying ~{nums['runway']} months of runway pre-round."),
            ("H1", "8. Deal Terms"),
            ("BULLET", f"Security: Series A Preferred Stock at ${nums['pps']:.2f}/share."),
            ("BULLET", "Liquidation preference: 1x non-participating."),
            ("BULLET", "Board: one VoLo seat; standard protective provisions."),
            ("BULLET", "Pro-rata rights and information rights granted."),
        ]
    else:
        blocks += [("H1", "7. Financials & Deal Terms"),
                   ("P", "[Sensitive financial terms and deal economics redacted in this "
                         "EXTERNAL copy. Refer to the internal memo for full detail.]")]
    blocks += [
        ("H1", "9. Risks & Mitigants"),
        ("BULLET", "Scale-up / manufacturing risk — staged milestone-based financing."),
        ("BULLET", "Customer concentration — active pipeline diversification."),
        ("BULLET", "Policy dependence — economics modeled with and without subsidies."),
        ("H1", "10. Impact Thesis"),
        ("P", f"At scale, {company} is projected to abate {nums['tco2']:,} tCO2e annually, "
              f"aligning with VoLo Earth's climate mandate."),
        ("H1", "11. Recommendation"),
        ("P", f"APPROVE a ${nums['check']/1e6:.1f}M investment, subject to confirmatory "
              f"diligence and final legal documentation."),
    ]
    if "draft" in tag:
        blocks.insert(1, ("P", "*** DRAFT — pre-decisional, not for circulation. Figures subject to change. ***"))
    return blocks


def captable_sheets(company: str, fd: int, pps: float, pre: int, post: int) -> dict:
    founders = int(fd * 0.50); seed = int(fd * 0.18); safe = int(fd * 0.07)
    seriesa = int(fd * 0.15); pool = fd - founders - seed - safe - seriesa
    summary = [
        ["Holder", "Security", "Shares", "Price/Share", "Investment", "Ownership % (FD)"],
        ["Founders (3)", "Common", founders, 0.001, round(founders*0.001), round(100*founders/fd, 2)],
        ["Pre-Seed SAFE", "SAFE→Pref", safe, round(pps*0.5, 4), round(safe*pps*0.5), round(100*safe/fd, 2)],
        ["Seed Preferred", "Series Seed", seed, round(pps*0.6, 4), round(seed*pps*0.6), round(100*seed/fd, 2)],
        ["Series A (VoLo + co.)", "Series A Pref", seriesa, pps, round(seriesa*pps), round(100*seriesa/fd, 2)],
        ["Option Pool (unissued)", "Options", pool, pps, 0, round(100*pool/fd, 2)],
        [],
        ["Fully diluted shares", "", fd, "", "", 100.0],
        ["Pre-money valuation", "", pre, "", "", ""],
        ["Post-money valuation", "", post, "", "", ""],
        ["Price per share", "", pps, "", "", ""],
    ]
    rounds = [
        ["Round", "Date", "Security", "Raised", "Pre-money", "Post-money", "Price/Share"],
        ["Pre-Seed", "2021-06", "SAFE", 750_000, "", 6_000_000, round(pps*0.5, 4)],
        ["Seed", "2022-09", "Series Seed Pref", 3_500_000, 12_000_000, 15_500_000, round(pps*0.6, 4)],
        ["Series A", "2024-02", "Series A Pref", post-pre, pre, post, pps],
    ]
    waterfall = [
        ["Exit value", "Pref return", "Common return", "VoLo proceeds (15%)"],
        [50_000_000, 8_000_000, 42_000_000, round(50_000_000*0.15)],
        [100_000_000, 8_000_000, 92_000_000, round(100_000_000*0.15)],
        [250_000_000, 8_000_000, 242_000_000, round(250_000_000*0.15)],
    ]
    return {"Summary": summary, "By Round": rounds, "Waterfall": waterfall}


def model_sheets(company: str, rev: list[int], arr: float, burn: int, cash: int, runway: int) -> dict:
    years = ["FY22", "FY23", "FY24", "FY25E", "FY26E"]
    revenue = rev + [int(rev[-1]*2.2), int(rev[-1]*4.6)]
    cogs = [int(r*0.55) for r in revenue]
    gp = [revenue[i]-cogs[i] for i in range(5)]
    rnd = [int(r*0.45)+800_000 for r in revenue]
    sm = [int(r*0.30)+400_000 for r in revenue]
    ga = [int(r*0.18)+300_000 for r in revenue]
    ebitda = [gp[i]-rnd[i]-sm[i]-ga[i] for i in range(5)]
    pnl = [["Line item ($)"] + years,
           ["Revenue"] + revenue,
           ["COGS"] + cogs,
           ["Gross Profit"] + gp,
           ["R&D"] + rnd, ["Sales & Marketing"] + sm, ["G&A"] + ga,
           ["EBITDA"] + ebitda,
           ["ARR (exit)"] + [int(arr*1e6*(0.5+0.25*i)) for i in range(5)]]
    cashflow = [["Cash ($)"] + years,
                ["Beginning cash"] + [cash + i*1_000_000 for i in range(5)],
                ["Net burn"] + [burn*12]*5,
                ["Runway (months)", runway, runway-3, runway-6, "—", "—"],
                ["Ending cash"] + [max(0, cash - burn*12*(i)) for i in range(5)]]
    head = [["Function"] + years,
            ["Engineering", 6, 14, 28, 45, 70],
            ["Commercial", 2, 5, 11, 20, 34],
            ["Operations", 2, 4, 8, 14, 22],
            ["G&A", 1, 3, 5, 8, 12]]
    assumptions = [["Assumption", "Value"],
                   ["Revenue CAGR (FY24-26E)", "115%"],
                   ["Gross margin target", "55%"],
                   ["Monthly net burn ($)", burn],
                   ["Cash on hand ($)", cash],
                   ["Runway (months)", runway]]
    return {"Assumptions": assumptions, "P&L": pnl, "Cash Flow": cashflow, "Headcount": head}


def board_blocks(company: str, period: str, nums: dict) -> list[tuple]:
    return [
        ("H2", "Agenda"),
        ("BULLET", "Approval of prior minutes"), ("BULLET", "CEO update & KPIs"),
        ("BULLET", "Financial review & runway"), ("BULLET", "Product roadmap"),
        ("BULLET", "Hiring plan"), ("BULLET", "Resolutions & asks"),
        ("H2", "KPI Dashboard"),
        ("TABLE", [["KPI", "Prior Q", "This Q", "Plan"],
                   ["Revenue ($K)", nums['rev_q'][0], nums['rev_q'][1], nums['rev_q'][2]],
                   ["ARR ($K)", nums['arr_q'][0], nums['arr_q'][1], nums['arr_q'][2]],
                   ["Cash ($M)", nums['cash_q'][0], nums['cash_q'][1], "—"],
                   ["Runway (mo)", nums['rw_q'][0], nums['rw_q'][1], "—"]]),
        ("H2", "Financial Review"),
        ("P", f"Revenue grew quarter-over-quarter; net burn held roughly flat. "
              f"Current cash provides approximately {nums['rw_q'][1]} months of runway. "
              f"Management recommends opening the next financing in two quarters."),
        ("H2", "Product Roadmap"),
        ("P", "Roadmap prioritizes commercial-scale reliability, cost-down initiatives, "
              "and two flagship customer deployments in the next two quarters."),
        ("H2", "Resolutions"),
        ("P", "RESOLVED, that the Board approves the expansion of the option pool by "
              "200,000 shares; and FURTHER RESOLVED, that the Company is authorized to "
              "enter into the customer agreements described herein."),
    ]


def term_sheet_blocks(company: str, nums: dict) -> list[tuple]:
    return [
        ("P", f"This non-binding term sheet summarizes the principal terms of a proposed "
              f"Series A Preferred Stock financing of {company} (the 'Company')."),
        ("H2", "1. Offering Terms"),
        ("TABLE", [["Term", "Value"],
                   ["Security", "Series A Preferred Stock"],
                   ["Amount raised", f"${nums['round']/1e6:.1f}M"],
                   ["Pre-money valuation", f"${nums['pre']/1e6:.0f}M"],
                   ["Post-money valuation", f"${nums['post']/1e6:.0f}M"],
                   ["Price per share", f"${nums['pps']:.2f}"],
                   ["Option pool", "10% post-financing, fully diluted"]]),
        ("H2", "2. Liquidation Preference"),
        ("P", "1x non-participating preference, with the right to convert to Common."),
        ("H2", "3. Anti-dilution & Conversion"),
        ("P", "Broad-based weighted-average anti-dilution. Automatic conversion on a "
              "qualified IPO above an agreed threshold."),
        ("H2", "4. Board & Protective Provisions"),
        ("P", "The Board shall consist of the CEO, one founder, and one VoLo designee. "
              "Customary protective provisions apply to Preferred."),
        ("H2", "5. Investor Rights"),
        ("BULLET", "Pro-rata participation rights in future financings."),
        ("BULLET", "Information rights: monthly and quarterly reporting."),
        ("BULLET", "Founder vesting: 4 years with a 1-year cliff."),
        ("H2", "6. Conditions & Confidentiality"),
        ("P", "Closing is subject to satisfactory legal, technical and financial "
              "diligence. The parties agree to keep these terms confidential."),
    ]


def investor_update_blocks(company: str, month: str, nums: dict) -> list[tuple]:
    return [
        ("P", f"Investor update for {month}. Prepared for VoLo Earth and other investors. "
              f"Confidential."),
        ("H2", "Key Metrics"),
        ("TABLE", [["Metric", "Last month", "This month", "MoM"],
                   ["MRR ($K)", nums['mrr'][0], nums['mrr'][1], f"+{nums['mrr_g']}%"],
                   ["ARR ($M)", nums['arr'][0], nums['arr'][1], ""],
                   ["Net burn ($K)", nums['burn'][0], nums['burn'][1], ""],
                   ["Cash ($M)", nums['cash'][0], nums['cash'][1], ""],
                   ["Runway (months)", nums['rw'][0], nums['rw'][1], ""],
                   ["Headcount", nums['hc'][0], nums['hc'][1], ""]]),
        ("H2", "Highlights"),
        ("BULLET", "Signed a flagship commercial customer with multi-year volume commitment."),
        ("BULLET", "Hit a key technical milestone ahead of plan."),
        ("H2", "Lowlights"),
        ("BULLET", "Supply-chain lead times slipped one deployment by ~6 weeks."),
        ("H2", "Asks"),
        ("BULLET", "Intros to two named strategic customers."),
        ("BULLET", "Candidates for VP of Manufacturing."),
    ]


# --------------------------------------------------------------------------- #
# Build the tree
# --------------------------------------------------------------------------- #

def build(root: Path):
    if root.exists():
        shutil.rmtree(root)
    root.mkdir(parents=True)

    def D(*parts) -> Path:
        d = root.joinpath(*parts); d.mkdir(parents=True, exist_ok=True); return d

    # ---------- Cambium: richly documented (biochar carbon removal) ----------
    cam = D("Cambium")
    cnums = dict(round=8_000_000, check=5_000_000, pre=34_000_000, post=42_000_000,
                 own=11.9, headcount=31, tam=12, rev=[0.4, 1.6, 4.2], arr=5.1,
                 cust=[3, 9, 22], burn=600_000, cash=9_000_000, runway=15,
                 pps=3.36, tco2=120_000)
    cam_memo_final = cam / "Investment Memo - Cambium - internal - final v2.docx"
    docx_doc(cam_memo_final, "Investment Memo — Cambium (Carbon Removal)",
             ic_memo_blocks("Cambium", "Carbon Removal", "internal final", cnums), "2024-02-10")
    docx_doc(cam / "Investment Memo - Cambium - internal - draft.docx",
             "Investment Memo — Cambium (Carbon Removal)",
             ic_memo_blocks("Cambium", "Carbon Removal", "internal draft", cnums), "2024-01-20")
    docx_doc(cam / "Investment Memo Cambium external redacted.docx",
             "Investment Memo — Cambium (Carbon Removal)",
             ic_memo_blocks("Cambium", "Carbon Removal", "external redacted", cnums), "2024-02-12")

    cam_cap = cam / "Cambium Cap Table 2024-03-15.xlsx"
    xlsx_book(cam_cap, captable_sheets("Cambium", 12_500_000, 3.36, 34_000_000, 42_000_000), "2024-03-15")
    xlsx_book(cam / "Cambium Cap Table 2023-09-01.xlsx",
              captable_sheets("Cambium", 9_800_000, 1.58, 12_000_000, 15_500_000), "2023-09-01")
    xlsx_book(cam / "Cambium Financial Model v3.xlsm",
              model_sheets("Cambium", [400_000, 1_600_000, 4_200_000], 5.1, 600_000, 9_000_000, 15), "2024-03-01")

    pdf_doc(cam / "Cambium Board Deck 2024-02 (final).pdf", "Cambium — Board Meeting", "February 2024",
            board_blocks("Cambium", "Q1-2024",
                         dict(rev_q=[900, 1300, 1500], arr_q=[3800, 5100, 6000],
                              cash_q=[10.2, 9.0], rw_q=[17, 15])), "2024-02-15")
    pdf_doc(cam / "Cambium Q1'24 Update (final) v2 #2.pdf", "Cambium — Investor Update", "Q1 2024",
            investor_update_blocks("Cambium", "March 2024",
                                   dict(mrr=[420, 470], mrr_g=12, arr=[5.0, 5.6], burn=[610, 600],
                                        cash=[9.0, 8.4], rw=[15, 14], hc=[29, 31])), "2024-04-05")
    pdf_doc(cam / "Cambium Term Sheet - executed.pdf", "Cambium — Series A Term Sheet", "Executed",
            term_sheet_blocks("Cambium", cnums), "2024-01-15")
    pptx_deck(cam / "Cambium pitch deck.pptx", "Cambium", "Carbon Removal",
              [("The Problem", "Gigatonne-scale carbon removal is needed; durable, verifiable methods are scarce."),
               ("The Solution", "Engineered biochar with measurable, permanent carbon storage."),
               ("Technology", "Continuous pyrolysis with proprietary feedstock handling."),
               ("Traction", "22 customers, $4.2M FY24 revenue, $5.1M ARR."),
               ("Market", "$12B+ durable CDR market by 2030."),
               ("Team", "Repeat climate-hardware founders."),
               ("Financials", "115% revenue CAGR; 55% gross margin target."),
               ("The Ask", "$5M Series A at $34M pre.")], "2024-01-08")
    txt(cam / "Cambium call notes 2024-01-10.txt",
        "CALL NOTES — Cambium <> VoLo Earth\nDate: 2024-01-10\nAttendees: VoLo (J. Z.), Cambium (CEO, CFO)\n\n"
        "Agenda: diligence kickoff, cap table walkthrough, customer pipeline.\n\n"
        "Discussion:\n - CEO walked through the pyrolysis scale-up plan and yield improvements.\n"
        " - CFO shared updated model; FY24 revenue $4.2M, runway ~15 months.\n"
        " - Pipeline: 3 enterprise CDR buyers in late-stage negotiation.\n\n"
        "Action items:\n - Cambium to send executed customer LOIs and 409A.\n"
        " - VoLo to circulate draft term sheet by 2024-01-15.\n"
        "Next steps: technical diligence call scheduled for next week.\n", "2024-01-10")
    pdf_doc(cam / "Cambium SBIR grant award.pdf", "Cambium — Grant Award Notice", "Department of Energy SBIR",
            [("H2", "Award"),
             ("TABLE", [["Field", "Value"], ["Award number", "DE-SC-SBIR-2023-1187"],
                        ["Agency", "U.S. Department of Energy"], ["Amount", "$1,150,000"],
                        ["Period of performance", "12 months"], ["Type", "Non-dilutive (Phase II SBIR)"]]),
             ("H2", "Scope"),
             ("P", "Funds development of continuous pyrolysis reactor controls and carbon "
                   "verification methodology. Non-dilutive grant award; no equity or IP "
                   "assignment to the funding agency beyond standard government use rights.")], "2023-11-01")
    pdf_doc(cam / "Cambium impact report 2023.pdf", "Cambium — 2023 Impact Report", "Carbon & LCA",
            [("H2", "Carbon Abatement"),
             ("P", "Cumulative durable carbon removal delivered in 2023 with third-party "
                   "verification, on a path to 120,000 tCO2e per year at commercial scale."),
             ("TABLE", [["Metric", "2023"], ["Durable CDR delivered (tCO2e)", "8,400"],
                        ["Permanence (years)", ">1,000"], ["Scope 1 emissions (tCO2e)", "1,200"],
                        ["Scope 2 emissions (tCO2e)", "640"], ["Net negativity ratio", "6.7x"]]),
             ("H2", "Life Cycle Assessment"),
             ("P", "Cradle-to-grave LCA confirms net-negative emissions across feedstock "
                   "sourcing, pyrolysis, and application. Methodology aligned with leading "
                   "CDR registries; Scope 1, 2 and 3 boundaries documented.")], "2023-12-31")
    # an opaque binary with no signals -> should land in 15/Pending for a human
    (cam / "random_unlabeled_thing.bin").write_bytes(os.urandom(4096))

    # ---------- Reframe: mixed; cap table is a screenshot (image-only) ----------
    ref = D("Reframe")
    rnums = dict(round=4_000_000, check=2_000_000, pre=13_000_000, post=15_000_000,
                 own=13.3, headcount=14, tam=8, rev=[0.2, 0.5, 1.1], arr=1.4,
                 cust=[1, 4, 9], burn=300_000, cash=3_000_000, runway=14,
                 pps=2.10, tco2=40_000)
    docx_doc(ref / "Reframe Investment Memo.docx", "Investment Memo — Reframe (Building Decarbonization)",
             ic_memo_blocks("Reframe", "Building Decarbonization", "internal", rnums), "2024-03-02")
    png_captable(ref / "Reframe cap table.png", "Reframe",
                 [["Holder", "Security", "Shares", "FD %"],
                  ["Founders", "Common", "5,200,000", "52%"],
                  ["Seed Fund", "Preferred", "2,100,000", "21%"],
                  ["VoLo Earth", "Preferred", "1,400,000", "14%"],
                  ["Option Pool", "Options", "1,300,000", "13%"]], "2024-02-20")
    pdf_doc(ref / "Reframe customer LOI.pdf", "Reframe — Letter of Intent", "Acme Properties",
            [("P", "This Letter of Intent is entered into between Reframe, Inc. and Acme "
                   "Properties LLC regarding a commercial deployment of Reframe's building "
                   "decarbonization platform."),
             ("H2", "Scope & Volume"),
             ("TABLE", [["Term", "Value"], ["Buildings (Year 1)", "12"],
                        ["Estimated contract value", "$640,000"], ["Term", "36 months"],
                        ["Exclusivity", "Non-exclusive"]]),
             ("H2", "Conditions"),
             ("P", "Non-binding; subject to a definitive master services agreement and "
                   "successful pilot. Commercial pipeline reference for diligence.")], "2024-03-10")
    xlsx_book(ref / "Reframe model.xlsx",
              model_sheets("Reframe", [200_000, 500_000, 1_100_000], 1.4, 300_000, 3_000_000, 14), "2024-03-05")
    txt(ref / "notes.txt", "Quick notes about Reframe — follow up on pilot results next week.", "2024-03-12")

    # ---------- Hestia: nearly empty (home electrification) ----------
    hes = D("Hestia")
    pdf_doc(hes / "Hestia one-pager.pdf", "Hestia — Company One-Pager", "Home Electrification",
            [("H2", "Overview"),
             ("P", "Hestia makes whole-home electrification turnkey for homeowners and "
                   "installers. Pre-seed stage; seeking design-partner utilities."),
             ("H2", "Product"),
             ("P", "Software + financing layer that bundles heat pumps, panels and storage."),
             ("BULLET", "Stage: Pre-seed"), ("BULLET", "Team: 4"),
             ("BULLET", "The problem: electrification is too complex for homeowners."),
             ("BULLET", "The solution: a single managed offering.")], "2023-10-01")

    # ---------- NxLite: older, still active; a monthly update SERIES ----------
    nx = D("NxLite")
    xlsx_book(nx / "NXLITE_CAPTABLE_FINAL.XLSX",
              captable_sheets("NxLite", 6_400_000, 2.81, 14_000_000, 18_000_000), "2022-06-01")
    for mon, d, mrr in [("January 2024", "2024-01-31", [180, 205]),
                        ("February 2024", "2024-02-29", [205, 232])]:
        pdf_doc(nx / f"NxLite investor update {mon.split()[0]} 2024.pdf",
                "NxLite — Investor Update", mon,
                investor_update_blocks("NxLite", mon,
                                       dict(mrr=mrr, mrr_g=14, arr=[2.4, 2.7], burn=[240, 250],
                                            cash=[5.1, 4.9], rw=[20, 19], hc=[18, 19])), d)
    pdf_doc(nx / "NxLite_term_sheet_2022.pdf", "NxLite — Seed Term Sheet", "2022",
            term_sheet_blocks("NxLite", dict(round=3_000_000, pre=11_000_000, post=14_000_000, pps=2.19)),
            "2022-05-15")

    # ---------- THEME folders: the competing 'second home' ----------
    captables = D("Cap Tables")
    xlsx_book(captables / "Veckta cap table 2024-01.xlsx",
              captable_sheets("Veckta", 7_100_000, 3.10, 19_000_000, 22_000_000), "2024-01-15")
    xlsx_book(captables / "NthCycle_captable.xlsx",
              captable_sheets("Nth Cycle", 8_300_000, 3.13, 23_000_000, 26_000_000), "2024-02-01")
    dup(cam_cap, captables / "Cambium Cap Table 2024-03-15.xlsx")               # exact duplicate

    memos = D("Investment Memos")
    docx_doc(memos / "Veckta investment memo.docx", "Investment Memo — Veckta (Microgrids)",
             ic_memo_blocks("Veckta", "Distributed Energy / Microgrids", "internal",
                            dict(round=6_000_000, check=3_000_000, pre=19_000_000, post=22_000_000,
                                 own=13.6, headcount=22, tam=15, rev=[0.6, 1.4, 3.0], arr=2.8,
                                 cust=[4, 8, 15], burn=400_000, cash=6_000_000, runway=15,
                                 pps=3.10, tco2=90_000)), "2024-01-10")
    docx_doc(memos / "Nth Cycle - IC memo.docx", "Investment Memo — Nth Cycle (Critical Minerals)",
             ic_memo_blocks("Nth Cycle", "Critical Minerals Recovery", "internal",
                            dict(round=9_000_000, check=4_000_000, pre=23_000_000, post=26_000_000,
                                 own=15.4, headcount=27, tam=20, rev=[0.5, 1.8, 5.0], arr=3.4,
                                 cust=[2, 5, 11], burn=550_000, cash=8_000_000, runway=15,
                                 pps=3.13, tco2=150_000)), "2024-02-05")
    dup(cam_memo_final, memos / "Cambium Investment Memo - internal - final v2.docx")  # exact duplicate

    dil = D("Diligence")
    pdf_doc(dil / "Veckta technical diligence.pdf", "Veckta — Technical Diligence", "Confidential",
            [("H2", "Technology Overview"),
             ("P", "Veckta's microgrid design-and-procurement platform optimizes distributed "
                   "energy resources for commercial and industrial sites."),
             ("H2", "Technology Readiness"),
             ("TABLE", [["Dimension", "Assessment"], ["TRL", "7-8 (field-validated)"],
                        ["Scale-up risk", "Low-moderate"], ["IP", "2 issued patents, 3 pending"]]),
             ("H2", "Key Risks"),
             ("BULLET", "Dependence on third-party hardware supply."),
             ("BULLET", "Project-based revenue lumpiness."),
             ("P", "Intellectual property and patent landscape reviewed; freedom-to-operate "
                   "opinion pending from counsel.")], "2024-01-20")
    pdf_doc(dil / "Nth Cycle patent summary.pdf", "Nth Cycle — Patent Portfolio Summary", "IP Diligence",
            [("H2", "Patent Families"),
             ("TABLE", [["Family", "Title", "Status", "Jurisdictions"],
                        ["P1", "Electro-extraction of metals", "Granted", "US, EP, CN"],
                        ["P2", "Continuous recovery cell", "Pending", "US, PCT"],
                        ["P3", "Selective leaching process", "Pending", "US"]]),
             ("H2", "Assessment"),
             ("P", "Strong intellectual property position around the core electro-extraction "
                   "process. Patent portfolio summary indicates broad claims with defensible "
                   "priority dates.")], "2024-02-08")

    updates = D("Investor Updates")
    pdf_doc(updates / "Veckta investor update Mar 2024.pdf", "Veckta — Investor Update", "March 2024",
            investor_update_blocks("Veckta", "March 2024",
                                   dict(mrr=[260, 300], mrr_g=15, arr=[2.6, 3.0], burn=[390, 400],
                                        cash=[6.2, 5.8], rw=[16, 15], hc=[20, 22])), "2024-03-31")

    boards = D("Board Decks")
    pdf_doc(boards / "Nth Cycle board deck 2024-03.pdf", "Nth Cycle — Board Meeting", "March 2024",
            board_blocks("Nth Cycle", "Q1-2024",
                         dict(rev_q=[1100, 1500, 1700], arr_q=[3000, 3400, 4000],
                              cash_q=[8.6, 8.0], rw_q=[16, 15])), "2024-03-20")

    # ---------- master investment tracker, six dated copies (root level) ----------
    cols = [["Company", "Sector", "Entry", "Stage", "Invested ($)", "Own %", "Last val ($)", "MOIC", "Status"],
            ["Cambium", "Carbon Removal", "2022-09", "Series A", 5_000_000, 11.9, 42_000_000, 1.4, "Active"],
            ["Reframe", "Buildings", "2023-05", "Seed", 2_000_000, 13.3, 15_000_000, 1.1, "Active"],
            ["Veckta", "Microgrids", "2023-08", "Series A", 3_000_000, 13.6, 22_000_000, 1.2, "Active"],
            ["Nth Cycle", "Crit. Minerals", "2023-11", "Series A", 4_000_000, 15.4, 26_000_000, 1.0, "Active"],
            ["NxLite", "Materials", "2022-05", "Seed", 1_500_000, 9.0, 18_000_000, 1.3, "Active"]]
    for i, d in enumerate(["2024-01-05", "2024-02-09", "2024-03-15",
                           "2024-04-20", "2024-05-18", "2024-06-21"]):
        sheets = {"Portfolio": [[f"Master Investment Tracker — as of {d}"]] + cols
                  + [["Total invested", "", "", "", 15_500_000 + i * 250_000, "", "", "", ""]]}
        xlsx_book(root / f"Master Investment Tracker {d}.xlsx", sheets, d)

    return root


def txt(p: Path, body: str, date: str):
    p.write_text(body, encoding="utf-8"); _set_date(p, date)


def main():
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).resolve().parents[1] / "sample_drive"
    root = build(out)
    n = sum(1 for _ in root.rglob("*") if _.is_file())
    print(f"Built sample drive at {root} ({n} files across {len(ROSTER)} companies + theme folders)")


if __name__ == "__main__":
    main()
